"""The `.pptx` writer's foundation: the ONE OPC-zip normalizer + the writer's shared constants.

WHY THIS MODULE EXISTS (Phase 2 / WKLY-01, promoting `tests/fixtures/weekly/_determinism.py`).
The Phase 1 determinism spike proved its finding inside `tests/`, which was right for a spike and
wrong for a renderer: the normalizer is the mechanism the *writer* depends on, so leaving it under
`tests/fixtures/weekly/` forced three importers to mutate `sys.path` to reach it (IN-03) and left
the repo one careless copy-paste away from a second implementation of "make this deterministic".
`.planning/notes/2026-08-29-pptx-determinism-decision.md` records the rule this promotion enforces:
**ONE normalizer contract.** A second implementation drifts from the first exactly as a second
epoch sentinel would, and the drift is invisible until a gate that trusted both goes red.

THE FIX. The five public names below moved here VERBATIM — bodies, comments and the explanatory
prose unchanged — so a reviewer can diff the move against the deleted fixture module and see that
nothing was "tidied". Nothing was parameterized either: `date_time` and `create_system` are the
exact axes the promotion exists to pin, so making them arguments would dissolve the contract while
appearing to generalize it.

SCOPE OF THE CLAIM (unchanged by the move). Full-file byte identity holds for a fixed
(python-pptx, zlib) pair; `part_digest` is the implementation-independent, cross-environment
assertion. The paragraphs below state this in full and are the canonical explanation in this repo.

THE WRITER READS THE REVIEW GATE AND CANNOT WRITE IT (D-01; the product's hardest rule).
`render_surface_pptx_bytes` **reads** `surface.is_published` — a computed property over
`review.state` — to decide the Draft watermark and the `cp:contentStatus` value. There is no
assignment to `surface.review`, to `surface.review.state`, or to ANY `Surface` field anywhere in
this module. The gate state is mirrored outward into the deck, never back into the record: there is
no read-then-fix path, and therefore no auto-publish path. The durable guard is the absence of the
code path plus `tests/test_pptx_writer.py::test_render_does_not_touch_the_gate` (a Draft Surface is
`model_dump()`-identical after a render) and the standing `git diff --exit-code --
src/newsletters/semantic.py` gate.

THE MARKER IS PROVENANCE, NOT AUTHENTICITY (threat T-02-07, accepted and recorded). `cp:category` is
operator-editable. Nobody may later treat an unmarked deck as proof it was not generated, or a
marked one as proof it was.

THIS MODULE CONSTRUCTS NO XML PARSER (threat T-02-03). python-pptx already parses with
``resolve_entities=False``; `copy.deepcopy` and ``getparent().remove()`` operate on ALREADY-PARSED
trees. Do not add a parse of an XML string here — that would reintroduce the XXE surface upstream
closed.

AI-OPTIONAL / BARE-INSTALL DISCIPLINE. Everything at module level here is **stdlib only**
(`copy`, `hashlib`, `io`, `pathlib`, `zipfile` — the `Surface` annotation is under `TYPE_CHECKING`
and never imported at runtime), so `newsletters.pptx_writer` imports on a bare `pip install .` with
no `[pptx]` extra — which is what lets the duplicate-member and idempotence contracts run on the
bare-install CI job. The writer obtains python-pptx **lazily**, inside its render
function, through the existing boundary `newsletters.adapters._pptx_loader._load_pptx()` — it does
not re-implement that boundary and it does not widen it. There must therefore NEVER be a column-0
``import pptx`` / ``from pptx ...`` line in this file. That is not a convention; it is enforced by
`tests/test_ai_optional.py::test_pptx_writer_has_no_toplevel_pptx_import` (column-0 match) and
`tests/test_ai_optional.py::test_pptx_writer_imports_without_pptx` (a real subprocess import with
`pptx` blocked on `sys.meta_path`). This module is deliberately NOT re-exported from
`newsletters/__init__.py`: that module imports its members eagerly, so widening it would widen the
bare-install blast radius for zero benefit (the `adapters` precedent — callers say
``from newsletters.pptx_writer import ...``).

--- the promoted contract, verbatim from the Phase 1 spike ---------------------------------------

WHY this exists. `python-pptx`'s `Presentation.save()` is deterministic in every respect the
ROADMAP worried about — part traversal order, rId allocation, media dedup, core properties — with
exactly ONE exception: `_ZipPkgWriter.write()` calls
``self._zipf.writestr(pack_uri.membername, blob)`` with a **str** arcname
(``pptx/opc/serialized.py:234-242``). When `zipfile` is handed a string rather than a `ZipInfo`, it
manufactures the `ZipInfo` itself and stamps ``date_time=time.localtime()[:6]``. So two saves of
IDENTICAL content, taken more than two seconds apart (DOS timestamps have 2-second granularity),
produce different bytes — while every unzipped part stays byte-identical. There is no python-pptx
API to override the arcname; post-processing the archive is the only API-level route (monkeypatching
the private `_ZipPkgWriter` was considered and rejected: it breaks silently on upgrade).

THE FIX. `normalize_opc_zip` rewrites the package with FIXED zip metadata — every entry gets a fresh
``ZipInfo(filename, date_time=DOS_EPOCH)``, ``compress_type=ZIP_DEFLATED``, ``create_system=0``
(0 = MS-DOS; the stdlib default is platform-dependent, 3 on unix) — while preserving each entry's
``external_attr`` and, deliberately, the **emitted entry order**. Part BYTES are never touched: no
XML is reformatted, no whitespace is normalized, no attribute is reordered. A determinism definition
that normalized XML would be unfalsifiable — a renderer bug that reordered attributes would pass.

WHY PRESERVE ORDER RATHER THAN SORT. python-pptx's emitted order is already deterministic (measured:
identical name order across time-separated writes) and it puts ``[Content_Types].xml`` first *by
construction*, which the OPC convention expects. Sorting happens to keep it first too (``[`` = 0x5B
sorts before ``_`` and all lowercase) — but only incidentally, which would be a second guarantee to
defend for no gain.

SCOPE OF THE BYTE-IDENTITY CLAIM. DEFLATE output is **implementation-dependent**: reproducible-builds
documented Fedora 40 (zlib-ng) and Debian (zlib) emitting different compressed streams — and
different CRCs — for identical input. Full-file byte identity therefore holds for a fixed
(python-pptx, zlib) pair, which is exactly the scope of an in-process double-render assertion. It is
NOT safe to assert across environments. That is what `part_digest` is for: sorted
``(part name, sha256(part bytes))`` rows are implementation-independent and strictly stronger than
"the zips look the same", because they address the CONTENT and ignore the container entirely. A
committed-artifact-vs-fresh-render gate must use `part_digest`; a double-render test may use bytes.

SECURITY PROPERTY (threat T-01-06, zip-slip). The normalizer NEVER joins an archive member name to
the filesystem. It reads with ``ZipFile.read(name)`` and writes with ``ZipFile.writestr(ZipInfo,
data)``, both entirely in memory — no ``open()``, no ``Path`` join, no ``extract``/``extractall``.
A malicious ``../../etc/passwd`` member in an operator-supplied template is carried through as an
opaque *name string* and cannot reach disk through this function. Phase 2 inherits that property by
reusing this module rather than rediscovering it.

DUPLICATE MEMBER NAMES ARE REFUSED, LOUDLY. The ZIP format permits multiple entries under one name,
and ``ZipFile.read(name)`` resolves to the LAST one — so a by-name pass over a shadowed archive
would silently rewrite part bytes (violating the contract above) and two archives with different
first-entry content would share a `part_digest` (a collision in the Phase-4 trust assertion).
Duplicate-name shadowing is a classic archive-smuggling trick precisely because different consumers
pick different entries. python-pptx never emits duplicates, so the only way one arrives is a
hand-crafted or malicious archive — every function here raises ``ValueError`` naming the duplicated
members rather than picking a winner. This mirrors the template contract's duplicate-SHAPE-name
raise: refuse the ambiguity, never resolve it silently.

Stdlib only (``hashlib``, ``io``, ``zipfile``) — this module must NOT import `pptx`, so it stays
importable on a bare install and needs no skip guard. `tests/test_pptx_determinism.py` carries the
`pptx` import behind the standard `pytestmark` skip.
"""

from __future__ import annotations

import copy
import hashlib
import io
import zipfile
from collections.abc import Iterator, Mapping, Sequence
from pathlib import Path
from typing import TYPE_CHECKING, Any, Union

if TYPE_CHECKING:  # pragma: no cover - annotation only; never imported at runtime
    from .semantic import Surface

__all__ = [
    "DOS_EPOCH",
    "normalize_opc_zip",
    "part_digest",
    "differing_parts",
    "differing_zipinfo_fields",
    "SLOT_PREFIX",
    "WATERMARK_NAME",
    "MARKER",
    "DRAFT_STATUS",
    "WATERMARK_TEXT",
    "bind_slots",
    "fill_slot",
    "render_surface_pptx_bytes",
    "render_surface_pptx",
]

# The earliest timestamp the DOS date format used inside a ZIP can represent. Any fixed instant
# would do; 1980-01-01 is the reproducible-builds convention and is unmistakably "not a real time".
DOS_EPOCH = (1980, 1, 1, 0, 0, 0)

# The ZipInfo fields compared by `differing_zipinfo_fields`. `date_time` is the one python-pptx
# leaves unstable; the rest are pinned so a future drift in ANY of them is reported, not hidden.
_COMPARED_FIELDS = (
    "date_time",
    "compress_type",
    "create_system",
    "external_attr",
    "CRC",
)

# --- the template contract's shared strings (decision note: "The template contract") ------------
#
# These are module constants for the same anti-drift reason `_pptx_loader.MISSING_PPTX_MESSAGE` is
# one: the writer, its tests and the fixture authors must assert against ONE spelling of each, not
# against a string literal copied into four files that silently diverge.

# The reserved shape-name prefix: only shapes named `NL_*` are renderer slots, so an operator's
# logo, footer and page number are not mistaken for unfilled slots (decision note, D-03).
SLOT_PREFIX = "NL_"

# The watermark shape's name — owned by the renderer, refused if the template already defines it
# (measured W14: python-pptx writes two shapes with one name and raises nothing).
WATERMARK_NAME = "NL_DRAFT_WATERMARK"

# The generated-by marker, written to `cp:category` (decision note: provenance, NOT authenticity —
# `cp:category` is operator-editable and an unmarked deck proves nothing).
MARKER = "generated-by:newsletters"

# The review-gate state written to `cp:contentStatus` while the Surface is not published (P-04:
# implemented verbatim per the decision note — `"draft"` if not published else `""`).
DRAFT_STATUS = "draft"

# The watermark's text. A fixed literal, like every other watermark property, so the watermark
# contributes nothing to non-determinism.
WATERMARK_TEXT = "DRAFT"


def _reject_duplicate_member_names(archive: zipfile.ZipFile) -> None:
    """Raise ``ValueError`` if the archive shadows any member name (see the module docstring).

    Every by-name read below is only unambiguous because this ran first: ``ZipFile.read(name)``
    silently resolves a duplicated name to its LAST entry, which would rewrite part bytes in
    `normalize_opc_zip` and collide `part_digest` — the two contracts this module exists to keep.
    """
    names = [info.filename for info in archive.infolist()]
    duplicated = sorted({name for name in names if names.count(name) > 1})
    if duplicated:
        raise ValueError(
            f"duplicate member names in archive: {duplicated!r} — ZIP permits shadowed "
            "entries and ZipFile.read(name) silently picks the last one, so a by-name pass "
            "would rewrite part bytes and collide part_digest. Refusing to touch this "
            "archive; if it is an operator-supplied template, rebuild it without duplicate "
            "entries."
        )


def normalize_opc_zip(raw: bytes) -> bytes:
    """Rewrite an OPC package with FIXED zip metadata. Part BYTES are untouched.

    Fully in memory: no member name ever reaches the filesystem (see the module docstring's
    zip-slip note). Idempotent — ``normalize_opc_zip(normalize_opc_zip(x)) == normalize_opc_zip(x)``.
    Raises ``ValueError`` on duplicate member names rather than silently keeping the last entry.
    """
    with zipfile.ZipFile(io.BytesIO(raw)) as zin:
        _reject_duplicate_member_names(zin)
        infos = zin.infolist()
        # preserve the EMITTED entry order — it is already deterministic
        names = [i.filename for i in infos]
        attrs = {i.filename: i.external_attr for i in infos}
        data = {n: zin.read(n) for n in names}

    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", compression=zipfile.ZIP_DEFLATED) as zout:
        for name in names:
            info = zipfile.ZipInfo(filename=name, date_time=DOS_EPOCH)
            info.compress_type = zipfile.ZIP_DEFLATED
            # 0 = MS-DOS; the stdlib default is platform-dependent (3 on unix)
            info.create_system = 0
            info.external_attr = attrs[name]  # preserve python-pptx's value verbatim
            zout.writestr(info, data[name])
    return buf.getvalue()


def part_digest(raw: bytes) -> str:
    """CONTENT identity of an OPC package: sha256 over sorted ``(name, sha256(part bytes))`` rows.

    No zip metadata, no XML normalization. This is the implementation-INDEPENDENT assertion: it is
    unaffected by which zlib built the archive, by entry order, and by every timestamp — while still
    failing loudly if one byte of one part changed. Raises ``ValueError`` on duplicate member
    names: a digest that silently read only the LAST of two shadowed entries would report two
    different-content archives as identical — the exact spoof the Phase-4 trust gate must catch.

    THE ROW ENCODING IS LENGTH-PREFIXED, AND THAT IS LOAD-BEARING (Phase-2 review WR-01). Member
    names are attacker-controlled bytes — a ZIP name may legally contain NUL and newline — so a
    delimiter-based row encoding (``name + b"\\0" + hex + b"\\n"``) is NOT injective: a single
    crafted member named ``"a\\0" + <64 hex> + "\\nb"`` can serialize to the same byte stream as a
    two-member archive, colliding the digest for archives with entirely different part content.
    Prefixing each name with its 8-byte big-endian length makes the row boundaries unambiguous for
    EVERY possible name, closing that construction. Do not "simplify" this back to delimiters.
    """
    with zipfile.ZipFile(io.BytesIO(raw)) as archive:
        _reject_duplicate_member_names(archive)
        rows = sorted(
            (name, hashlib.sha256(archive.read(name)).hexdigest())
            for name in archive.namelist()
        )
    digest = hashlib.sha256()
    for name, part in rows:
        encoded = name.encode("utf-8")
        digest.update(len(encoded).to_bytes(8, "big"))
        digest.update(encoded)
        digest.update(part.encode("ascii"))
    return digest.hexdigest()


def differing_parts(a: bytes, b: bytes) -> list[str]:
    """Sorted names of parts whose unzipped BYTES differ between two packages.

    A name present in only one package is reported too — a part that appeared or vanished is a
    content difference, not a metadata one. Raises ``ValueError`` if either package shadows a
    member name (a by-name comparison over duplicates would compare only the LAST entries).
    """
    with zipfile.ZipFile(io.BytesIO(a)) as za, zipfile.ZipFile(io.BytesIO(b)) as zb:
        _reject_duplicate_member_names(za)
        _reject_duplicate_member_names(zb)
        names_a, names_b = set(za.namelist()), set(zb.namelist())
        differing = set(names_a ^ names_b)
        for name in names_a & names_b:
            if za.read(name) != zb.read(name):
                differing.add(name)
    return sorted(differing)


def differing_zipinfo_fields(a: bytes, b: bytes) -> list[str]:
    """Sorted names of the ZIP METADATA fields that differ between two packages.

    Compares ``date_time``, ``compress_type``, ``create_system``, ``external_attr`` and ``CRC``
    entry by entry, plus ``filename`` order (reported as ``"filename"`` when the emitted entry-name
    sequences differ). Returns ``[]`` when the two archives carry identical metadata — and
    ``["date_time"]`` for a raw python-pptx double write across a DOS-time boundary, which is the
    measured shape of the ONE non-determinism. Raises ``ValueError`` if either package shadows a
    member name (the by-name pairing below is last-wins over duplicates).
    """
    with zipfile.ZipFile(io.BytesIO(a)) as za, zipfile.ZipFile(io.BytesIO(b)) as zb:
        _reject_duplicate_member_names(za)
        _reject_duplicate_member_names(zb)
        infos_a, infos_b = za.infolist(), zb.infolist()
        names_a = [i.filename for i in infos_a]
        names_b = [i.filename for i in infos_b]
        differing: set[str] = set()
        if names_a != names_b:
            differing.add("filename")
        by_name_b = {i.filename: i for i in infos_b}
        for info_a in infos_a:
            info_b = by_name_b.get(info_a.filename)
            if info_b is None:
                continue  # a missing entry is a CONTENT difference — differing_parts reports it
            for field in _COMPARED_FIELDS:
                if getattr(info_a, field) != getattr(info_b, field):
                    differing.add(field)
    return sorted(differing)


# ==================================================================================================
# EVERYTHING ABOVE THIS LINE IS STDLIB AND BARE-IMPORTABLE.
# EVERYTHING BELOW REQUIRES THE OPTIONAL ``[pptx]`` EXTRA **AT CALL TIME** — never at import time.
#
# The functions below take/return python-pptx objects, but they never import `pptx` at column 0.
# The one lazy import per render lives inside `bind_slots` (the `MSO_SHAPE_TYPE` recursion
# predicate); `render_surface_pptx_bytes` obtains the `pptx` module itself through the EXISTING
# boundary `newsletters.adapters._pptx_loader._load_pptx()` rather than re-implementing it. Keep it
# that way: `tests/test_ai_optional.py::test_pptx_writer_imports_without_pptx` imports this module
# in a subprocess with `pptx` blocked on `sys.meta_path`, and it must keep passing.
#
# python-pptx `Presentation` and shape objects are typed `Any` throughout: python-pptx ships no
# stubs and this repo deliberately does not add a `types-*` package (the `_pptx_loader.py`
# precedent). Everything else here is fully annotated.
# ==================================================================================================


def _walk(shapes: Any, group_type: Any) -> Iterator[Any]:
    """Every shape in a shapes collection, DESCENDING into group shapes.

    The recursion is load-bearing, not defensive. `slide.shapes` is **top-level only** (measured,
    02-RESEARCH W17, and re-proved in this repo by
    `tests/test_pptx_writer.py::test_group_nesting_hides_a_slot_from_slide_shapes`): a textbox named
    `NL_INSIDE_GROUP` nested inside a group is simply absent from it. So a flat
    ``{shape.name: shape for shape in slide.shapes}`` comprehension is BLIND to a slot the operator
    can see in their Selection Pane — content bound to that name would be rejected as *unknown*, and
    the unfilled-slot refusal would never fire for it. Grouping boxes is ordinary PowerPoint
    hygiene, so this is not an exotic case.

    `group_type` is `MSO_SHAPE_TYPE.GROUP`, passed in rather than imported here so the predicate is
    MEASURED (an enum comparison) rather than duck-typed, and so there is exactly ONE lazy
    python-pptx import per render — in `bind_slots`, its only caller.
    """
    for shape in shapes:
        yield shape
        if shape.shape_type == group_type:
            yield from _walk(shape.shapes, group_type)


def bind_slots(
    prs: Any, slots: Mapping[str, Union[str, Sequence[str]]]
) -> dict[str, Any]:
    """Map every shape name in the deck to its shape, refusing every ambiguity rather than guessing.

    Binds over ``slide.shapes`` (recursively — see :func:`_walk`) and NOT ``slide.placeholders``:
    the latter holds only real placeholders keyed by ``idx``, and operator-added textboxes and
    pictures are absent from it entirely (decision note, "The template contract").

    Five refusals, in this order, each a ``ValueError`` in the house three-part teaching voice —
    what was found, why it cannot be resolved silently, what the operator does next:

    1. a **duplicate `NL_`-prefixed shape name** (legal in OOXML; copy-paste is how an operator
       makes one — a last-wins map would silently DROP a slot, the exact failure D-03 exists to
       prevent). SCOPED TO THE RESERVED PREFIX ON PURPOSE (Phase-2 review WR-02, reproduced live):
       PowerPoint and python-pptx auto-name shapes PER SLIDE ("Title 1", "TextBox 1"), so a
       deck-wide refusal over ALL shapes rejects every ordinary multi-slide deck with default
       names. Unprefixed shapes are not renderer slots (only ``NL_*`` names are — the SLOT_PREFIX
       contract), so a collision between two decorative shapes drops nothing — the map keeps the
       first-seen shape, deterministically, and the operator's business stays the operator's;
    2. the template already defining **WATERMARK_NAME** (the renderer owns that name and adds the
       watermark itself; leaving it in the template writes two shapes under one name — measured,
       W14, with no error — and defeats this map on the next render);
    3. an **unknown content name** — content bound to a name the template does not contain;
    4. an **unfilled reserved slot** — an ``NL_``-prefixed shape with no matching content (a
       reserved-prefix slot left empty would ship a blank box to a reader). The ``NL_`` prefix is
       what keeps this direction usable on a real deck: without it every operator logo, footer and
       page number would be rejected;
    5. a **slot that cannot hold text** — a group or a picture reports ``has_text_frame == False``,
       so filling it would raise ``AttributeError``: a stack trace instead of a teaching error.

    Returns the name→shape map. Mutates nothing.
    """
    # Lazy on purpose: `pptx` is the optional [pptx] extra and this module must stay importable on
    # a bare install (see the banner above). One import per render, here, for `_walk`'s predicate.
    # (No `# type: ignore[import-untyped]` needed on this one — mypy resolves the enum module. The
    # trailing comment is kept to ONE pragma so the line clears isort's width; see DEF-15.)
    from pptx.enum.shapes import MSO_SHAPE_TYPE  # noqa: PLC0415

    by_name: dict[str, Any] = {}
    for slide in prs.slides:
        for shape in _walk(slide.shapes, MSO_SHAPE_TYPE.GROUP):
            if shape.name in by_name and shape.name.startswith(SLOT_PREFIX):
                raise ValueError(
                    f"template has two shapes named {shape.name!r} — the name->shape binding is "
                    "ambiguous. Duplicate shape names are legal in OOXML and copy-pasting a box in "
                    "PowerPoint is how one is made, so a last-wins map would silently DROP a slot. "
                    "Rename one in PowerPoint's Selection Pane (Alt+F10). Refusing to guess which "
                    "slot the content belongs in."
                )
            # First-seen wins for UNPREFIXED duplicates — deliberate, not an oversight (WR-02):
            # PowerPoint auto-names shapes per slide ("TextBox 1"), so any real multi-slide deck
            # carries unprefixed collisions. Unprefixed shapes are not renderer slots (the
            # SLOT_PREFIX contract), so keeping the first is inconsequential and deterministic;
            # raising here rejected every ordinary operator deck.
            by_name.setdefault(shape.name, shape)

    if WATERMARK_NAME in by_name:
        raise ValueError(
            f"the template already defines {WATERMARK_NAME!r}. That name is owned by the renderer, "
            "which adds the Draft watermark itself; leaving it in the template would produce two "
            "shapes with one name (measured: python-pptx writes both and raises nothing) and would "
            "defeat the name->shape binding on the next render. Remove it from the template."
        )

    unknown = sorted(set(slots) - set(by_name))
    if unknown:
        raise ValueError(
            f"Surface content is bound to placeholder name(s) {unknown!r} that this template does "
            f"not contain. Named shapes in this template: {sorted(by_name)!r}. The renderer never "
            "invents layout — add the shape to the template (name it in the Selection Pane, "
            "Alt+F10) or fix the name in the content mapping."
        )

    unfilled = sorted(
        name for name in by_name if name.startswith(SLOT_PREFIX) and name not in slots
    )
    if unfilled:
        raise ValueError(
            f"template placeholder(s) {unfilled!r} have no matching Surface content. A "
            f"{SLOT_PREFIX!r}-prefixed slot left empty would ship a blank box to a reader — "
            "refusing. Either populate it or remove the reserved prefix from the shape's name "
            "(unprefixed shapes are the operator's: logos, footers, page numbers are never slots)."
        )

    for name in slots:
        shape = by_name[name]
        if not getattr(shape, "has_text_frame", False):
            raise ValueError(
                f"slot {name!r} is a {shape.shape_type} and cannot hold text (measured: group and "
                "picture shapes report has_text_frame == False, so filling one raises "
                "AttributeError rather than anything a reader of the traceback could act on). "
                "Point the slot at a text box, or place this content as an asset."
            )

    return by_name


def fill_slot(text_frame: Any, lines: Union[str, Sequence[str]]) -> None:
    """Fill a template slot with N lines, PRESERVING the operator's paragraph and run formatting.

    A BARE ``str`` IS ONE PARAGRAPH (Phase-2 review WR-03). ``str`` is itself a ``Sequence[str]``
    of its characters, so neither mypy nor any refusal used to catch ``fill_slot(tf, "hello")`` —
    it silently shipped five paragraphs, ``['h', 'e', 'l', 'l', 'o']``, from the single most likely
    caller typo. A string is what an author means by "one line", so it is treated as ``[lines]``
    here (covering both this function and ``render_surface_pptx_bytes``'s ``slots`` mapping) rather
    than refused: the input is not ambiguous, it has exactly one sane reading.

    python-pptx exposes **no bullet API at all** — `_Paragraph`'s entire public surface is
    ``add_line_break, add_run, alignment, clear, font, level, line_spacing, part, runs, space_after,
    space_before, text`` (measured, W6) — so a bullet can only ever be *inherited*, never
    synthesized. That single fact decides the primitive: paragraph 0 of the operator's slot is the
    formatting carrier, and every additional line is a deep copy of it with its run's text replaced.
    The moment the writer starts *building* formatting it has started inventing layout, which is the
    one thing D-03 forbids.

    THE TWO RULES THIS ENCODES — both are the ones a future editor is most tempted to "simplify":

    (a) **Never assign ``text_frame.text`` on an operator's shape.** Measured (W3) on a box authored
        20pt bold with a ``buChar`` bullet: the setter erases the run's ``rPr`` (20pt bold →
        ``None``) AND the paragraph's ``pPr`` (the bullet vanishes), collapsing the frame to one
        paragraph. The deck still renders — with the operator's formatting silently downgraded to
        theme defaults. That is a visual regression NO automated check in this repo can see except
        the fidelity tests in `tests/test_pptx_writer.py`. ``tf.clear()`` + ``add_paragraph()`` is
        only half-safe and is therefore worse: it keeps p0's ``pPr`` but loses every run's ``rPr``,
        and added paragraphs get no ``pPr`` at all (W4). A half-preserved deck is harder to notice
        than an obviously-wrong one.

    (b) **Never re-apply enumerated font properties after a ``tf.text`` assignment.** That copies
        only the properties somebody thought to list; the deep copy carries the WHOLE ``rPr``/``pPr``
        subtree, including theme references and language tags nobody would think to enumerate.

    STANDING BAN — ``TextFrame.fit_text()`` must NEVER be used anywhere in this module (T-02-15).
    Its ``_best_fit_font_size`` path locates a TrueType font file *installed on the current system*
    (W7), which would make the output machine-dependent: the single most tempting non-determinism in
    the library. Overflow is disclosed to a human (02-RESEARCH Pitfall 3), never "fixed" by shrinking
    the font — faithful, not suggestive.

    The lxml idiom ``el.getparent().remove(el)`` prunes surplus paragraphs and runs. It operates on
    an ALREADY-PARSED tree and constructs no XML parser, so python-pptx's upstream
    ``resolve_entities=False`` mitigation is not weakened (T-02-03). ``copy`` is stdlib, so this
    function adds nothing to the module's bare-install surface.

    Raises:
        ValueError: if ``lines`` is empty (an empty slot ships a blank box to a reader), or if
            paragraph 0 carries no run (nothing to inherit — Pitfall 6, an OBSERVED failure: a text
            box with no typed characters really does have a paragraph with zero runs).
    """
    if isinstance(lines, str):
        # One paragraph, not one per character — see the docstring (WR-03).
        lines = [lines]
    if not lines:
        raise ValueError(
            "refusing to fill a slot with no lines — an empty slot ships a blank box to a reader, "
            "which is the same silent gap the unfilled-slot refusal exists to prevent. Either "
            "supply at least one line or do not bind content to this slot."
        )

    p0 = text_frame.paragraphs[0]
    for extra in list(text_frame.paragraphs)[1:]:
        extra._p.getparent().remove(extra._p)

    runs = p0.runs
    if not runs:
        raise ValueError(
            "template slot's first paragraph carries no run, so there is no formatting to inherit "
            "(measured: a text box with no typed characters has a paragraph with ZERO runs). The "
            "renderer will not invent a font, size or bullet to put in its place. Type one "
            "character of placeholder text into the shape in PowerPoint and re-save the template."
        )
    for surplus in runs[1:]:
        surplus._r.getparent().remove(surplus._r)

    # One deep copy of the carrier per EXTRA line, chained with `addnext` so document order matches
    # `lines` order. The clone carries the whole rPr/pPr subtree — that is the entire point.
    previous = p0._p
    for _ in lines[1:]:
        clone = copy.deepcopy(p0._p)
        previous.addnext(clone)
        previous = clone

    for paragraph, line in zip(text_frame.paragraphs, lines):
        paragraph.runs[0].text = line


def _add_watermark(slide: Any, pptx: Any) -> Any:
    """Add the Draft watermark textbox to one slide, LAST so it sits at the top of z-order.

    Every property is a fixed literal — no clock, no randomness — so the watermark contributes
    nothing to non-determinism (measured W12: normalized bytes and `part_digest` equal across a
    3-second gap, and it reads back last in `slide.shapes` with `rotation == 315.0`).

    NOTE THE DELIBERATE CONTRAST WITH :func:`fill_slot`: assigning ``text_frame.text`` is CORRECT
    here. The writer created this shape and owns every property on it, so there is no operator
    formatting to preserve. Use the cheap API where you own the shape; use the preserving primitive
    where the operator does. That is the rule, not an inconsistency.

    TWO RECORDED NON-CHOICES, so nobody reopens them:

    (i) **This is not true transparency.** python-pptx has no alpha API and a genuinely see-through
        watermark needs raw ``a:alpha`` XML. A light-grey rotated run is visible, deterministic and
        API-supported, and is not worth trading for that fragility.
    (ii) **The watermark is ADDED, never toggled off an operator-supplied element.** Removal was
        measured deterministic and exactly reversible (W13), so the decision note's *mechanical*
        objection is disproved — but the binding objection stands: a toggle would make correct gate
        behaviour depend on operator template content, and "the watermark is missing" would become a
        template bug instead of a renderer bug. Adding is unconditional in mechanism and conditional
        only on the gate.
    """
    box = slide.shapes.add_textbox(
        pptx.util.Inches(1.5),
        pptx.util.Inches(2.5),
        pptx.util.Inches(7),
        pptx.util.Inches(2),
    )
    box.name = WATERMARK_NAME
    box.rotation = 315.0
    box.text_frame.text = WATERMARK_TEXT
    run = box.text_frame.paragraphs[0].runs[0]
    run.font.size = pptx.util.Pt(96)
    run.font.bold = True
    run.font.color.rgb = pptx.dml.color.RGBColor(0xD0, 0xD0, 0xD0)
    return box


def render_surface_pptx_bytes(
    surface: "Surface",
    *,
    template: Union[str, Path],
    slots: Mapping[str, Union[str, Sequence[str]]],
) -> bytes:
    """Render a Surface into deterministic `.pptx` bytes through an operator-supplied template.

    Fills the operator's EXISTING slides — it never calls ``add_slide`` (which regenerates
    placeholder names and would reject every one of the operator's names on first contact). The
    Surface→`slots` derivation belongs to the composer (Phase 3): only it knows which authored block
    belongs in which Selection-Pane name, so `slots` is a required keyword argument here (P-03).

    THE GATE IS READ, NEVER WRITTEN. ``surface.is_published`` decides the watermark and
    ``cp:contentStatus``; nothing here assigns to any `Surface` field (see the module docstring).

    Returns normalized bytes: ``prs.save(BytesIO)`` is the ONLY clock in this function, and
    :func:`normalize_opc_zip` removes its trace. Nothing is written to disk — see
    :func:`render_surface_pptx` for that.
    """
    # Lazy on purpose ([pptx] extra). We call the EXISTING boundary rather than writing a second
    # `try: import pptx`: `_load_pptx` already raises the teaching ImportError naming
    # `pip install '.[pptx]'` and is already covered by
    # `test_pptx_loader_raises_teaching_error_without_pptx`. A second implementation of "import the
    # extra safely" drifts from the first exactly as a second normalizer would.
    from .adapters._pptx_loader import _load_pptx  # noqa: PLC0415
    from .adapters._timestamps import EPOCH_ZERO  # noqa: PLC0415

    pptx = _load_pptx()
    prs = pptx.Presentation(str(template))

    by_name = bind_slots(prs, slots)
    # Fill order does not affect the output bytes (measured, W18: text edits are in-place on
    # existing elements), so iterate `slots` directly — sorting would imply a guarantee we do not
    # need and would hide the real ordering variable, which is image ADD order (Phase 3's problem).
    for name, lines in slots.items():
        fill_slot(by_name[name].text_frame, lines)

    if not surface.is_published:
        # Every slide, never `prs.slides[0]`: a two-slide deck with one watermark is a deck whose
        # second page reads as approved.
        for slide in prs.slides:
            _add_watermark(slide, pptx)

    core = prs.core_properties
    core.category = MARKER
    # P-04, the decision note's mapping VERBATIM. `ReviewState` has three members, so an IN_REVIEW
    # deck is also labelled "draft" — deliberately not deviated from here; the tri-state amendment
    # is raised in the PR body rather than changed silently.
    core.content_status = "" if surface.is_published else DRAFT_STATUS
    # Serializes as `dc:identifier`, NOT `cp:identifier` (W20 — a wording correction to the
    # decision note's field table, with no code consequence). Makes a deck found on somebody's
    # desktop self-locating: you can get from the file back to the reviewed record it came from.
    core.identifier = surface.id
    # `dcterms:created` serializes as W3CDTF but python-pptx reads it back tz-NAIVE, while the
    # repo-wide `EPOCH_ZERO` sentinel is tz-aware UTC — a plain `==` fails (Pitfall 8). Strip the tz
    # AT THE OPC BOUNDARY rather than minting a second epoch constant.
    core.created = EPOCH_ZERO.replace(tzinfo=None)
    core.modified = EPOCH_ZERO.replace(tzinfo=None)

    buf = io.BytesIO()
    prs.save(buf)  # the ONLY clock in this function
    return normalize_opc_zip(buf.getvalue())


def render_surface_pptx(
    surface: "Surface",
    *,
    template: Union[str, Path],
    slots: Mapping[str, Union[str, Sequence[str]]],
    out_path: Union[str, Path],
) -> Path:
    """Render a Surface to a `.pptx` on disk in ONE write of complete, already-normalized bytes.

    Never ``prs.save(path)`` followed by a rewrite in place: a raised normalization would leave an
    un-normalized deck on disk, and the next reader would have no way to tell.

    SECURITY (threat T-02-08, path traversal): `out_path` is **caller-supplied and never derived
    from Surface content**. Joining `surface.id` — authored data — to a directory would be a
    path-traversal primitive. The caller owns the path; this function only writes to it.
    """
    raw = render_surface_pptx_bytes(surface, template=template, slots=slots)
    out = Path(out_path)
    out.write_bytes(raw)
    return out
