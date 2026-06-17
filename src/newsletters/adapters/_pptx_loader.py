"""Lazy python-pptx boundary — the optional ``[pptx]`` extra, imported only on use (ADAPT-04).

WHY THIS MODULE EXISTS (CONTEXT decision 1 / RESEARCH, threat T-06-03 + T-06-05):
``python-pptx`` is a third-party dependency that backs the PowerPoint adapter (06-03). It is NOT
AI, so the ``forbid-ai-in-core`` import-linter contract is unaffected — but the AI-optional /
minimal-core invariant still demands that a bare ``pip install .`` (no ``[pptx]``) can
``import newsletters``, ``import newsletters.adapters``, and run the deterministic spine with zero
python-pptx (and zero of its transitive C-extensions lxml/Pillow/XlsxWriter). We achieve that by
mirroring the ``[ai]`` / ``[excel]`` lazy-import discipline: **python-pptx is never imported at
module top-level anywhere reachable from ``import newsletters``.** It is imported INSIDE
:func:`_load_pptx`, which the PPTX adapter (Plan 06-03) calls only when it actually parses a
``.pptx``. Absence raises a clear teaching :class:`ImportError` pointing at ``pip install '.[pptx]'``.

This module itself must therefore have NO top-level ``import pptx`` / ``from pptx ...`` — the
bare-install gate (``tests/test_ai_optional.py``) asserts column-0 import count 0 for those edges.
"""

from __future__ import annotations

import io
from typing import Any, Union

# NOTE on typing: python-pptx ships NO inline type stubs and there is no maintained ``types-*``
# stub package. python-pptx is the ONLY new dependency this phase adds (CONTEXT decision 1), so we
# do NOT pull a stub package just for annotations — the ``Presentation`` objects are typed as
# ``Any`` (mypy then treats them opaquely, which is correct: the faithful per-shape extraction
# logic lives in the 06-03 adapter, not here). This mirrors the Excel ``_openpyxl_loader`` decision.
Presentation = Any  # alias used in return annotations below; kept ``Any`` to avoid a stub dependency

# The exact teaching message a user sees if they reach the PowerPoint adapter without the extra.
# Exposed as a module constant so tests can assert against it without string-duplication drift.
MISSING_PPTX_MESSAGE = (
    "The PowerPoint adapter requires the optional 'python-pptx' dependency. "
    "Install it with: pip install '.[pptx]'  (or: pip install newsletters[pptx]). "
    "The deterministic spine runs without it — python-pptx is needed only for .pptx extraction "
    "(AI-optional / minimal-core: third-party adapter deps live behind extras)."
)


def _load_pptx() -> Any:
    """Import and return the ``pptx`` module, lazily.

    The import lives INSIDE this function (never at module top) so that importing
    ``newsletters`` / ``newsletters.adapters`` never requires the ``[pptx]`` extra. If python-pptx
    is not installed, re-raise a teaching :class:`ImportError` naming the extra and the install
    command, and stating the spine runs without it.

    Returns:
        The imported ``pptx`` module.

    Raises:
        ImportError: if ``python-pptx`` is not installed, with an actionable message.
    """
    try:
        # python-pptx ships no stubs; we deliberately do NOT add a types-* package (python-pptx is
        # the ONLY new dep this phase adds) -> ignore the missing-stub error. PLC0415: lazy on
        # purpose (optional [pptx] extra, T-06-03).
        import pptx  # type: ignore[import-untyped]  # noqa: PLC0415
    except ImportError as exc:
        raise ImportError(MISSING_PPTX_MESSAGE) from exc
    return pptx


def load_presentation(raw: Union[str, bytes, "io.BytesIO"]) -> "Presentation":
    """Open a ``.pptx`` and return a python-pptx ``Presentation``.

    The PPTX adapter (Plan 06-03) consumes the returned ``Presentation`` to walk slides, shapes,
    tables and notes. python-pptx's ``Presentation()`` accepts either a filesystem path or a
    file-like object; raw bytes are wrapped in a fresh :class:`io.BytesIO` (a stream is consumed by
    the open, so bytes need their own buffer).

    Args:
        raw: a filesystem path (``str``), raw ``.pptx`` bytes, or a ``BytesIO``.

    Returns:
        A python-pptx ``Presentation`` instance.

    Raises:
        ImportError: if ``python-pptx`` is not installed (via :func:`_load_pptx`).
    """
    pptx = _load_pptx()
    if isinstance(raw, (bytes, bytearray)):
        return pptx.Presentation(io.BytesIO(bytes(raw)))
    return pptx.Presentation(raw)


__all__ = ["_load_pptx", "load_presentation", "MISSING_PPTX_MESSAGE"]
