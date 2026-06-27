r"""The PowerPoint ``.pptx`` adapter (ADAPT-04) — the third file-extraction ``DistillPort`` backend.

One ``.pptx`` deck -> one ``Source`` whose ``transcript`` is a canonical text serialization of the
deck: one line ``Slide N / <shape>\t<text>\n`` per non-empty *extractable* paragraph/cell/note,
slides in deck order, shapes in document order, groups RECURSED. The adapter slices the verbatim
text units out of that transcript and hands them — IN TRANSCRIPT ORDER — to the shared
``normalize()`` (ADAPT-01), which mints content-addressed ``Claim(+Trace)``. The adapter hand-mints
NO hashes; the faithful-extraction rule lives only in ``normalize()`` (CONTEXT decision 2).

WHAT IS EXTRACTED (faithful, per-paragraph atoms). Text frames (titles / body placeholders / text
boxes / auto-shapes with text) emit one verbatim unit per non-empty paragraph (``paragraph.text``,
keeping ``\v`` soft breaks VERBATIM). Tables emit one verbatim unit per non-empty cell, row-major.
Speaker notes emit per paragraph with the ``Slide N / notes`` locator. Each emitted text is the SAME
string in BOTH the transcript line and ``units`` — guaranteeing it is a substring ``normalize()``
locates.

ZERO SILENT DROPS (criterion 2). Every LEAF shape is exactly one of: (a) >=1 emitted unit, (b) a
skipped-empty (an empty text frame / a placeholder with no readable payload — like Excel's blank
cell), or (c) one ``unextracted[]`` entry. A GROUP node itself contributes nothing — it is a
container, recursed into, never reported. The unreadable taxonomy (SmartArt, chart, picture, media,
OLE, and any unknown graphic frame) routes to ``unextracted[]`` and is never silently lost. A
malformed/unreadable ``.pptx`` is caught and disclosed as a whole-source ``unextracted[]`` entry
with an empty transcript, never an unhandled crash (V5 input validation).

SMARTART (the silent-drop trap, L2). python-pptx exposes NO high-level SmartArt API; a SmartArt
diagram is a ``graphicFrame`` whose ``shape_type`` is ``None``, ``has_table``/``has_chart`` both
``False``. It is identified ONLY by its ``a:graphicData/@uri`` being the diagram namespace. We read
that URI via the 06-02-confirmed accessor ``shape._element.graphicData_uri`` with an lxml
``<a:graphicData>/@uri`` fallback, so a future python-pptx rename cannot crash us — worst case an
unknown URI is reported as "unknown graphic frame", NEVER dropped.

HARD RULE (mirrors ``manual.py`` / ``excel_adapter.py``): this returns truth only. It NEVER calls
``Surface.publish``, sets ``ReviewState.PUBLISHED``, or builds a published ``Review``.

SECURITY (V5 / T-06-06/07/08). ``.pptx`` is UNTRUSTED input. ``Presentation()`` is wrapped in
try/except -> a whole-source disclosure on failure (no crash, no decompression amplification). The
adapter reads only shape TYPE and text — it NEVER reads embedded image/OLE/media BYTES, so no code
path touches a macro/OLE payload and there is no zip-bomb decompression of blobs. No network, no
external-link following.

LAZY PYTHON-PPTX. python-pptx is the optional ``[pptx]`` extra. It is imported ONLY through
``_pptx_loader`` (never at module top) and, for the enums/namespace helpers, INSIDE the functions
that need them — so ``import newsletters`` / ``import newsletters.adapters`` (which
``adapters/__init__.py`` runs to ``register()`` this adapter) stays extra-free. The deterministic
spine runs with zero python-pptx (AI-optional / minimal-core).
"""

from __future__ import annotations

import re
from datetime import datetime
from typing import Any

from ..distill.coverage import Coverage, Unextracted
from ..distill.ports import DistillationResult
from ..locators import FreeLocator
from ..semantic import Distillation, Source
from ._coverage_codec import (
    decode_coverage,
    encode_coverage,
    not_reconstructable_marker,
)
from ._pptx_loader import load_presentation
from ._timestamps import deterministic_timestamp
from .normalize import normalize

# The transcript separator (mirror Excel; RESEARCH "Transcript layout"). It only needs to be
# unambiguous for a HUMAN reader of the transcript — ``normalize()`` does pure substring matching,
# so a paragraph/cell value that itself contains a tab, newline, or ``\v`` is still emitted VERBATIM
# (never escaped) and stays locatable; the transcript line simply spans more than one physical line.
# Escaping the value would break ``claim.text == transcript slice`` (the Phase-3 span gate).
SEP = "\t"

# The OOXML DrawingML diagram namespace — a SmartArt graphicFrame's ``a:graphicData/@uri``.
_DIAGRAM_URI = "http://schemas.openxmlformats.org/drawingml/2006/diagram"

# ---- reason strings (the unextracted[] contract; Plan 04's golden corpus pins these) ----------
# Copied verbatim from 06-RESEARCH.md "unextracted[] taxonomy" so 06-04 can pin them exactly.
_R_CHART = "chart not extracted (chart data is out of scope)"
_R_PICTURE = "image not extracted (picture content out of scope)"
_R_MEDIA = "media (audio/video) not extracted"
_R_OLE = "embedded/linked OLE object not extracted"
_R_SMARTART = "SmartArt/diagram not extracted (no high-level API)"
_R_UNKNOWN_FRAME = "unknown graphic frame ({uri}) not extracted"
_R_UNREADABLE = (
    "presentation could not be read by python-pptx ({error}) — not extractable"
)


def _graphic_data_uri(shape: Any) -> str | None:
    """Return a graphic frame's ``a:graphicData/@uri`` (or ``None``), crash-free.

    Primary accessor is the 06-02-confirmed ``shape._element.graphicData_uri`` (a python-pptx
    internal helper, stable in 1.0.2). The lxml fallback reads ``<a:graphicData>/@uri`` directly so
    a future rename of that helper cannot crash the walk — worst case the caller classifies an
    unknown URI as "unknown graphic frame" (reported, never dropped). Total: any failure -> ``None``.
    """
    uri = getattr(shape._element, "graphicData_uri", None)
    if isinstance(uri, str):
        return uri
    try:
        from pptx.oxml.ns import qn  # noqa: PLC0415 — lazy, optional [pptx] extra

        gd = shape._element.find(".//" + qn("a:graphicData"))
        if gd is not None:
            value = gd.get("uri")
            return value if isinstance(value, str) else None
    except Exception:  # noqa: BLE001 — never crash on a malformed/renamed element; report instead
        return None
    return None


def _classify_unreadable(shape: Any) -> str | None:
    """Classify a leaf shape with no readable text frame/table into a taxonomy reason, or ``None``.

    Returns the verbatim reason string for picture / media / OLE / chart (by ``shape_type``) and for
    SmartArt / unknown graphic frame (by the diagram ``@uri``). Returns ``None`` for a shape genuinely
    NOT in the unreadable taxonomy (an empty placeholder, a connector/line with no payload) — the
    caller then treats it as a skipped-empty (no unit, no drop), mirroring Excel's blank-cell skip.
    """
    # Lazy enum import (optional [pptx] extra) — never at module top.
    from pptx.enum.shapes import MSO_SHAPE_TYPE  # noqa: PLC0415

    st = shape.shape_type
    if st == MSO_SHAPE_TYPE.PICTURE:
        return _R_PICTURE
    if st == MSO_SHAPE_TYPE.MEDIA:
        return _R_MEDIA
    if st in (MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT, MSO_SHAPE_TYPE.LINKED_OLE_OBJECT):
        return _R_OLE
    if st == MSO_SHAPE_TYPE.CHART:
        return _R_CHART
    # SmartArt: shape_type is None on its graphicFrame -> identify by the diagram URI.
    uri = _graphic_data_uri(shape)
    if uri == _DIAGRAM_URI:
        return _R_SMARTART
    if uri is not None:
        # A graphic frame we do not recognise — reported, NEVER dropped (zero-silent-drops, T-06-09).
        return _R_UNKNOWN_FRAME.format(uri=uri)
    return None  # not in the unreadable taxonomy -> skipped-empty


def _drop(loc: str, reason: str) -> Unextracted:
    """Build one ``Unextracted`` with a human-readable ``Slide N / <shape>`` content locator."""
    return Unextracted(locator=FreeLocator(text=loc), reason=reason)


def _emit_or_report(
    shape: Any,
    slide_no: int,
    units: list[str],
    lines: list[str],
    drops: list[Unextracted],
) -> None:
    r"""The leaf decision: emit verbatim text units, or report the shape to ``unextracted[]``.

    The public ladder (06-RESEARCH "Per-shape emit/report decision"): tables -> per-cell verbatim,
    row-major; charts -> report; text frames -> per non-empty paragraph (``paragraph.text``, keeping
    ``\v`` verbatim) appended to BOTH the transcript line and ``units`` (same string -> guaranteed
    substring); an empty text frame is skipped-empty; otherwise ``_classify_unreadable`` decides a
    taxonomy drop or a genuine skipped-empty. A GROUP never reaches here (``_walk`` recurses it).
    """
    loc = f"Slide {slide_no} / {shape.name}"

    # 1) Tables (graphic frame) -> extract every non-empty cell verbatim, row-major.
    if shape.has_table:
        table = shape.table
        for r, row in enumerate(table.rows):
            for c, cell in enumerate(row.cells):
                text = cell.text
                if text:
                    lines.append(f"{loc} [r{r}c{c}]{SEP}{text}\n")
                    units.append(text)
        return

    # 2) Charts -> report (faithful: chart data is structured, not prose; we do not cherry-pick).
    if shape.has_chart:
        drops.append(_drop(loc, _R_CHART))
        return

    # 3) Text frames (text boxes, titles, body placeholders, auto-shapes with text).
    if shape.has_text_frame:
        emitted = False
        for para in shape.text_frame.paragraphs:
            text = para.text  # \v for soft breaks (kept VERBATIM); paragraphs are ordered
            if text:
                lines.append(f"{loc}{SEP}{text}\n")
                units.append(text)
                emitted = True
        if emitted:
            return
        # An empty text frame is a skipped-empty (like a blank cell), NOT a drop. Fall through in
        # case the shape also matches the unreadable taxonomy (it will not for a plain text box).

    # 4) Unreadable taxonomy (no readable text frame / table) — or a genuine skipped-empty.
    reason = _classify_unreadable(shape)
    if reason is not None:
        drops.append(_drop(loc, reason))
    # else: a genuinely empty placeholder / connector with no payload -> skipped-empty (no drop).


def _walk(
    shapes: Any,
    slide_no: int,
    units: list[str],
    lines: list[str],
    drops: list[Unextracted],
) -> None:
    """Walk ``shapes`` in document order; recurse into GROUPS; decide each LEAF via ``_emit_or_report``.

    A ``GroupShape`` is recursed into (``group.shapes``), NEVER reported as one drop — that would
    lose the readable text of its members and break the accounting identity. Handles nested groups
    by recursing. The accounting identity then holds WITH nesting: every leaf shape contributes
    exactly one of {>=1 unit, skipped-empty, one drop}; a GROUP node contributes nothing.
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE  # noqa: PLC0415 — lazy, optional [pptx] extra

    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            _walk(shape.shapes, slide_no, units, lines, drops)  # recurse — never report the group
            continue
        _emit_or_report(shape, slide_no, units, lines, drops)


def _serialize(prs: Any) -> tuple[str, list[str], list[Unextracted]]:
    r"""Walk the deck slide-by-slide and build ``(transcript, units, drops)`` deterministically.

    Slides in deck order; within a slide, shapes in document order (groups recursed); speaker notes
    LAST (after the slide's shapes). Each emitted text yields a transcript line
    ``Slide N / <shape>\t<text>\n`` (or ``[r{r}c{c}]`` for a cell, ``/ notes`` for a note) and
    appends the SAME text string to ``units`` (so every unit is, by construction, an exact substring
    of the transcript at its own offset; the structural prefix separates adjacent values, so
    duplicate values get distinct, ordered, locatable spans via ``normalize()``'s forward cursor).
    """
    lines: list[str] = []
    units: list[str] = []
    drops: list[Unextracted] = []

    for idx, slide in enumerate(prs.slides):
        slide_no = idx + 1
        _walk(slide.shapes, slide_no, units, lines, drops)
        # Speaker notes last (a documented, stable position): one verbatim unit per non-empty note
        # paragraph, with the ``Slide N / notes`` locator.
        if slide.has_notes_slide:
            notes_tf = slide.notes_slide.notes_text_frame
            for para in notes_tf.paragraphs:
                text = para.text
                if text:
                    lines.append(f"Slide {slide_no} / notes{SEP}{text}\n")
                    units.append(text)

    return "".join(lines), units, drops


class PptxAdapter:
    """Parse ``.pptx`` bytes into a ``Source`` + verbatim units, then ``normalize()`` -> result.

    Registered under the name ``"pptx"`` (see ``adapters/__init__.py``). ``distill(sources)`` is
    signature-exact with ``DistillPort``; ``parse(raw, path)`` is the public entrypoint that turns
    raw ``.pptx`` bytes into the ``Source`` this backend distills. The adapter is STATELESS across
    ``parse()``/``distill()`` (R1): per-Source ``unextracted[]`` travels on the Source itself
    (``Source.extraction`` via the shared codec), never in instance memory.
    """

    name = "pptx"

    # ------------------------------------------------------------------ #
    # Parse: raw .pptx bytes -> (Source, transcript-order units, unextracted[])
    # ------------------------------------------------------------------ #
    def parse(
        self, raw: bytes, path: str
    ) -> tuple[Source, list[str], list[Unextracted]]:
        """Parse ``raw`` ``.pptx`` bytes into a ``Source`` + transcript-order units + drops.

        Lazily loads a single ``Presentation`` via ``load_presentation``, walks slides/shapes
        (groups recursed) + notes into the canonical transcript + units, and collects the unreadable
        taxonomy into ``drops``. Sets ``source.extraction = encode_coverage(drops)`` so the drops
        travel WITH the Source (R1).

        Robust to hostile input (V5 / T-06-06): a malformed/unreadable ``.pptx`` (python-pptx
        raises) is caught and disclosed as a SINGLE whole-source ``unextracted[]`` entry with an
        empty transcript — never an unhandled crash. The adapter reads only shape TYPE and text, so
        no embedded image/OLE/media bytes are ever touched (T-06-08).
        """
        try:
            prs = load_presentation(raw)
        except Exception as exc:  # noqa: BLE001 — untrusted input; disclose, never crash
            drop = _drop(path, _R_UNREADABLE.format(error=type(exc).__name__))
            source = Source(
                id=path,
                context=path,
                transcript="",
                extraction=encode_coverage([drop]),
                # A whole-source-unreadable deck has no intrinsic timestamp; pass the deterministic
                # sentinel explicitly (L1) so even this error path never falls back to wall-clock.
                timestamp=deterministic_timestamp(None),
            )
            return source, [], [drop]

        transcript, units, drops = _serialize(prs)

        # A deterministic, document-intrinsic timestamp (mirrors EmailAdapter sourcing it from the
        # Date header, ExcelAdapter from docProps created — NOT now()). python-pptx's
        # ``core_properties.created`` is the analog AND, unlike openpyxl, it does NOT fabricate a
        # wall-clock value: it returns ``None`` faithfully when the deck has no ``dcterms:created``
        # (probed against python-pptx 1.0.2). So we read it directly — no raw-XML workaround needed.
        # ``deterministic_timestamp`` maps ``None`` -> EPOCH_ZERO and coerces a tz-naive value to UTC,
        # so a .pptx with no intrinsic ``created`` parses to a byte-identical Source twice (L1).
        created = getattr(prs.core_properties, "created", None)
        timestamp = deterministic_timestamp(
            created if isinstance(created, datetime) else None
        )

        source = Source(
            id=path,
            context=path,
            transcript=transcript,
            extraction=encode_coverage(drops),
            timestamp=timestamp,
        )
        return source, units, drops

    # ------------------------------------------------------------------ #
    # The DistillPort entrypoint
    # ------------------------------------------------------------------ #
    def distill(self, sources: list[Source]) -> DistillationResult:
        """Mint claims via ``normalize()`` and merge adapter-drops with normalize-drops.

        Each ``Source`` was produced by ``parse()`` (transcript + carried ``extraction``). We
        re-derive the units from the transcript the SAME deterministic way ``parse`` did, so a
        ``Source`` that round-tripped through JSON still distills identically on a FRESH adapter.
        Returns truth only — never publishes (HARD RULE, ``manual.py``).
        """
        all_claims = []
        merged_unextracted: list[Unextracted] = []
        traces: list[Source] = []

        for source in sources:
            # R2 safety-net (belt-and-suspenders): a Source this adapter did not produce carries no
            # ``extraction`` record, and its adapter-side drops are NOT reconstructable from the
            # transcript alone (a picture/SmartArt leaves no transcript line). Record an explicit
            # 'coverage-not-reconstructable' marker (forces Coverage.complete=False) — honest
            # uncertainty over a false complete=True. A Source this adapter parse()d has
            # ``extraction`` set, so the marker does NOT fire.
            if source.extraction is None:
                merged_unextracted.append(not_reconstructable_marker(source.id))

            units, adapter_unx = self._units_for(source)
            claims, norm_unx = normalize(source, units)
            all_claims.extend(claims)
            merged_unextracted.extend(adapter_unx)
            merged_unextracted.extend(norm_unx)
            traces.append(source)

        coverage = Coverage(
            complete=(len(merged_unextracted) == 0),
            unextracted=merged_unextracted,
            cost_hint="free",
            effort_hint="deterministic",
        )
        return DistillationResult(
            distillation=Distillation(claims=all_claims, traces=traces),
            coverage=coverage,
            backend=self.name,
        )

    def _units_for(self, source: Source) -> tuple[list[str], list[Unextracted]]:
        r"""Re-derive the text units from a built transcript ``Source`` + recover its drops.

        ``parse()`` put the canonical transcript on the ``Source`` as ``Slide N / <shape>\t<text>\n``
        records. Here we recover the verbatim text units by splitting on the unambiguous RECORD
        BOUNDARY (the anchored ``_RECORD_PREFIX`` ``^Slide \d+ / ... \t``) — NOT on ``\n``/``\t``,
        because a value may itself contain ``\n``, ``\v``, or ``\t`` (faithfulness — values are never
        escaped). This is the SAME serialization rule ``parse`` used, so each unit stays an exact
        substring of ``source.transcript`` and ``normalize()`` locates every one.

        The adapter-side ``unextracted[]`` (picture/chart/SmartArt/etc.) is NOT reconstructable from
        the transcript alone, so it is recovered from the typed carrier ``parse()`` set on the Source
        (``source.extraction``) via the shared codec — so a round-tripped Source distills with
        identical coverage on a FRESH adapter (R1).
        """
        units = _split_transcript_units(source.transcript)
        adapter_unx = decode_coverage(source.extraction)
        return units, adapter_unx


# A record prefix is exactly ``Slide {n} / ...{SEP}`` at a LINE START. A shape name cannot contain a
# newline (OOXML attribute) and the notes/cell variants share the ``Slide {n} / `` stem, so the
# prefix never contains a ``\n``. The ``SEP`` ends the prefix. ``[^\n]+?`` is non-greedy up to the
# FIRST ``SEP`` on the line (a shape name may itself contain ``/`` but not a tab in practice; if it
# did, the first tab still bounds the prefix and the value stays verbatim). This anchored pattern is
# how we tell a TRUE record start from a value's own continuation line (a wrapped value line begins
# with the value's own text, not a ``Slide N / ...{SEP}`` prefix — see _split_transcript_units).
_RECORD_PREFIX = re.compile(rf"Slide \d+ / [^\n]*?{re.escape(SEP)}")


def _split_transcript_units(transcript: str) -> list[str]:
    r"""Recover the ordered text units from a serialized transcript (the inverse of ``_serialize``).

    The transcript is the concatenation of ``"Slide N / <shape>{SEP}{value}\n"`` records. A value
    may itself contain ``\n``, ``\v``, and/or ``SEP`` (faithfulness — values are NEVER escaped), so
    we cannot split on ``\n`` or ``SEP`` alone. We instead split on the unambiguous RECORD BOUNDARY:
    a record starts at a LINE whose prefix matches ``Slide N / ...{SEP}`` (the anchored
    ``_RECORD_PREFIX``). A value's wrapped continuation line never matches that prefix (it begins
    with the value's own text), so each unit is exactly the text from after its prefix's ``SEP`` up
    to the ``\n`` that immediately precedes the NEXT record prefix (or end of transcript). Each
    recovered unit is thus a verbatim substring of the transcript, in order — so ``normalize()``
    locates every one and duplicates get distinct, forward-only offsets.

    Returns ``[]`` for an empty transcript (a whole-source-unreadable Source carries no units).
    """
    if not transcript:
        return []
    # Find the byte offset where every TRUE record starts (a line-start prefix match).
    starts: list[int] = []
    for m in _RECORD_PREFIX.finditer(transcript):
        at_line_start = m.start() == 0 or transcript[m.start() - 1] == "\n"
        if at_line_start:
            starts.append(m.start())
    units: list[str] = []
    for i, start in enumerate(starts):
        # The record's SEP-bounded prefix: find the first SEP at or after this record start. The
        # regex matched a prefix ending at the FIRST SEP, so re-find it from ``start``.
        match = _RECORD_PREFIX.match(transcript, start)
        # ``match`` is guaranteed (we found ``start`` via the same pattern). Its end is just past SEP.
        assert match is not None
        value_start = match.end()
        # The value runs up to the '\n' that precedes the NEXT record start (or to end-of-text for
        # the last record); strip exactly the serializer's single terminal '\n'.
        next_start = starts[i + 1] if i + 1 < len(starts) else len(transcript)
        value_end = next_start
        if value_end > value_start and transcript[value_end - 1] == "\n":
            value_end -= 1
        units.append(transcript[value_start:value_end])
    return units
