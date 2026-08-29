"""WKLY-05 tier 2 + SC-3 — the committed weekly deck, proven against a FRESH render.

WHY THIS MODULE IS SEPARATE FROM ``tests/test_weeklysite.py``. That module is stdlib-only by
construction so it can run in the ``site-integrity`` CI job, which installs ``[test,config]`` and
has NO ``0 skipped`` assertion. Everything here needs the optional ``[pptx]`` extra, so it lives
in its own module wired into the ``weekly`` job — the one job that installs ``[pptx]`` AND fails
the build on any reported skip. Putting a ``[pptx]``-gated test in ``site-integrity`` would make it
skip forever and read as green: that is W21, and this repo has already paid for it twice.

THE GUARD IS A ``skipif``, NOT A MODULE-LEVEL ``importorskip``, and the difference is the whole
point. ``importorskip`` makes the module *vanish* on an install without the extra — pytest reports
one skipped file and the ``0 skipped`` grep in the weekly job would catch that, but a reader of the
log sees nothing named. The ``skipif`` idiom (bind ``_pptx = None`` in an ``except ImportError``,
mark each test) keeps the module importable everywhere, so a missing extra shows up as NAMED skips
against NAMED tests in the one job that treats a skip as a failure.

The four things proven here, none of which the corpus can prove about itself:

1. **Tier 2 — the committed deck IS the deck.** ``content/weekly/deck/*.pptx`` equals a fresh
   render of the same record under ``part_digest``. Tier 1 (``tests/test_weeklysite.py``) proves
   the committed binary matches its committed sidecar — that catches a hand-edit of the binary,
   but it cannot catch the pair drifting together away from what the CODE now produces. Only a
   fresh render can, and only this job installs the extra that makes one possible.
2. **Determinism across a REAL time boundary.** Two renders separated by an actual
   ``time.sleep(3)``, on the corpus's OWN inputs, give equal ``part_digest``.
3. **SC-3 — the deck reads back Draft.** The WRITTEN file is reopened with python-pptx, and the
   watermark and the generated-by marker are asserted FROM the reopened file, never from the
   writer's return value.
4. **The gate is read, never written.** ``surface.model_dump()`` is identical before and after a
   render, and ``weeklysite.py`` names no gate-advancing call at all.

NEVER ``read_bytes() == read_bytes()`` ON A ZIP. DEFLATE output is zlib-implementation-dependent
(zlib vs zlib-ng), so a raw-byte assertion between the committed deck and a runner's fresh render
is green locally and red in CI for a reason that has nothing to do with the record's content.
``part_digest`` is sorted, length-prefixed ``(part name, sha256(part))`` rows — it is exactly the
implementation-INDEPENDENT assertion a cross-environment gate needs. And nothing here compares a
rendered deck to the TEMPLATE: python-pptx's LOAD path re-serializes empty core properties and
re-orders parts, so that comparison is guaranteed-wrong (02-03 decision).
"""

from __future__ import annotations

import pathlib
import time

import pytest

from newsletters import pptx_writer, weeklysite
from newsletters.semantic import ReviewState
from newsletters.weeklyspec import weekly_slots

# The writer's contract strings and its two comparison helpers are reached through the MODULE, so
# this file's imports stay one-per-line: a parenthesized multi-line `from ... import (...)` is
# exactly the shape isort and black disagree about until isort gets `profile = "black"` (DEF-15,
# maintainer-gated). Same reason tests/test_weeklysite.py reaches weeklyspec's constants this way.
DRAFT_STATUS = pptx_writer.DRAFT_STATUS
MARKER = pptx_writer.MARKER
WATERMARK_NAME = pptx_writer.WATERMARK_NAME
part_digest = pptx_writer.part_digest
differing_parts = pptx_writer.differing_parts

try:  # noqa: SIM105 — the guard needs the bound name, not just the suppression
    import pptx as _pptx
except ImportError:  # pragma: no cover — a bare install, never the weekly job
    _pptx = None

requires_pptx = pytest.mark.skipif(
    _pptx is None, reason="optional [pptx] extra (python-pptx) not installed"
)

REPO_ROOT = pathlib.Path(__file__).resolve().parent.parent
CORPUS = REPO_ROOT / "content" / "weekly"
DECK_DIR = CORPUS / "deck"
CORPUS_TEMPLATE = CORPUS / "template.pptx"


def _committed_deck() -> pathlib.Path:
    """The single committed deck, DISCOVERED rather than named (it tracks the authored week)."""
    decks = sorted(DECK_DIR.glob("*.pptx"))
    assert len(decks) == 1, f"expected exactly one committed deck, found {decks}"
    return decks[0]


def _load_surface_and_slots():
    """The load, its composed surface, and the slots derived from THE PAIR.

    Reached through ``weeklysite._load_and_compose`` — the module-private helper that
    ``build_weekly_surfaces`` and ``build_weekly_deck`` BOTH call — rather than through a second
    load written here. ``weekly_slots(load, surface)`` validates its disclosure lines against
    ``surface.missing``, so a surface composed from one load and slots derived from another would
    be checking one record's slides against another record's honesty panel. A test that rebuilt
    the sequence by hand would be asserting against its own copy of the pipeline, not the shipped
    one, and would keep passing after the shipped one changed.

    ``root=REPO_ROOT`` so this resolves the committed corpus regardless of the cwd pytest ran in.
    """
    load, surface = weeklysite._load_and_compose(
        root=REPO_ROOT, spec_path=None, lanes_path=None, author=None
    )
    return load, surface, weekly_slots(load, surface)


# --------------------------------------------------------------------------- #
# 1 — tier 2: the committed deck equals a fresh render
# --------------------------------------------------------------------------- #


@requires_pptx
def test_committed_deck_equals_fresh_render(tmp_path: pathlib.Path) -> None:
    """TIER 2 (T-04-04): a FRESH render of the committed record reproduces the committed deck.

    Rendered through ``build_weekly_deck`` — the same entry point ``newsletters weekly`` calls and
    the same one that produced the committed binary — with every path left to the corpus's own
    structural discovery, so this asserts the SHIPPED path rather than a test-local reconstruction
    of it. Compared under ``part_digest`` (see the module docstring on why never raw bytes).

    On failure the message names ``differing_parts``, so a red teaches which PART moved instead of
    only reporting that two 64-hex strings are unequal.
    """
    committed = _committed_deck()
    committed_bytes = committed.read_bytes()

    fresh_path = weeklysite.build_weekly_deck(tmp_path / committed.name, root=REPO_ROOT)
    fresh_bytes = fresh_path.read_bytes()

    assert part_digest(fresh_bytes) == part_digest(committed_bytes), (
        f"the committed deck has drifted from what the code now renders. Differing parts: "
        f"{differing_parts(fresh_bytes, committed_bytes)}. Regenerate it with `newsletters "
        f"weekly` in the SAME commit as the change — never by editing the binary or its sidecar."
    )


@requires_pptx
def test_fresh_render_reproduces_the_committed_digest_sidecar(
    tmp_path: pathlib.Path,
) -> None:
    """The sidecar the fresh render writes equals the COMMITTED sidecar, byte for byte.

    Tier 1 (stdlib-only, ``tests/test_weeklysite.py``) checks the committed deck against the
    committed sidecar; this closes the remaining hole by tying BOTH committed artifacts to a fresh
    render, so a deck and a sidecar that drifted together cannot pass as a matched pair.
    """
    committed = _committed_deck()
    committed_sidecar = committed.with_suffix(committed.suffix + ".digest")

    fresh_path = weeklysite.build_weekly_deck(tmp_path / committed.name, root=REPO_ROOT)
    fresh_sidecar = fresh_path.with_suffix(fresh_path.suffix + ".digest")

    assert fresh_sidecar.is_file(), "build_weekly_deck wrote no digest sidecar"
    assert fresh_sidecar.read_bytes() == committed_sidecar.read_bytes(), (
        "the freshly written digest sidecar differs from the committed one — the committed "
        "deck/sidecar PAIR has drifted from the code"
    )


# --------------------------------------------------------------------------- #
# 2 — determinism across a real time boundary, on the corpus's own inputs
# --------------------------------------------------------------------------- #


@requires_pptx
def test_double_render_is_stable_across_a_real_time_boundary() -> None:
    """Two renders separated by an ACTUAL 3 seconds give the same ``part_digest``.

    Phase 1 recorded the property and Phase 2 measured it on fixtures; this re-proves it on the
    corpus's own committed inputs, which is the only place it matters for the shipped sample. The
    sleep is real on purpose: an in-process double render can cross no DOS-time boundary, and the
    ONE measured non-determinism in a raw python-pptx write is exactly a ``date_time`` field. A
    mocked clock would prove that the mock worked.
    """
    _load, surface, slots = _load_surface_and_slots()
    render = pptx_writer.render_surface_pptx_bytes

    first = render(surface, template=CORPUS_TEMPLATE, slots=slots)
    time.sleep(3)
    second = render(surface, template=CORPUS_TEMPLATE, slots=slots)

    assert part_digest(first) == part_digest(second), (
        "two renders three seconds apart differ in PART CONTENT — the writer embedded a clock. "
        f"Differing parts: {differing_parts(first, second)}"
    )


# --------------------------------------------------------------------------- #
# 3 — SC-3: the written deck reads back Draft-watermarked and marked
# --------------------------------------------------------------------------- #


@requires_pptx
def test_rendered_deck_reads_back_draft_watermarked_and_marked(
    tmp_path: pathlib.Path,
) -> None:
    """SC-3: REOPEN the written file and assert the Draft watermark + the generated-by marker.

    Every assertion here is made against the reopened FILE, never against the writer's return
    value: "the writer says it watermarked" and "the deck on disk is watermarked" are two
    different claims, and only the second one is what a reader who opens the deck experiences.

    The watermark is asserted on EVERY slide, not just the first: a two-slide deck with one
    watermark is a deck whose second page reads as approved.
    """
    _load, surface, _slots = _load_surface_and_slots()
    assert not surface.is_published, (
        "fixture invariant: the sample must be Draft, or this test would be asserting the "
        "watermark of a published record"
    )

    out_path = weeklysite.build_weekly_deck(tmp_path / "readback.pptx", root=REPO_ROOT)

    written = _pptx.Presentation(str(out_path))  # reopen the FILE that was written
    core = written.core_properties
    assert core.category == MARKER, core.category
    assert core.content_status == DRAFT_STATUS, core.content_status
    assert core.identifier == surface.id, core.identifier

    slides = list(written.slides)
    assert slides, "the reopened deck has no slides"
    for index, slide in enumerate(slides):
        names = [shape.name for shape in slide.shapes]
        assert WATERMARK_NAME in names, (
            f"slide {index} of a DRAFT weekly carries no watermark — an unreviewed page that "
            f"does not look unreviewed. Shapes: {names}"
        )


# --------------------------------------------------------------------------- #
# 4 — the gate is READ, never written
# --------------------------------------------------------------------------- #


@requires_pptx
def test_render_leaves_the_surface_draft(tmp_path: pathlib.Path) -> None:
    """The product's hardest rule on the deck path: a render is a READ of the gate.

    Asserted on the WHOLE ``model_dump()`` before and after, not just on ``review.state``. The
    inverted form is the point: a writer that advanced the gate would fail a state check, but a
    writer that quietly wrote back ANY other field would pass one — and it would still be writing
    into a reviewed record.
    """
    _load, surface, slots = _load_surface_and_slots()
    before = surface.model_dump()

    pptx_writer.render_surface_pptx(
        surface,
        template=CORPUS_TEMPLATE,
        slots=slots,
        out_path=tmp_path / "gate.pptx",
    )

    assert surface.model_dump() == before, "the render mutated the Surface"
    assert surface.review.state is ReviewState.DRAFT
    assert not surface.is_published


def test_no_gate_transition_in_the_weekly_build_path() -> None:
    """T-04-08 / the first hard rule: the weekly builder names NO gate-advancing call.

    A source-level assertion, deliberately: it holds even for a code path no test happens to
    execute, and it is the same shape as the Phase-3 guard over ``weeklyspec.py``. NOT
    ``[pptx]``-gated — reading source needs no optional extra, and the one rule that must never
    break should not be gated behind an install.

    The rule: no `Surface` reaches `Published` except through the review gate with a recorded
    reviewer. A corpus builder that could publish would put the whole product's promise inside
    a build script.
    """
    code = pathlib.Path(weeklysite.__file__).read_text(encoding="utf-8")
    for call in (".publish(", ".approve(", ".open_pull_request("):
        assert call not in code, (
            f"{call!r} appears in weeklysite.py — the corpus builder must only ever produce "
            "Draft. No auto-publish, ever: a Surface reaches Published through the review gate "
            "with a recorded reviewer, and nowhere else."
        )
