"""Risk-A1 probe: confirm the graphic-frame ``@uri`` accessor python-pptx exposes (ADAPT-04).

WHY THIS FILE EXISTS (RESEARCH risk A1 + Open Question 1):
The PPTX adapter (06-03) must detect SmartArt diagrams so they are *reported* in
``unextracted[]`` rather than silently dropped (the zero-silent-drops invariant). python-pptx
1.0.2 has NO high-level SmartArt API: a SmartArt diagram is a ``graphicFrame`` whose
``shape_type`` is ``None``, ``has_table``/``has_chart`` are both ``False``, and the ONLY reliable
signal is the ``<a:graphicData>`` element's ``@uri`` attribute equalling the diagram namespace.

RESEARCH (risk A1) flagged the exact ACCESSOR NAME as the one uncertain edge: it claimed
``shape._element.graphicData_uri`` is reachable from a high-level ``GraphicFrame`` but warned the
name is internal and could be renamed, recommending an lxml ``@uri`` fallback so the worst case is
"reported as an unknown graphic frame", never a crash and never a silent drop.

This probe RESOLVES that uncertainty against the installed python-pptx so 06-03 can wire the
verified recipe with confidence. The CONFIRMED recipe (python-pptx 1.0.2, this test green) is:

    DIAGRAM_URI = "http://schemas.openxmlformats.org/drawingml/2006/diagram"

    def graphic_frame_uri(shape) -> str | None:
        # Primary accessor — CONFIRMED present & reliable on the underlying CT_GraphicalObjectFrame:
        uri = getattr(shape._element, "graphicData_uri", None)
        if uri is not None:
            return uri
        # lxml fallback (used only if the private accessor is ever renamed/absent):
        from pptx.oxml.ns import qn
        gd = shape._element.find(".//" + qn("a:graphicData"))
        return gd.get("uri") if gd is not None else None

    # SmartArt iff: shape.shape_type is None AND graphic_frame_uri(shape) == DIAGRAM_URI
    # (a non-None, non-diagram uri on a shape_type-None frame -> "unknown graphic frame", reported)

Both paths were observed to AGREE on the diagram URI here. Skips cleanly when python-pptx is
absent so the bare-install gate (tests/test_ai_optional.py) is unaffected.
"""

from __future__ import annotations

import copy
import io

import pytest

from newsletters.adapters._pptx_loader import load_presentation

# The drawingml diagram (SmartArt) namespace — the single signal that identifies SmartArt.
DIAGRAM_URI = "http://schemas.openxmlformats.org/drawingml/2006/diagram"
# A real, non-diagram graphicData uri (a table) — used to prove the accessor returns the *actual*
# uri, not a hardcoded diagram constant.
TABLE_URI = "http://schemas.openxmlformats.org/drawingml/2006/table"

# Skip the whole module on a bare env (no [pptx]); the loader's lazy boundary is proven separately
# in tests/test_ai_optional.py. This probe needs python-pptx to inspect real OOXML.
pptx = pytest.importorskip("pptx", reason="python-pptx not installed (bare env); A1 probe needs it")


def _diagram_pptx_bytes() -> bytes:
    """Author a minimal in-memory ``.pptx`` whose one shape is a *diagram* graphicFrame.

    There is no high-level SmartArt authoring API, so we follow the RESEARCH fixture recipe in
    miniature: create a normal table graphicFrame, then XML-mutate its ``<a:graphicData>`` ``@uri``
    to the diagram namespace and strip the ``<a:tbl>`` child — yielding a graphicFrame that python
    -pptx classifies exactly like SmartArt (``shape_type is None``, ``has_table``/``has_chart``
    False), which is what the 06-03 detection must handle. Round-tripped through ``save()`` so we
    parse the same way :func:`load_presentation` will in production.
    """
    from pptx.oxml.ns import qn
    from pptx.util import Inches

    prs = pptx.Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    gf = slide.shapes.add_table(1, 1, Inches(1), Inches(1), Inches(2), Inches(1))
    graphic_data = gf._element.find(".//" + qn("a:graphicData"))
    graphic_data.set("uri", DIAGRAM_URI)
    tbl = graphic_data.find(qn("a:tbl"))
    if tbl is not None:
        graphic_data.remove(tbl)  # no longer a real table -> closer to a bare SmartArt frame

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _graphic_frame_uri_via_accessor(shape: object) -> object:
    """The PRIMARY accessor 06-03 should use: the private-but-stable ``graphicData_uri``."""
    return getattr(shape._element, "graphicData_uri", None)  # type: ignore[attr-defined]


def _graphic_frame_uri_via_lxml(shape: object) -> object:
    """The lxml FALLBACK: read ``@uri`` off the ``<a:graphicData>`` element directly."""
    from pptx.oxml.ns import qn

    graphic_data = shape._element.find(".//" + qn("a:graphicData"))  # type: ignore[attr-defined]
    return graphic_data.get("uri") if graphic_data is not None else None


def test_diagram_graphic_frame_has_smartart_signature() -> None:
    """A diagram graphicFrame matches the RESEARCH SmartArt signature: shape_type None, no table/chart."""
    prs = load_presentation(_diagram_pptx_bytes())
    shape = prs.slides[0].shapes[0]
    # The defining SmartArt signature (RESEARCH Pitfall 1 / taxonomy): the public ladder is blind.
    assert shape.shape_type is None, f"expected shape_type None for a diagram frame, got {shape.shape_type}"
    assert shape.has_table is False, "diagram frame must not be classed as a table"
    assert shape.has_chart is False, "diagram frame must not be classed as a chart"


def test_graphic_frame_uri_accessor_and_lxml_fallback_agree_on_diagram() -> None:
    """A1 resolved: ``graphicData_uri`` and the lxml ``@uri`` fallback BOTH return the diagram URI.

    This is the load-bearing assertion 06-03 relies on. If the private accessor is ever renamed,
    the accessor branch returns ``None`` and the lxml fallback still returns the diagram URI, so
    detection degrades to "unknown graphic frame" (reported), never a silent drop — but here, on
    python-pptx 1.0.2, the accessor IS present and the two paths agree.
    """
    prs = load_presentation(_diagram_pptx_bytes())
    shape = prs.slides[0].shapes[0]

    via_accessor = _graphic_frame_uri_via_accessor(shape)
    via_lxml = _graphic_frame_uri_via_lxml(shape)

    # The lxml fallback MUST always find the diagram URI (it reads the raw @uri attribute).
    assert via_lxml == DIAGRAM_URI, f"lxml fallback failed to read diagram @uri: {via_lxml!r}"

    if via_accessor is not None:
        # On 1.0.2 the accessor is present -> the two paths must AGREE (records the confirmed name).
        assert via_accessor == DIAGRAM_URI, f"accessor returned wrong uri: {via_accessor!r}"
        assert via_accessor == via_lxml, (
            f"accessor ({via_accessor!r}) and lxml fallback ({via_lxml!r}) disagree"
        )
    else:
        # Worst case A1 warned about: accessor renamed/absent. The lxml fallback carries detection;
        # 06-03 must therefore ALWAYS keep the fallback. (Not expected on 1.0.2 — fails loudly if so.)
        pytest.fail(
            "graphicData_uri accessor is ABSENT on this python-pptx — 06-03 must rely on the "
            "lxml @uri fallback (which still works); update the confirmed-accessor docs."
        )


def test_accessor_returns_actual_uri_not_hardcoded_diagram() -> None:
    """The accessor returns the FRAME's real uri (a table reads as the table uri, not diagram).

    Guards against a false-positive detection: a real table/chart graphicFrame must NOT be flagged
    as SmartArt. Proves the accessor is reading the live ``@uri``, so the diagram-vs-other branch
    in 06-03's taxonomy is sound.
    """
    from pptx.util import Inches

    prs = pptx.Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    table_frame = slide.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(4), Inches(2))

    assert _graphic_frame_uri_via_accessor(table_frame) == TABLE_URI
    assert _graphic_frame_uri_via_lxml(table_frame) == TABLE_URI
    assert table_frame.shape_type is not None, "a real table must report a concrete shape_type"
    assert table_frame.has_table is True
