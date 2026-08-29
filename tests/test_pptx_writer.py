"""The `.pptx` writer's proof module (WKLY-01; ROADMAP Phase 2's five success criteria).

WHAT THIS MODULE PROVES, once the phase is complete:

1. **SC-1 — the template contract.** A composed `Surface(REPORT)` fills the operator's EXISTING
   slides through `NL_`-named shapes, and the binding fails LOUD in both directions (content bound
   to a name the template lacks; an `NL_` slot with no content), on a duplicate shape name, on a
   slot that cannot hold text, and on a template that already owns the watermark name.
2. **SC-2 — determinism.** Two in-process renders of the same Surface are byte-identical after
   `normalize_opc_zip`, with the negative control living next door in `test_pptx_determinism.py`.
3. **SC-3 — the generated-by marker**, asserted by reading the WRITTEN file back.
4. **SC-4 — the review gate is untouched.** Rendering a Draft Surface leaves it Draft.
5. **SC-5 — AI-optional.** The writer's bare-install boundary is guarded in
   `tests/test_ai_optional.py` (NOT here: this module skips itself without the `[pptx]` extra, so a
   guard proving "imports without the extra" would never run where it matters).

Plan 02-01 lands the scaffold: the deck builders every later test binds against, and the self-test
that proves the decks are what they claim to be.

WHY THE DECKS ARE BUILT IN-TEST AND THE COMMITTED TEMPLATE IS NOT REGENERATED (decision P-06).
`tests/fixtures/weekly/template.pptx` stays BYTE-UNCHANGED this phase.
`.planning/notes/2026-08-29-pptx-determinism-evidence.json` records `part_digest_a` / `part_digest_b`
computed from that exact template, and both are in the recorder's `CHECKED_FIELDS` — so regenerating
it would turn this phase's own `_record_determinism_evidence.py --check` gate red and make the
decision note's cited digests stale. Retiring committed evidence to buy test convenience is exactly
the trade this repo does not make. The rich and pathological decks the fail-loud contract needs are
therefore authored into `tmp_path` here, which is also how 02-RESEARCH measured W17 and W14 in the
first place.

EVERY BUILDER IS A FIXED POINT OF THE NORMALIZER. Each saves to `io.BytesIO`, routes the bytes
through `normalize_opc_zip`, and does ONE `write_bytes` — the same discipline the writer itself uses.
A deck built here therefore contributes no non-determinism to any test that renders through it, and
the self-test below re-proves that rather than asserting it in prose.

EVERY BUILDER WRITES ONLY INTO ITS `directory` ARGUMENT. `FIXTURE_DIR` is read-only in this module
(plan 02-03 renders through the committed template); a builder that wrote into the fixture directory
would break `test_pptx_determinism.py::test_weekly_fixture_corpus_is_exactly_the_committed_template`.

ALL FIXTURE CONTENT IS FABRICATED AND NEUTRAL — no person, no organisation, no third-party tool
name. `Presentation()` with no argument loads python-pptx's bundled default template, whose
`docProps/core.xml` ships a third party's name and a foreign tool's marketing string; every builder
overwrites the whole core-properties block for the same reason `_author_template.py` does.
"""

from __future__ import annotations

import io
import pathlib
import time
import zipfile
from datetime import datetime
from typing import NamedTuple, Union

import pytest

# `importorskip` raises Skipped at module level, which pytest converts into a module skip — so a
# bare install (no [pptx] extra) SKIPS this file instead of erroring at collection. A `pytestmark`
# skipif can NOT guard a module-level import: the marker is evaluated per collected test item,
# AFTER pytest has already imported the module. Same pattern as tests/test_pptx_determinism.py.
pptx = pytest.importorskip(
    "pptx", reason="optional [pptx] extra (python-pptx) not installed"
)
Presentation = pptx.Presentation

from lxml import etree  # noqa: E402
from pptx.enum.shapes import MSO_SHAPE_TYPE  # noqa: E402
from pptx.enum.text import MSO_AUTO_SIZE  # noqa: E402
from pptx.oxml.ns import qn  # noqa: E402
from pptx.util import Inches, Pt  # noqa: E402

from newsletters import (  # noqa: E402
    REPORT,
    Claim,
    ClaimsBlock,
    ProseBlock,
    Review,
    ReviewState,
    Surface,
    Trace,
    pptx_writer,
)
from newsletters.adapters._timestamps import EPOCH_ZERO  # noqa: E402
from newsletters.pptx_writer import (  # noqa: E402
    DRAFT_STATUS,
    MARKER,
    SLOT_PREFIX,
    WATERMARK_NAME,
    WATERMARK_TEXT,
    bind_slots,
    differing_parts,
    differing_zipinfo_fields,
    fill_slot,
    normalize_opc_zip,
    part_digest,
    render_surface_pptx,
    render_surface_pptx_bytes,
)

FIXTURE_DIR = pathlib.Path(__file__).parent / "fixtures" / "weekly"
# READ-ONLY here. Plan 02-03 renders through it; nothing in this module writes to FIXTURE_DIR.
COMMITTED_TEMPLATE = FIXTURE_DIR / "template.pptx"

# A fixed instant pinned into every built deck's docProps/core.xml so no wall-clock is embedded.
# DELIBERATELY NOT `EPOCH_ZERO` — same falsifiability control as `_author_template._FIXED` (P-05):
# the writer is specified to pin `dcterms:created`/`.modified` to EPOCH_ZERO, so a TEMPLATE that
# already carried EPOCH_ZERO would let plan 02-02's marker read-back pass on a deck the writer
# never touched.
_FIXED = datetime(2026, 1, 1, 0, 0, 0)

# Fabricated, neutral core properties — every one overwrites a stock-template value. Mirrors
# `_author_template._CORE_PROPERTIES`. `category` and `content_status` are deliberately EMPTY: the
# generated-by marker and the gate state are the WRITER's to set, and a template that pre-carried
# them would make the 02-02 read-back assertion unfalsifiable.
_CORE_PROPERTIES = {
    "author": "newsletters fixture author",
    "last_modified_by": "newsletters fixture author",
    "title": "In-test weekly deck (fabricated fixture)",
    "subject": "Phase 2 renderer test deck",
    "comments": "Fabricated in-test deck; built by tests/test_pptx_writer.py",
    "category": "",
    "content_status": "",
    "keywords": "",
    "identifier": "",
    "language": "",
    "version": "",
}

# The happy-path content mapping for the rich deck: every `NL_` name it carries, including the one
# nested in a group and the one on slide 2. `NL_HIGHLIGHTS` carries three lines so the multi-line
# clone path (02-RESEARCH Pattern 3) is exercised, and one line puts Unicode and XML metacharacters
# on the HAPPY path (W11) rather than in a special case nobody runs.
RICH_SLOTS: dict[str, list[str]] = {
    "NL_WEEK_TITLE": ["2026-W35"],
    "NL_MODULE": ["module-a"],
    "NL_HIGHLIGHTS": [
        "shipped the thing",
        'café — "smart" … 🎯 <&>',
        "closed the loop with the reviewer",
    ],
    "NL_LOWLIGHTS": ["the other thing slipped"],
    "NL_INSIDE_GROUP": ["a slot the operator grouped with its neighbours"],
    "NL_NEXT_WEEK": ["finish the thing that slipped"],
}

# The content mapping for the COMMITTED synthetic template — exactly the four `NL_` names it
# carries (`Footer` is deliberately absent: it has no reserved prefix, so it is not a slot).
# Fabricated and neutral, like everything else in this module. `NL_HIGHLIGHTS` carries more than one
# line so the end-to-end render exercises the clone path rather than only the reuse path.
COMMITTED_SLOTS: dict[str, list[str]] = {
    "NL_WEEK_TITLE": ["2026-W35"],
    "NL_MODULE": ["module-a"],
    "NL_HIGHLIGHTS": [
        "shipped the fabricated thing",
        "closed the loop with the reviewer",
    ],
    "NL_LOWLIGHTS": ["the other fabricated thing slipped"],
}

# Every shape name the rich deck contains, across both slides and INCLUDING group members. The
# group container and the footer are deliberately unprefixed: they are not slots.
RICH_SHAPE_NAMES = {
    "NL_WEEK_TITLE",
    "NL_MODULE",
    "NL_HIGHLIGHTS",
    "NL_LOWLIGHTS",
    "NL_INSIDE_GROUP",
    "NL_NEXT_WEEK",
    "Footer",
    "Narrative Group",
}


def _blank_slide(prs):
    """Add a slide on the BLANK layout (index 6 in python-pptx's default template)."""
    return prs.slides.add_slide(prs.slide_layouts[6])


def _scrub(prs) -> None:
    """Overwrite every core property with a fabricated, neutral value; pin the timestamps."""
    core = prs.core_properties
    for field, value in _CORE_PROPERTIES.items():
        setattr(core, field, value)
    core.created = _FIXED
    core.modified = _FIXED
    core.revision = 1


def _textbox(
    slide,
    name: str,
    left: float,
    top: float,
    width: float,
    height: float,
    body: str = "",
):
    """Add a named textbox with the overflow-safe text-frame settings (P-07 / Pitfall 3).

    `add_textbox`'s defaults are the worst combination — `word_wrap = False` with
    `auto_size = SHAPE_TO_FIT_TEXT` — i.e. no wrapping plus a stored size PowerPoint only corrects
    when a human edits the box, so overflowing text escapes the slide silently. `word_wrap = True`
    with `auto_size = NONE` makes overflow CLIP instead, which a reviewer can see. This is carried
    forward to Phase 4's `docs/weekly.md` operator recipe for the shipped template.
    """
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    box.name = name
    frame = box.text_frame
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.NONE
    frame.text = body
    return box


def _finalize(prs, directory: pathlib.Path, filename: str) -> pathlib.Path:
    """Save, NORMALIZE, then do ONE `write_bytes` into `directory`. Returns the written path."""
    buf = io.BytesIO()
    prs.save(buf)
    # Normalize BEFORE the file is written, so the built deck is a fixed point of the normalizer
    # and never contributes non-determinism to the tests that render through it.
    path = directory / filename
    path.write_bytes(normalize_opc_zip(buf.getvalue()))
    return path


def build_rich_template(directory: pathlib.Path) -> pathlib.Path:
    """The happy-path deck: two slides, a grouped slot, a decorative shape, real formatting.

    Slide 1 carries the four flat slots, an unprefixed `Footer`, and an unprefixed group
    (`Narrative Group`) holding `NL_INSIDE_GROUP`. **The group is deliberately NOT `NL_`-prefixed:**
    it is a container, not a slot, and giving it the reserved prefix would make it an unfilled
    reserved slot in every happy-path render. Slide 2 carries `NL_NEXT_WEEK` so "the watermark is on
    EVERY slide" is not a vacuous assertion on a one-slide deck.

    `NL_HIGHLIGHTS` is the formatting-fidelity carrier: `Pt(20)` bold with a `buChar` bullet, so a
    fill primitive that flattened the operator's formatting (02-RESEARCH W3: `text_frame.text = ...`
    does exactly that) is caught by a read-back assertion rather than by a human noticing.
    """
    prs = Presentation()
    _scrub(prs)

    slide = _blank_slide(prs)
    # Authored in a FIXED order so shape order is deterministic.
    _textbox(slide, "NL_WEEK_TITLE", 0.5, 0.4, 9.0, 1.0, "NL_WEEK_TITLE")
    _textbox(slide, "NL_MODULE", 0.5, 1.4, 9.0, 0.6, "NL_MODULE")
    highlights = _textbox(slide, "NL_HIGHLIGHTS", 0.5, 2.2, 4.3, 3.4, "NL_HIGHLIGHTS")
    _textbox(slide, "NL_LOWLIGHTS", 5.2, 2.2, 4.3, 3.4, "NL_LOWLIGHTS")
    # NOT a renderer slot: no NL_ prefix. Proves the reserved prefix discriminates.
    _textbox(
        slide, "Footer", 0.5, 6.0, 9.0, 0.4, "fabricated deck - not a renderer slot"
    )

    # The operator's formatting on the fidelity carrier: one run, sized and bold, plus a bullet.
    paragraph = highlights.text_frame.paragraphs[0]
    run = paragraph.runs[0]
    run.font.size = Pt(20)
    run.font.bold = True
    # Injecting bullet XML is legitimate HERE AND ONLY HERE: this builder stands in for an operator
    # authoring their deck in PowerPoint, which has no python-pptx API. 02-RESEARCH's ban on
    # SYNTHESIZING bullet XML applies to the WRITER, which must inherit the operator's formatting
    # and never construct it. `etree.SubElement` operates on an ALREADY-PARSED tree and constructs
    # no XML parser, so python-pptx's upstream `resolve_entities=False` mitigation is not weakened
    # (threat T-02-03) — do not replace this with a parse of an XML string.
    etree.SubElement(paragraph._p.get_or_add_pPr(), qn("a:buChar"), char="•")

    # `add_group_shape` MOVES existing shapes into the group, so the member is authored FIRST.
    inner = _textbox(slide, "NL_INSIDE_GROUP", 0.5, 5.0, 4.3, 0.8, "NL_INSIDE_GROUP")
    group = slide.shapes.add_group_shape([inner])
    group.name = "Narrative Group"

    second = _blank_slide(prs)
    _textbox(second, "NL_NEXT_WEEK", 0.5, 0.4, 9.0, 1.0, "NL_NEXT_WEEK")

    return _finalize(prs, directory, "rich.pptx")


def build_duplicate_name_template(directory: pathlib.Path) -> pathlib.Path:
    """Two top-level textboxes sharing the name `NL_WEEK_TITLE`.

    Duplicate shape names are LEGAL in OOXML and are exactly how an operator's copy-paste presents.
    A naive `{sh.name: sh for sh in ...}` comprehension is last-wins and silently DROPS a slot —
    the failure the fail-loud contract exists to prevent, so the binding map must raise instead.
    """
    prs = Presentation()
    _scrub(prs)
    slide = _blank_slide(prs)
    _textbox(slide, "NL_WEEK_TITLE", 0.5, 0.4, 9.0, 1.0, "first copy")
    _textbox(slide, "NL_WEEK_TITLE", 0.5, 1.6, 9.0, 1.0, "pasted copy")
    return _finalize(prs, directory, "duplicate_name.pptx")


def build_watermark_owning_template(directory: pathlib.Path) -> pathlib.Path:
    """A template that already defines the renderer-owned watermark name (02-RESEARCH W14).

    Measured: adding the watermark to such a deck writes TWO shapes with one name and raises
    nothing — so a "helpful" operator who copies the watermark into their template silently defeats
    the binding map. The writer must refuse the name at bind time.
    """
    prs = Presentation()
    _scrub(prs)
    slide = _blank_slide(prs)
    _textbox(slide, "NL_WEEK_TITLE", 0.5, 0.4, 9.0, 1.0, "NL_WEEK_TITLE")
    _textbox(slide, WATERMARK_NAME, 1.5, 2.5, 7.0, 2.0, WATERMARK_TEXT)
    return _finalize(prs, directory, "watermark_owning.pptx")


def build_nontext_slot_template(directory: pathlib.Path) -> pathlib.Path:
    """An `NL_`-prefixed shape that CANNOT hold text: an empty group named `NL_NOT_TEXT`.

    A group (like a picture) reports `has_text_frame == False`, so a naive
    `by_name[name].text_frame.text = ...` raises `AttributeError` — a stack trace, not a teaching
    error. The fill must check first and say what the operator should do instead.
    """
    prs = Presentation()
    _scrub(prs)
    slide = _blank_slide(prs)
    _textbox(slide, "NL_WEEK_TITLE", 0.5, 0.4, 9.0, 1.0, "NL_WEEK_TITLE")
    group = slide.shapes.add_group_shape()
    group.name = "NL_NOT_TEXT"
    return _finalize(prs, directory, "nontext_slot.pptx")


def build_empty_run_template(directory: pathlib.Path) -> pathlib.Path:
    """An `NL_` slot the operator "left blank on purpose": paragraph 0 with ZERO runs.

    A text box with no typed characters has no run, and therefore no formatting carrier for the
    fill primitive to reuse and clone (02-RESEARCH Pitfall 6). The writer must raise a teaching
    error telling the operator to type one placeholder character — not crash on an IndexError.
    """
    prs = Presentation()
    _scrub(prs)
    slide = _blank_slide(prs)
    _textbox(slide, "NL_WEEK_TITLE", 0.5, 0.4, 9.0, 1.0, "")
    return _finalize(prs, directory, "empty_run.pptx")


def _walk(shapes):
    """Every shape on a slide, DESCENDING into groups.

    `slide.shapes` is TOP-LEVEL ONLY (02-RESEARCH W17), which the second test below re-proves.
    """
    for shape in shapes:
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _walk(shape.shapes)


def _all_shape_names(prs) -> set[str]:
    """The recursive shape-name set across every slide."""
    return {shape.name for slide in prs.slides for shape in _walk(slide.shapes)}


def test_test_decks_have_the_expected_shape_inventory(tmp_path: pathlib.Path) -> None:
    """The rich deck is what the rest of this suite will assume it is — re-proved, not asserted.

    A test deck nobody checks is a fixture that quietly stops carrying the case it was built for:
    the group flattens, the bold run is lost to a refactor, the bullet never made it into the XML —
    and every fidelity assertion downstream goes green while proving nothing. This is the self-test
    that keeps the later reds attributable.
    """
    path = build_rich_template(tmp_path)
    written = Presentation(
        str(path)
    )  # reopen the WRITTEN file, never trust the builder

    assert len(written.slides) == 2, (
        "the rich deck no longer has two slides — 'the watermark is added to EVERY slide' becomes "
        "a vacuous assertion on a one-slide deck"
    )
    assert _all_shape_names(written) == RICH_SHAPE_NAMES, (
        "the rich deck's shape inventory changed; every fail-loud and fidelity assertion in this "
        f"phase is written against the documented set. Found: {sorted(_all_shape_names(written))}"
    )

    carrier = next(
        shape
        for shape in _walk(written.slides[0].shapes)
        if shape.name == "NL_HIGHLIGHTS"
    )
    paragraph = carrier.text_frame.paragraphs[0]
    assert paragraph.runs, (
        "the formatting carrier lost its run in the write/normalize round trip — there is nothing "
        "left for the fill primitive to reuse, so the fidelity tests would prove nothing"
    )
    assert paragraph.runs[0].font.size == Pt(20), (
        "the carrier's run size did not survive the round trip; a fidelity test that cannot see "
        "the operator's formatting cannot detect a writer that destroys it"
    )
    assert (
        paragraph.runs[0].font.bold is True
    ), "the carrier's bold did not survive the round trip — same consequence as the size above"
    assert (
        paragraph._p.find(qn("a:pPr")) is not None
    ), "the carrier's paragraph properties vanished; the injected bullet lives inside them"
    assert paragraph._p.find(qn("a:pPr")).find(qn("a:buChar")) is not None, (
        "the injected buChar bullet is gone. It stands in for operator formatting python-pptx has "
        "no API for; without it, nothing proves the writer INHERITS bullets rather than dropping "
        "them (02-RESEARCH W3)"
    )

    # FALSIFIABILITY CONTROL (P-05 applied to the in-test decks). The built template must NOT
    # already carry the marker, the gate state or the writer's pinned timestamps — otherwise plan
    # 02-02's read-back assertion would pass on a deck the writer never touched.
    core = written.core_properties
    assert core.category != MARKER, (
        "the test deck already carries the generated-by marker, so 02-02's read-back assertion "
        "would be green even if the writer never wrote it"
    )
    assert core.content_status != DRAFT_STATUS, (
        "the test deck already carries the gate state, so the contentStatus read-back would be "
        "green even if the writer never wrote it"
    )
    assert core.created != EPOCH_ZERO.replace(tzinfo=None), (
        "the test deck's dcterms:created is already EPOCH_ZERO, so the timestamp read-back would "
        "be green even if the writer never pinned it (this is why _FIXED is 2026-01-01)"
    )

    # The builder is a fixed point of the normalizer and is itself deterministic: a test deck that
    # drifted between two builds would inject non-determinism into every determinism assertion
    # that renders through it.
    raw = path.read_bytes()
    assert normalize_opc_zip(raw) == raw, (
        "the built deck is not its own normalization — the builder must route bytes through "
        "normalize_opc_zip BEFORE writing, or every deck it produces is a moving target"
    )
    second_dir = tmp_path / "again"
    second_dir.mkdir()
    again = build_rich_template(second_dir)
    assert part_digest(again.read_bytes()) == part_digest(raw), (
        "two builds of the rich deck have different CONTENT digests — a clock, a random value or "
        "an unstable ordering has entered the builder"
    )


def test_group_nesting_hides_a_slot_from_slide_shapes(tmp_path: pathlib.Path) -> None:
    """W17, re-proved in-repo: `slide.shapes` does NOT descend into groups.

    This is the fact the group-recursive binding map in plan 02-02 exists for, so it is owned here
    as a repo fact rather than cited as a research claim. If python-pptx ever changed this, the
    recursion would become dead code and this test — not a mystery in an operator's deck — is what
    would say so. The operator's failure mode without the recursion: they can SEE the box in the
    Selection Pane, and the renderer calls its name unknown.
    """
    path = build_rich_template(tmp_path)
    written = Presentation(str(path))
    slide = written.slides[0]

    top_level = [shape.name for shape in slide.shapes]
    assert "NL_INSIDE_GROUP" not in top_level, (
        "`slide.shapes` now yields shapes nested inside groups. If that is real, the `_walk` "
        f"recursion in the writer is dead code and the duplicate-name check may double-count. "
        f"Top level: {top_level}"
    )
    assert "Narrative Group" in top_level, (
        "the group container is not a top-level shape — the deck builder no longer groups the "
        "inner slot, so this test is not testing group nesting at all"
    )
    assert "NL_INSIDE_GROUP" in _all_shape_names(written), (
        "the grouped slot is invisible to the RECURSIVE walk too, which means the builder lost it "
        "rather than nesting it — the negative assertion above would then be vacuous"
    )

    group = next(shape for shape in slide.shapes if shape.name == "Narrative Group")
    assert group.has_text_frame is False, (
        "a group shape now reports has_text_frame True; the writer's 'this slot cannot hold text' "
        "teaching error is built on the measured opposite (a group and a picture have no text "
        "frame, so `.text_frame` raises AttributeError)"
    )
    assert not group.name.startswith(SLOT_PREFIX), (
        "the group container acquired the reserved NL_ prefix — it would then be an unfilled "
        "reserved slot in every happy-path render, and the fail-loud contract would reject the "
        "deck it is supposed to accept"
    )


# --- SC-1: the binding map, and the five ambiguities it refuses ---------------------------------
#
# Every test below asserts the OFFENDER'S NAME is in the message. A refusal that does not name the
# offender is not a teaching error: it tells the operator that something is wrong with a deck they
# authored in PowerPoint and leaves them to bisect it by hand.


def test_group_nested_slot_is_bound(tmp_path: pathlib.Path) -> None:
    """The positive half of W17: a slot inside a group is BOUND, not silently dropped.

    `test_group_nesting_hides_a_slot_from_slide_shapes` above proves the flat view cannot see
    `NL_INSIDE_GROUP`. This proves `bind_slots` can — which is the whole reason `_walk` recurses.
    Without the recursion this deck would fail twice over: the content would be rejected as an
    unknown name, and "remove it from the content mapping" would then ship an empty box.
    """
    path = build_rich_template(tmp_path)
    prs = Presentation(str(path))

    bound = bind_slots(prs, RICH_SLOTS)

    assert "NL_INSIDE_GROUP" in bound, (
        "the group-nested slot is missing from the binding map — `_walk` stopped recursing into "
        f"groups, so an operator's grouped box is invisible to the renderer. Bound: {sorted(bound)}"
    )
    assert set(bound) == RICH_SHAPE_NAMES, (
        "the binding map is not the deck's full recursive shape inventory. Both fail-loud "
        f"directions are computed from it, so a gap here is a silent drop. Bound: {sorted(bound)}"
    )
    assert bound["NL_INSIDE_GROUP"].has_text_frame is True, (
        "the grouped slot bound to something that cannot hold text — `_walk` is yielding the group "
        "container instead of its member"
    )


def test_unprefixed_shape_is_not_a_slot(tmp_path: pathlib.Path) -> None:
    """The reserved prefix discriminates: the operator's own shapes are bound but never demanded.

    This is what makes D-03 usable on a real deck. Without the `NL_` prefix, direction (b) of the
    fail-loud contract would reject every operator logo, footer and page number as an unfilled slot.
    `Footer` and `Narrative Group` are bound (so a duplicate of either is still caught) and neither
    triggers the unfilled refusal.
    """
    path = build_rich_template(tmp_path)
    prs = Presentation(str(path))

    bound = bind_slots(
        prs, RICH_SLOTS
    )  # does NOT raise, though two shapes carry no content

    for unprefixed in ("Footer", "Narrative Group"):
        assert unprefixed in bound, (
            f"{unprefixed!r} is not in the binding map — unprefixed shapes must still be BOUND "
            "(the map is the deck's full recursive inventory; unprefixed shapes are visible in "
            "it, they are just never demanded, filled, or subject to the duplicate refusal)"
        )
        assert unprefixed not in RICH_SLOTS, (
            f"{unprefixed!r} acquired content, so this test no longer proves that an unprefixed "
            "shape can be left unfilled"
        )
        assert not unprefixed.startswith(SLOT_PREFIX), (
            f"{unprefixed!r} now carries the reserved prefix; the discrimination this test asserts "
            "has been lost in the deck builder"
        )


def test_unknown_slot_name_raises(tmp_path: pathlib.Path) -> None:
    """Fail-loud direction (a): content bound to a name the template does not contain.

    The message must name the offender AND list the names the template does contain — otherwise
    the operator's next move is a guess. `never invents layout` is the sentence that tells them the
    fix is theirs (add or rename the shape), not the renderer's (fill by position).
    """
    path = build_rich_template(tmp_path)
    prs = Presentation(str(path))
    slots = {**RICH_SLOTS, "NL_NOPE": ["content for a slot that does not exist"]}

    with pytest.raises(ValueError, match="NL_NOPE") as caught:
        bind_slots(prs, slots)

    message = str(caught.value)
    assert "NL_WEEK_TITLE" in message, (
        "the refusal does not list the names the template DOES contain, so the operator cannot see "
        f"whether they mistyped a name or omitted a shape. Message: {message}"
    )
    assert "never invents layout" in message, (
        "the refusal no longer states why the renderer will not resolve this itself (D-03). "
        f"Message: {message}"
    )


def test_content_bound_to_unprefixed_name_raises(tmp_path: pathlib.Path) -> None:
    """WR-06: content bound to an UNPREFIXED shape name is refused, never silently filled.

    `Footer` exists in the deck, so every other refusal passes it — before this fix the render
    would overwrite the operator's own footer with renderer text, the exact "operator's
    logo/footer modified" failure the committed-template test asserts against, with no teaching
    error. The `NL_` prefix must discriminate in BOTH directions: unprefixed shapes are never
    demanded AND never written.
    """
    path = build_rich_template(tmp_path)
    prs = Presentation(str(path))
    slots = {**RICH_SLOTS, "Footer": ["renderer text targeting the operator's footer"]}

    with pytest.raises(ValueError, match="Footer") as caught:
        bind_slots(prs, slots)

    message = str(caught.value)
    assert "never written" in message, (
        "the refusal no longer states the contract (unprefixed shapes are the operator's and are "
        f"never written), which is what tells the composer the bug is theirs. Message: {message}"
    )
    assert "reserved prefix" in message, (
        "the refusal no longer offers the fix (rename the shape with the reserved prefix if it is "
        f"meant to be a slot). Message: {message}"
    )


def test_unfilled_reserved_slot_raises(tmp_path: pathlib.Path) -> None:
    """Fail-loud direction (b): an `NL_` slot with no matching content, named in the message.

    A reserved-prefix slot left empty would ship a blank box to a reader — the renderer's version
    of publishing something nobody wrote.
    """
    path = build_rich_template(tmp_path)
    prs = Presentation(str(path))
    slots = {k: v for k, v in RICH_SLOTS.items() if k != "NL_LOWLIGHTS"}

    with pytest.raises(ValueError, match="NL_LOWLIGHTS") as caught:
        bind_slots(prs, slots)

    assert "blank box" in str(caught.value), (
        "the refusal no longer says what the operator would have shipped; the consequence is the "
        f"reason this direction exists. Message: {caught.value}"
    )


def test_duplicate_shape_name_raises(tmp_path: pathlib.Path) -> None:
    """A template with two shapes of one name is REFUSED, not last-wins-dropped (T-02-04).

    Duplicate names are legal in OOXML and copy-paste is how an operator makes one. The naive
    comprehension would bind the second and silently discard the first — the deck would render,
    with one of the operator's slots quietly empty.
    """
    path = build_duplicate_name_template(tmp_path)
    prs = Presentation(str(path))

    with pytest.raises(ValueError, match="two shapes named 'NL_WEEK_TITLE'") as caught:
        bind_slots(prs, {"NL_WEEK_TITLE": ["a line"]})

    assert "Selection Pane" in str(caught.value), (
        "the refusal no longer tells the operator WHERE to rename the shape (Alt+F10); a teaching "
        f"error that stops at the diagnosis is half an error. Message: {caught.value}"
    )


def test_default_auto_named_multi_slide_deck_is_accepted() -> None:
    """WR-02's regression test: PowerPoint's per-slide auto-names must not be refused.

    PowerPoint (and python-pptx) auto-name shapes PER SLIDE — "TextBox 1" on slide 1 and
    "TextBox 1" on slide 2 — so ANY real two-slide operator deck with default names carries an
    unprefixed duplicate. The Phase-2 review reproduced the old deck-wide refusal rejecting such a
    deck even with an EMPTY slots mapping. The duplicate refusal is scoped to `NL_`-prefixed names
    (the ones a last-wins map would actually silently drop); unprefixed decorative duplicates are
    the operator's business and bind first-seen, deterministically.
    """
    prs = Presentation()
    _scrub(prs)
    for index in range(2):
        slide = _blank_slide(prs)
        box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.4), Inches(9.0), Inches(1.0)
        )
        box.text_frame.text = f"decorative default-named box on slide {index}"

    names = [shape.name for slide in prs.slides for shape in slide.shapes]
    assert len(names) == 2 and len(set(names)) == 1, (
        "python-pptx no longer auto-names shapes per slide, so this deck carries no collision "
        f"and the test is not testing the WR-02 case at all. Names: {names}"
    )

    bound = bind_slots(prs, {})  # must NOT raise — this is an ordinary operator deck

    assert names[0] in bound, (
        "the auto-named shape vanished from the binding map entirely — first-seen binding for "
        f"unprefixed duplicates has broken. Bound: {sorted(bound)}"
    )


def test_template_owning_watermark_name_raises(tmp_path: pathlib.Path) -> None:
    """A template that already defines `NL_DRAFT_WATERMARK` is refused (T-02-05 / W14).

    Measured: adding the watermark onto such a deck writes TWO shapes with one name and raises
    nothing — so a "helpful" operator who copies the watermark into their template would silently
    defeat the binding map on the next render.
    """
    path = build_watermark_owning_template(tmp_path)
    prs = Presentation(str(path))

    with pytest.raises(ValueError, match=WATERMARK_NAME) as caught:
        bind_slots(prs, {"NL_WEEK_TITLE": ["a line"]})

    assert "owned by the renderer" in str(caught.value), (
        "the refusal no longer explains that the name is reserved, which is the one thing the "
        f"operator needs to know to fix their template. Message: {caught.value}"
    )


def test_slot_without_text_frame_raises(tmp_path: pathlib.Path) -> None:
    """A slot that cannot hold text gets a teaching error, never an `AttributeError` (T-02-13).

    A group (like a picture) reports `has_text_frame == False`. Without this check the caller gets
    a stack trace from deep inside python-pptx; with it, they get told what the shape is and what
    to do instead.
    """
    path = build_nontext_slot_template(tmp_path)
    prs = Presentation(str(path))
    slots = {"NL_WEEK_TITLE": ["a line"], "NL_NOT_TEXT": ["content for a group"]}

    with pytest.raises(ValueError, match="NL_NOT_TEXT") as caught:
        bind_slots(prs, slots)

    message = str(caught.value)
    assert (
        "cannot hold text" in message
    ), f"the refusal no longer says what is wrong with the shape. Message: {message}"
    assert "text box" in message and "asset" in message, (
        "the refusal no longer offers the two things the operator can actually do (point the slot "
        f"at a text box, or place the content as an asset). Message: {message}"
    )


# --- Formatting fidelity: the only automated way to see 02-RESEARCH Pitfall 1 -------------------
#
# Every test below fills, SAVES, NORMALIZES and REOPENS THE WRITTEN BYTES before asserting. Never
# assert against the in-memory object the fill just mutated: that object would report the operator's
# formatting intact even if the write dropped it, which is precisely the failure being hunted.


def _fill_and_reread(path: pathlib.Path, slots) -> "Presentation":
    """Bind + fill every slot on the deck at `path`, then reopen the WRITTEN, normalized bytes."""
    prs = Presentation(str(path))
    bound = bind_slots(prs, slots)
    for name, lines in slots.items():
        fill_slot(bound[name].text_frame, lines)

    buf = io.BytesIO()
    prs.save(buf)
    return Presentation(io.BytesIO(normalize_opc_zip(buf.getvalue())))


def _paragraphs(prs, name: str):
    """The paragraphs of the named slot, found through the RECURSIVE walk."""
    shape = next(
        candidate
        for slide in prs.slides
        for candidate in _walk(slide.shapes)
        if candidate.name == name
    )
    return list(shape.text_frame.paragraphs)


def test_fill_preserves_operator_formatting_on_every_line(
    tmp_path: pathlib.Path,
) -> None:
    """Three lines into a 20pt bold slot read back as THREE 20pt bold lines (W3 / W5).

    `NL_HIGHLIGHTS` carries three lines, so this exercises the clone path — paragraph 0 is reused
    and paragraphs 1..n are deep copies of it. The size is asserted EXPLICITLY rather than "is not
    None" because 02-RESEARCH Pitfall 1 names `font.size is None` as the only automated way to see
    the silent-downgrade failure: the deck still opens, still passes every byte and digest
    assertion, and has quietly lost the operator's formatting.
    """
    path = build_rich_template(tmp_path)
    written = _fill_and_reread(path, RICH_SLOTS)

    paragraphs = _paragraphs(written, "NL_HIGHLIGHTS")
    assert len(paragraphs) == 3, (
        "the three filled lines did not survive as three paragraphs — the clone chain lost a line "
        f"or `tf.text` collapsed the frame. Found {len(paragraphs)}"
    )
    for index, paragraph in enumerate(paragraphs):
        assert paragraph.runs, (
            f"paragraph {index} came back with no run at all; there is nothing left carrying the "
            "operator's formatting"
        )
        run = paragraph.runs[0]
        assert run.font.size == Pt(20), (
            f"paragraph {index}'s run size read back as {run.font.size!r}, not Pt(20)/254000 EMU. "
            "A None here means the fill STOPPED INHERITING the operator's formatting: the deck "
            "still renders, with their 20pt type silently downgraded to the theme default. That is "
            "a visual regression no other check in this repo can see (02-RESEARCH Pitfall 1)."
        )
        assert run.font.bold is True, (
            f"paragraph {index}'s bold read back as {run.font.bold!r} — same silent-downgrade "
            "failure as the size above, same invisible-to-everything-else consequence"
        )


def test_fill_preserves_bullet_on_every_line(tmp_path: pathlib.Path) -> None:
    """The `buChar` bullet survives on EVERY filled line, read back off the written bytes (W6).

    python-pptx has no bullet API, so a bullet the fill dropped could never be put back by the
    writer — it can only be inherited. This is the paragraph-level half of the fidelity claim (the
    test above is the run-level half): `tf.clear()` + `add_paragraph()` would keep p0's bullet and
    lose it on lines 2..n, which is exactly the half-preserved deck the primitive exists to avoid.
    """
    path = build_rich_template(tmp_path)
    written = _fill_and_reread(path, RICH_SLOTS)

    for index, paragraph in enumerate(_paragraphs(written, "NL_HIGHLIGHTS")):
        properties = paragraph._p.find(qn("a:pPr"))
        assert properties is not None, (
            f"paragraph {index} has no pPr at all — the clone did not carry the operator's "
            "paragraph properties, so its bullet, level and spacing are all gone"
        )
        assert properties.find(qn("a:buChar")) is not None, (
            f"paragraph {index} lost its buChar bullet. python-pptx has NO bullet API, so a bullet "
            "the fill drops cannot be restored by the writer — it can only be inherited (W6)"
        )


def test_unicode_and_xml_metacharacters_roundtrip(tmp_path: pathlib.Path) -> None:
    """`café — "smart" … 🎯 <&>` reads back character for character (W11, threat T-02-14).

    Slot content crosses from authored data into XML text nodes. python-pptx escapes it on the way
    in and unescapes it on the way out; this asserts the round trip rather than assuming it, and it
    is on the HAPPY path (the line lives in `RICH_SLOTS`) rather than in a special case nobody runs.
    """
    path = build_rich_template(tmp_path)
    written = _fill_and_reread(path, RICH_SLOTS)

    read_back = [
        paragraph.runs[0].text for paragraph in _paragraphs(written, "NL_HIGHLIGHTS")
    ]
    assert read_back == RICH_SLOTS["NL_HIGHLIGHTS"], (
        "the filled lines did not round-trip exactly. If only the metacharacter line differs, XML "
        f"escaping is the suspect; if the order differs, the clone chain is. Read back: {read_back}"
    )


def test_empty_run_slot_raises(tmp_path: pathlib.Path) -> None:
    """A slot whose paragraph 0 has ZERO runs raises a teaching error (Pitfall 6).

    This is an OBSERVED failure, not defensive padding: a text box the operator "left blank on
    purpose" has no typed characters and therefore no run, so the primitive has nothing to inherit
    from. Without the raise the caller gets an `IndexError` from `paragraph.runs[0]`.
    """
    path = build_empty_run_template(tmp_path)
    prs = Presentation(str(path))
    bound = bind_slots(prs, {"NL_WEEK_TITLE": ["a line"]})

    with pytest.raises(ValueError, match="carries no run") as caught:
        fill_slot(bound["NL_WEEK_TITLE"].text_frame, ["a line"])

    assert "one character of placeholder text" in str(caught.value), (
        "the refusal no longer tells the operator the one thing that fixes it (type a character "
        f"into the shape and re-save the template). Message: {caught.value}"
    )


def test_fill_with_no_lines_raises(tmp_path: pathlib.Path) -> None:
    """An empty line list is refused: an empty slot ships a blank box to a reader."""
    path = build_rich_template(tmp_path)
    prs = Presentation(str(path))
    bound = bind_slots(prs, RICH_SLOTS)

    with pytest.raises(ValueError, match="no lines"):
        fill_slot(bound["NL_WEEK_TITLE"].text_frame, [])


def test_fill_with_only_blank_lines_raises(tmp_path: pathlib.Path) -> None:
    """IN-02: `[""]` (and whitespace-only lines) are refused like `[]` — same blank box.

    The empty-list refusal exists because "an empty slot ships a blank box to a reader";
    reproduced live in the review, `[""]` passed and shipped the identical visually blank box.
    If the blank box is the harm, the refusal covers the equivalent input — the asymmetry was an
    accident, not a spacer allowance.
    """
    path = build_rich_template(tmp_path)
    prs = Presentation(str(path))
    bound = bind_slots(prs, RICH_SLOTS)

    with pytest.raises(ValueError, match="only blank/whitespace"):
        fill_slot(bound["NL_WEEK_TITLE"].text_frame, [""])
    with pytest.raises(ValueError, match="only blank/whitespace"):
        fill_slot(bound["NL_WEEK_TITLE"].text_frame, ["   ", "\t"])
    # A blank SPACER line among real content is still legitimate — only all-blank is refused.
    fill_slot(bound["NL_WEEK_TITLE"].text_frame, ["real content", ""])


def test_bare_str_fills_one_paragraph_not_one_per_character(
    tmp_path: pathlib.Path,
) -> None:
    """WR-03, half one: `fill_slot(tf, "hello")` is ONE paragraph, never `['h','e','l','l','o']`.

    `str` IS a `Sequence[str]` of its characters, so neither mypy nor any refusal used to catch
    the single most likely caller typo — reproduced live in the review as five one-character
    paragraphs, a silently misrendered deck. A bare string means one line; it is normalized to
    `[line]` rather than refused, because the input has exactly one sane reading.
    """
    path = build_rich_template(tmp_path)
    prs = Presentation(str(path))
    bound = bind_slots(prs, RICH_SLOTS)

    fill_slot(bound["NL_WEEK_TITLE"].text_frame, "2026-W35")

    paragraphs = list(bound["NL_WEEK_TITLE"].text_frame.paragraphs)
    assert len(paragraphs) == 1, (
        f"a bare str filled {len(paragraphs)} paragraphs — the per-character explosion is back. "
        f"Texts: {[p.runs[0].text if p.runs else None for p in paragraphs]}"
    )
    assert paragraphs[0].runs[0].text == "2026-W35", paragraphs[0].runs[0].text


def test_bare_str_slot_value_renders_one_paragraph(tmp_path: pathlib.Path) -> None:
    """WR-03, half two: a bare-str VALUE in the `slots` mapping renders as one paragraph.

    `render_surface_pptx_bytes` is the entry point a composer actually calls, so the same typo is
    proven safe end to end, asserted off the WRITTEN bytes per this module's fidelity convention.
    """
    path = build_rich_template(tmp_path)
    str_slots: dict[str, Union[str, list[str]]] = {
        **RICH_SLOTS,
        "NL_WEEK_TITLE": "2026-W35",
    }

    rendered = render_surface_pptx_bytes(
        _sample_weekly_surface(), template=path, slots=str_slots
    )

    written = Presentation(io.BytesIO(rendered))
    paragraphs = _paragraphs(written, "NL_WEEK_TITLE")
    assert len(paragraphs) == 1, (
        f"a bare-str slot value shipped {len(paragraphs)} paragraphs through the render path. "
        f"Texts: {[p.runs[0].text if p.runs else None for p in paragraphs]}"
    )
    assert paragraphs[0].runs[0].text == "2026-W35", paragraphs[0].runs[0].text


# --- SC-3 (marker + watermark) and SC-4 (the gate is untouched) ---------------------------------


def _sample_weekly_surface(state: ReviewState = ReviewState.DRAFT) -> Surface:
    """A fabricated `Surface(REPORT)` — D-01: the weekly REUSES the Report, no new semantic kind.

    Built from EXISTING block kinds: the four weekly kinds do not exist yet (Phase 3 owns them), and
    the writer takes an explicit `slots` mapping anyway (P-03), so the blocks here only have to make
    the Surface real — they are not the deck's content.

    For the PUBLISHED case the `Review` must satisfy the policy validator or the model refuses to
    construct at all (`_published_requires_satisfied_policy` — there is no auto-publish path, even
    in a test fixture). `REPORT.review_policy` is `light()`: one approval, no peer required, so an
    author plus one approval suffices. Without this, the published half of SC-3 would be untestable
    and "no watermark when Published" would go unasserted.
    """
    review = Review(
        state=state,
        policy=REPORT.review_policy,
        author="fixture author",
        approvals=["fixture reviewer"] if state is ReviewState.PUBLISHED else [],
    )
    return Surface(
        id="weekly-2026-w35",
        template=REPORT,
        title="Week 35 — fabricated sample",
        blocks=[
            ProseBlock(
                heading="Narrative", text="A fabricated week, written for a test."
            ),
            ClaimsBlock(
                claims=[
                    Claim(
                        text="the checklist got shorter",
                        evidence=[Trace(source_id="fixture-source")],
                    )
                ]
            ),
        ],
        review=review,
    )


def _watermarks(prs) -> list:
    """Every watermark-named shape in the deck, found through the RECURSIVE walk."""
    return [
        shape
        for slide in prs.slides
        for shape in _walk(slide.shapes)
        if shape.name == WATERMARK_NAME
    ]


def test_marker_reads_back_off_the_written_file(tmp_path: pathlib.Path) -> None:
    """SC-3: the generated-by marker, asserted by reopening the WRITTEN FILE.

    The assertion block is the decision note's, verbatim. It reopens `out_path` with
    `Presentation(str(out_path))` — never the writer's return value and never the in-memory object:
    a writer that returned success without marking would pass any assertion made against itself.
    The falsifiability control that makes this red-able is that the TEMPLATE carries none of these
    values (`_FIXED` is 2026-01-01, `category`/`content_status` are empty — P-05).
    """
    template = build_rich_template(tmp_path)
    surface = _sample_weekly_surface()
    out_path = tmp_path / "weekly.pptx"

    render_surface_pptx(surface, template=template, slots=RICH_SLOTS, out_path=out_path)

    written = Presentation(str(out_path))  # reopen the FILE that was written
    cp = written.core_properties
    assert cp.category == MARKER, cp.category
    assert cp.content_status == DRAFT_STATUS, cp.content_status
    # dcterms:created serializes as W3CDTF and reads back tz-NAIVE (02-RESEARCH Pitfall 8).
    assert cp.created == EPOCH_ZERO.replace(tzinfo=None), cp.created
    assert cp.modified == EPOCH_ZERO.replace(tzinfo=None), cp.modified
    assert cp.identifier == surface.id, cp.identifier


def test_draft_surface_is_watermarked_on_every_slide(tmp_path: pathlib.Path) -> None:
    """SC-3: a Draft deck carries the watermark on EVERY slide, last in z-order.

    The rich deck has two slides precisely so "every slide" is not a vacuous claim: a writer that
    watermarked `prs.slides[0]` only would ship a second page that reads as approved.
    """
    template = build_rich_template(tmp_path)
    out_path = tmp_path / "draft.pptx"

    render_surface_pptx(
        _sample_weekly_surface(),
        template=template,
        slots=RICH_SLOTS,
        out_path=out_path,
    )

    written = Presentation(str(out_path))
    assert len(written.slides) == 2, "the template lost a slide in the render"
    for index, slide in enumerate(written.slides):
        names = [shape.name for shape in slide.shapes]
        assert WATERMARK_NAME in names, (
            f"slide {index} of a DRAFT deck carries no watermark — an unreviewed page that does "
            f"not look unreviewed. Shapes: {names}"
        )
        assert names[-1] == WATERMARK_NAME, (
            f"the watermark is not LAST in slide {index}'s shape order, so it is not at the top of "
            f"z-order and the operator's own boxes can paint over it. Shapes: {names}"
        )
    for shape in _watermarks(written):
        assert shape.rotation == 315.0, (
            f"the watermark's rotation read back as {shape.rotation!r}, not the fixed 315.0 — a "
            "watermark property that is not a literal is a determinism risk"
        )


def test_published_surface_has_no_watermark_and_empty_content_status(
    tmp_path: pathlib.Path,
) -> None:
    """SC-3's INVERTED half. Without it, "while the Surface is not Published" is half-asserted.

    A writer that watermarked unconditionally would pass every Draft assertion above and still be
    wrong — it would brand approved work as unreviewed forever.
    """
    template = build_rich_template(tmp_path)
    surface = _sample_weekly_surface(ReviewState.PUBLISHED)
    assert (
        surface.is_published
    ), "the fixture is not actually Published; the inversion is vacuous"
    out_path = tmp_path / "published.pptx"

    render_surface_pptx(surface, template=template, slots=RICH_SLOTS, out_path=out_path)

    written = Presentation(str(out_path))
    assert not _watermarks(written), (
        "a PUBLISHED deck carries the Draft watermark — the gate state is being ignored, so the "
        "watermark says nothing about whether a human approved this"
    )
    assert written.core_properties.content_status == "", (
        "a PUBLISHED deck still reports a draft contentStatus; the field is written FROM the gate, "
        f"so this means it is not. Found: {written.core_properties.content_status!r}"
    )


def test_render_does_not_touch_the_gate(tmp_path: pathlib.Path) -> None:
    """SC-4: rendering a Draft Surface leaves it Draft and `model_dump()`-identical.

    The product's hardest rule is that nothing publishes without a human. The writer READS
    `is_published`; there is no assignment to any `Surface` field. This is the positive proof of
    that absence — the negative proof is `git diff --exit-code -- src/newsletters/semantic.py`.
    """
    template = build_rich_template(tmp_path)
    surface = _sample_weekly_surface()
    before = surface.model_dump()

    render_surface_pptx(
        surface, template=template, slots=RICH_SLOTS, out_path=tmp_path / "gate.pptx"
    )

    assert surface.review.state is ReviewState.DRAFT, (
        "rendering moved the review gate. There must be no write path from the renderer to the "
        f"gate at all. State is now: {surface.review.state!r}"
    )
    assert surface.model_dump() == before, (
        "the Surface is not bit-for-bit what it was before the render. Even a benign-looking "
        "mutation here means the writer has a write path into the reviewed record"
    )


def test_render_surface_pptx_writes_the_same_bytes_it_returns(
    tmp_path: pathlib.Path,
) -> None:
    """The disk path is ONE write of complete, already-normalized bytes.

    `prs.save(path)` followed by a rewrite in place would leave an un-normalized deck on disk if the
    normalization raised, and nothing downstream could tell. Asserting the file equals the in-memory
    render is how that shape stays enforced rather than merely intended.
    """
    template = build_rich_template(tmp_path)
    surface = _sample_weekly_surface()
    out_path = tmp_path / "written.pptx"

    returned = render_surface_pptx(
        surface, template=template, slots=RICH_SLOTS, out_path=out_path
    )
    in_memory = render_surface_pptx_bytes(surface, template=template, slots=RICH_SLOTS)

    assert (
        returned == out_path
    ), f"the returned path is not the one asked for: {returned}"
    assert out_path.read_bytes() == in_memory, (
        "the bytes on disk differ from the bytes the byte-returning entry point produces for the "
        "same inputs — the two entry points have drifted, or the disk path is doing more than one "
        "write"
    )


# --- SC-2: the determinism battery, with its negative control -----------------------------------
#
# THE CONSTRAINT THAT GOVERNS EVERY ASSERTION BELOW: **no test here may compare a rendered deck to
# the template** — not its bytes, not its `part_digest`, not its part order. Two measured properties
# of python-pptx's LOAD path make any such comparison a false red:
#
#   * opening the committed template and saving it UNCHANGED already yields a different
#     `part_digest`. The `""`-valued core properties serialize as `<cp:keywords></cp:keywords>` and
#     come back as `<cp:keywords/>` — semantically identical XML, different bytes.
#   * part EMISSION order differs between a freshly built package and a reopened one, even when
#     `_rels/.rels` is byte-identical.
#
# Neither is a regression and neither is anything the writer can or should "fix". Every assertion in
# this battery therefore compares a render to ANOTHER RENDER. The consequence Phase 4 inherits: its
# golden deck must be produced by THE WRITER, never by re-saving the template.


class _Renders(NamedTuple):
    """One Surface rendered twice across a real time boundary, in both forms.

    The NORMALIZED pair is what the writer returns and what ships. The RAW pair is the writer's own
    pre-normalization bytes, which only the negative control needs — and which it needs in order to
    be a control at all.
    """

    normalized_a: bytes
    normalized_b: bytes
    raw_a: bytes
    raw_b: bytes


@pytest.fixture(scope="module")
def time_separated_renders(tmp_path_factory: pytest.TempPathFactory) -> _Renders:
    """Two renders of ONE Surface through ONE template, separated by a real 3-second gap.

    **The sleep is load-bearing, not incidental.** DOS timestamps inside a ZIP have 2-SECOND
    granularity, so two writes landing in a single wall-clock second are ALREADY byte-identical: a
    tight-loop double render would "prove" a stability that does not exist, and the whole
    determinism claim would rest on luck. Three seconds guarantees the boundary is crossed. The
    fixture is MODULE-scoped so the suite sleeps once, not once per assertion.

    HOW THE RAW PAIR IS OBTAINED, and why it is the same code path. `render_surface_pptx_bytes`
    returns normalized bytes by construction, so the raw pair is intercepted one call earlier: the
    module-level `normalize_opc_zip` the writer calls is temporarily wrapped, and the wrapper records
    the payload it was handed — i.e. exactly what `prs.save(BytesIO)` produced inside the writer —
    before handing it to the real normalizer. The recorded bytes are therefore the writer's own
    output and are **not** normalized. That matters: a raw pair that had accidentally routed through
    the normalizer would be byte-equal for the wrong reason, and the negative control below would be
    vacuously green. Rebuilding an "identical" presentation in the fixture instead would be the other
    failure mode — a second implementation of the writer, drifting from the first.
    """
    directory = tmp_path_factory.mktemp("determinism")
    template = build_rich_template(directory)
    surface = _sample_weekly_surface()

    raw: list[bytes] = []

    def _capture_then_normalize(payload: bytes) -> bytes:
        # `normalize_opc_zip` here is THIS module's import of the real function, bound at import
        # time, so the patch below cannot make this wrapper recurse into itself.
        raw.append(payload)
        return normalize_opc_zip(payload)

    with pytest.MonkeyPatch.context() as patch:
        patch.setattr(pptx_writer, "normalize_opc_zip", _capture_then_normalize)
        normalized_a = render_surface_pptx_bytes(
            surface, template=template, slots=RICH_SLOTS
        )
        time.sleep(3)
        normalized_b = render_surface_pptx_bytes(
            surface, template=template, slots=RICH_SLOTS
        )

    assert len(raw) == 2, (
        "the writer no longer routes its saved bytes through the module-level `normalize_opc_zip`, "
        f"so this fixture captured {len(raw)} raw payload(s) instead of 2. Until the capture is "
        "retargeted at whatever replaced it, the negative control below is not measuring the "
        "writer's un-normalized output and proves nothing."
    )
    return _Renders(normalized_a, normalized_b, raw[0], raw[1])


def test_double_render_is_byte_identical(time_separated_renders: _Renders) -> None:
    """(B) SC-2: two renders of the same Surface, three seconds apart, are byte-identical.

    The failure message diagnoses itself off the RAW pair, because the two lists separate the two
    possible causes: a differing PART means the deck's content moved (a clock, an unstable part
    order or an unstable rel id leaked into the writer, and no zip normalization can repair that);
    a differing zip FIELD beyond `date_time` means the normalizer is pinning too little.
    """
    renders = time_separated_renders

    assert renders.normalized_a == renders.normalized_b, (
        "two renders of the SAME Surface, 3s apart, are NOT byte-identical — the recorded "
        "determinism definition (byte-stable via post-save zip normalization) does not hold for "
        "the writer on this python-pptx/zlib pair. Differing parts (un-normalized): "
        f"{differing_parts(renders.raw_a, renders.raw_b)}. Differing zip fields (un-normalized): "
        f"{differing_zipinfo_fields(renders.raw_a, renders.raw_b)}."
    )


def test_unnormalized_double_render_is_not_byte_equal(
    time_separated_renders: _Renders,
) -> None:
    """(C) THE NEGATIVE CONTROL — the byte-equality above must be able to FAIL.

    DO NOT DELETE THIS AS "redundant now that the writer exists". It is the assertion most likely to
    be dropped in a tidy-up, and dropping it silently converts `test_double_render_is_byte_identical`
    from a determinism proof into a test that is green whenever two writes land in the same second.
    """
    renders = time_separated_renders

    assert renders.raw_a != renders.raw_b, (
        "the negative control has stopped controlling: two UN-normalized renders 3s apart are "
        "byte-equal. Either both renders landed inside one DOS timestamp tick, or python-pptx "
        "stopped stamping the clock, or the raw capture accidentally routed through the "
        "normalizer. Until this inequality holds, the byte equality in "
        "test_double_render_is_byte_identical is NOT attributable to the normalizer and proves "
        "nothing."
    )
    assert differing_parts(renders.raw_a, renders.raw_b) == [], (
        "unzipped part CONTENT differs across two renders of identical input — that is a real "
        "non-determinism in the writer (a clock, an unstable part order, or an unstable rel id), "
        "not a zip-metadata artefact the normalizer can fix"
    )
    assert differing_zipinfo_fields(renders.raw_a, renders.raw_b) == ["date_time"], (
        "a zip metadata field other than date_time drifted across two renders; the normalizer pins "
        "date_time, compress_type, create_system and external_attr, so a new offender here means "
        "the recorded determinism mechanism is incomplete"
    )


def test_part_digest_is_stable_across_time_separated_renders(
    time_separated_renders: _Renders,
) -> None:
    """(A) CONTENT IDENTITY, computed on the UN-normalized bytes.

    Deliberately not normalized first: the point is that content identity holds WITHOUT any zip
    metadata fix. This is the implementation-INDEPENDENT assertion Phase 4's committed==fresh gate
    inherits, because DEFLATE output is zlib-implementation-dependent (zlib vs zlib-ng) — a
    full-file hash there would be green locally and red in CI with byte-identical part content.
    """
    renders = time_separated_renders

    assert part_digest(renders.raw_a) == part_digest(renders.raw_b), (
        "the part-content digest differs across two renders of identical content — the deck's "
        "CONTENT is non-deterministic, which no zip normalization can repair. This is the "
        "assertion Phase 4's committed==fresh gate is built on, so a red here disqualifies that "
        "gate too"
    )


def test_rendered_archive_is_valid_and_normalization_is_idempotent(
    time_separated_renders: _Renders,
) -> None:
    """A rendered deck is a valid archive, OPC-conventional, and a fixed point of the normalizer.

    Idempotence is what lets any later stage (an assembler, a publisher) re-normalize without
    changing the bytes; without it, a second pass anywhere in the pipeline would silently produce a
    different artifact than the one that was reviewed.
    """
    normalized = time_separated_renders.normalized_a

    assert normalize_opc_zip(normalized) == normalized, (
        "normalize_opc_zip is not idempotent on the writer's output — normalizing an "
        "already-normalized package must be a no-op, or a second pass anywhere in the pipeline "
        "would change the bytes"
    )

    with zipfile.ZipFile(io.BytesIO(normalized)) as archive:
        assert (
            archive.testzip() is None
        ), "the rendered archive fails its own CRC check — the render produced a corrupt deck"
        names = archive.namelist()
    assert names[0] == "[Content_Types].xml", (
        "[Content_Types].xml is no longer the first entry of a rendered deck; the OPC convention "
        "expects it there, and preserving the emitted order is supposed to keep it there by "
        f"construction. Got: {names[0]!r}"
    )


def test_fill_order_does_not_affect_output_bytes(tmp_path: pathlib.Path) -> None:
    """Fill order is NOT a determinism variable — measured (W18), asserted rather than assumed.

    Text edits are in-place on elements that already exist in the operator's template, so iterating
    `slots` in any order writes the same XML. This is why the writer iterates `slots.items()`
    directly instead of sorting: sorting would imply a guarantee that is not needed and would hide
    the ordering variable that IS real — image ADD order, which Phase 3 pins to Weekly Spec file
    order.
    """
    template = build_rich_template(tmp_path)
    surface = _sample_weekly_surface()
    reversed_slots = dict(reversed(list(RICH_SLOTS.items())))
    assert list(reversed_slots) != list(RICH_SLOTS), (
        "the reversed mapping iterates in the same order as the original, so this test is not "
        "varying anything"
    )

    forward = render_surface_pptx_bytes(surface, template=template, slots=RICH_SLOTS)
    backward = render_surface_pptx_bytes(
        surface, template=template, slots=reversed_slots
    )

    assert forward == backward, (
        "the order the slots are filled in changed the output bytes. Either the fill stopped "
        "editing existing elements in place (it is appending or re-creating them), or something "
        "downstream of the fill has become order-sensitive. Differing parts: "
        f"{differing_parts(forward, backward)}"
    )


# --- SC-5: the SHIPPED template renders a real Surface, end to end ------------------------------
#
# Everything above renders through a deck this module built. These two render through
# `tests/fixtures/weekly/template.pptx` — the artifact the repo actually ships, and therefore the
# only one whose rendering proves the phase's claim. Both write ONLY into `tmp_path`: the fixture
# directory stays read-only here (a stray write would break
# `test_pptx_determinism.py::test_weekly_fixture_corpus_is_exactly_the_committed_template`).


def test_sample_surface_renders_through_the_committed_template(
    tmp_path: pathlib.Path,
) -> None:
    """SC-5, and the one place all five phase criteria are asserted against ONE artifact.

    A real `Surface(REPORT, Draft)` goes through the COMMITTED synthetic template and out to disk;
    every assertion below is made by REOPENING THE WRITTEN FILE. In order: the file exists and is a
    valid archive (SC-2's container half); the four slots read back line for line (SC-1); the
    unprefixed `Footer` is untouched, which is what makes the reserved prefix a discriminator rather
    than a decoration; the generated-by marker and the pinned timestamps (SC-3); the watermark is
    present and last in z-order (SC-3's gate mirror); and the Surface is still Draft (SC-4).
    """
    surface = _sample_weekly_surface()
    out_path = tmp_path / "weekly-from-committed-template.pptx"
    template_footer = next(
        shape
        for shape in Presentation(str(COMMITTED_TEMPLATE)).slides[0].shapes
        if shape.name == "Footer"
    ).text_frame.text

    render_surface_pptx(
        surface,
        template=COMMITTED_TEMPLATE,
        slots=COMMITTED_SLOTS,
        out_path=out_path,
    )

    assert out_path.is_file() and out_path.stat().st_size > 0, (
        "the render reported success but wrote no usable file — the disk entry point is the one "
        "an operator actually calls"
    )
    with zipfile.ZipFile(out_path) as archive:
        assert archive.testzip() is None, (
            "the deck rendered from the SHIPPED template fails its own CRC check; whatever the "
            "in-test decks prove, this is the artifact the repo tells operators to start from"
        )

    written = Presentation(str(out_path))  # reopen the FILE that was written
    for name, lines in COMMITTED_SLOTS.items():
        read_back = [paragraph.runs[0].text for paragraph in _paragraphs(written, name)]
        assert read_back == lines, (
            f"slot {name!r} did not read back as written. If a line is missing the clone chain "
            f"dropped it; if the order differs the clone chain reversed it. Read back: {read_back}"
        )

    footer = next(shape for shape in written.slides[0].shapes if shape.name == "Footer")
    assert footer.text_frame.text == template_footer, (
        "the renderer modified `Footer`, which carries no `NL_` prefix and is therefore NOT a slot. "
        "An operator's logo, footer and page numbers must survive a render untouched, or the "
        f"reserved prefix is not discriminating. Found: {footer.text_frame.text!r}"
    )

    cp = written.core_properties
    assert cp.category == MARKER, cp.category
    assert cp.content_status == DRAFT_STATUS, cp.content_status
    # dcterms:created serializes as W3CDTF and reads back tz-NAIVE (02-RESEARCH Pitfall 8). The
    # committed template pins 2026-01-01 (`_author_template._FIXED`), so this is falsifiable: a
    # writer that never touched the timestamps would read back that date, not the epoch.
    assert cp.created == EPOCH_ZERO.replace(tzinfo=None), cp.created
    assert cp.modified == EPOCH_ZERO.replace(tzinfo=None), cp.modified
    assert cp.identifier == surface.id, cp.identifier

    names = [shape.name for shape in written.slides[0].shapes]
    assert WATERMARK_NAME in names, (
        "the Draft deck rendered from the shipped template carries no watermark — an unreviewed "
        f"page that does not look unreviewed. Shapes: {names}"
    )
    assert names[-1] == WATERMARK_NAME, (
        "the watermark is not LAST in shape order, so it is not at the top of z-order and the "
        f"operator's own boxes can paint over it. Shapes: {names}"
    )

    assert surface.review.state is ReviewState.DRAFT, (
        "rendering moved the review gate. Nothing publishes without a human, and the renderer has "
        f"no write path to the gate at all. State is now: {surface.review.state!r}"
    )


def test_draft_and_published_renders_differ() -> None:
    """The cheap falsifiability check on the whole gate wiring.

    Same template, same content, two gate states — the decks must be genuinely different artifacts
    (the watermark, and `cp:contentStatus`). If a future refactor stopped reading
    `surface.is_published`, both renders would collapse into one and every Draft assertion in this
    module would stay green. This is the test that notices.

    `part_digest` rather than raw bytes on purpose: it is the implementation-independent comparison,
    so a red here means the CONTENT is identical, not that two zips happened to compress alike.
    """
    draft = render_surface_pptx_bytes(
        _sample_weekly_surface(),
        template=COMMITTED_TEMPLATE,
        slots=COMMITTED_SLOTS,
    )
    published = render_surface_pptx_bytes(
        _sample_weekly_surface(ReviewState.PUBLISHED),
        template=COMMITTED_TEMPLATE,
        slots=COMMITTED_SLOTS,
    )

    assert part_digest(draft) != part_digest(published), (
        "a Draft render and a Published render of the same content produced identical part "
        "content. The writer has stopped reading the review gate, so a Draft deck is now "
        "indistinguishable from an approved one — the watermark and cp:contentStatus say nothing"
    )
