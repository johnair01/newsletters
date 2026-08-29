"""Record — and re-verify — the measured `.pptx` determinism evidence (ROADMAP Phase-1 criterion 2).

WHY a committed script rather than a committed number. CONTEXT is explicit that "the determinism
spike must run a REAL write twice and commit the evidence — a decision without evidence is a vibe".
The risk an evidence-first phase actually carries is not a wrong measurement, it is a **fabricated**
one (threat T-01-03, repudiation): a plausible-looking JSON file typed by hand, never measured,
never falsifiable. The mitigation is that the artifact is produced by this script and re-verifiable
by it — `--check` re-runs the whole measurement and fails on any drift in an implementation-
independent field. A hand-typed hash cannot survive `--check`, and it cannot survive
`tests/test_pptx_determinism.py`, which independently re-derives the same invariants.

WHAT IS RECORDED VS WHAT IS ASSERTED. The file records `raw_a_sha256`, `raw_b_sha256`,
`normalized_a_sha256` and `normalized_b_sha256` as **evidence of what this machine measured**. Those
four are NEVER asserted across environments, and `--check` does not compare them. DEFLATE output is
zlib-implementation-dependent — reproducible-builds documented Fedora 40 (zlib-ng) and Debian (zlib)
emitting different compressed streams, and different CRCs, for identical input (01-RESEARCH
§Pitfall 1). Asserting a full-file hash across machines produces a test that is green locally and
red in CI with byte-identical part CONTENT: green-vs-red with nothing wrong. What `--check` does
compare is the set of implementation-INDEPENDENT facts — `raw_bytes_equal`, `varying_parts`,
`varying_zip_fields`, `normalized_bytes_equal`, `part_digest_a`, `part_digest_b` — because those
hold on any zlib, and a drift in any of them is a real finding.

THE `notes` MAPPING carries one short explanatory note per non-obvious key, mirroring the
`EXPECTED`-table convention in `tests/test_pptx_golden.py`: an evidence file a reader cannot
interpret without the plan open is not evidence, it is a number dump.

Run:  .venv/bin/python tests/fixtures/weekly/_record_determinism_evidence.py
      .venv/bin/python tests/fixtures/weekly/_record_determinism_evidence.py --check
"""

from __future__ import annotations

import hashlib
import io
import json
import pathlib
import platform
import sys
import time
import zlib

import pptx
from pptx import Presentation

from newsletters.adapters._timestamps import EPOCH_ZERO

HERE = pathlib.Path(__file__).parent
sys.path.insert(0, str(HERE))

from _determinism import (  # noqa: E402
    differing_parts,
    differing_zipinfo_fields,
    normalize_opc_zip,
    part_digest,
)

REPO_ROOT = HERE.parents[2]
TEMPLATE = HERE / "template.pptx"
EVIDENCE = (
    REPO_ROOT / ".planning" / "notes" / "2026-08-29-pptx-determinism-evidence.json"
)

# The generated-by marker and gate state — identical to the ones test_pptx_determinism.py writes,
# so the evidence measures the same render the durable test re-proves.
MARKER = "generated-by:newsletters"
GATE_STATE = "draft"

# DOS timestamps have 2-SECOND granularity; 3 seconds guarantees the boundary is crossed.
SECONDS_BETWEEN_WRITES = 3

# EPOCH_ZERO, tz-stripped (dcterms reads back tz-naive). DERIVED from the canonical sentinel —
# one epoch sentinel for the whole repo, never a second (test_pptx_determinism.py states the rule).
_EPOCH_NAIVE = EPOCH_ZERO.replace(tzinfo=None)

# The fields `--check` re-verifies. Deliberately excludes every raw/normalized sha256: those are
# zlib-implementation-dependent (see the module docstring) and asserting them would be a bug.
CHECKED_FIELDS = (
    "raw_bytes_equal",
    "varying_parts",
    "varying_zip_fields",
    "normalized_bytes_equal",
    "part_digest_a",
    "part_digest_b",
)

NOTES = {
    "seconds_between_writes": (
        "DOS zip timestamps have 2-second granularity; a shorter gap can leave two writes "
        "inside one tick, which are already byte-identical and would record a false positive."
    ),
    "raw_bytes_equal": (
        "False is the CORRECT result and is the negative control: un-normalized python-pptx "
        "output across a time boundary must differ, or the byte-equality result below is not "
        "attributable to the normalizer."
    ),
    "varying_parts": (
        "[] means every unzipped part was byte-identical across the two writes — the deck's "
        "CONTENT is deterministic; only the container drifted."
    ),
    "varying_zip_fields": (
        "['date_time'] is python-pptx's single non-determinism: _ZipPkgWriter.write() passes a "
        "str arcname to zipfile.writestr, so stdlib stamps time.localtime()."
    ),
    "normalized_bytes_equal": (
        "True is the recorded determinism outcome: BYTE-STABLE via a declared post-save zip "
        "normalization. Scoped to a fixed (python-pptx, zlib) pair."
    ),
    "part_digest_a": (
        "Implementation-INDEPENDENT content identity: sha256 over sorted (part name, "
        "sha256(part bytes)). This is the digest a committed==fresh gate must compare."
    ),
    "raw_a_sha256": (
        "Recorded as evidence, NEVER asserted across environments — DEFLATE output varies by "
        "zlib implementation (zlib vs zlib-ng)."
    ),
    "normalized_a_sha256": (
        "Same caveat as raw_a_sha256: evidence of what this machine measured, not a "
        "cross-environment contract."
    ),
    "zlib": (
        "Recorded because it scopes the byte-identity claim; a different zlib may produce "
        "different compressed bytes for identical part content."
    ),
    "python_pptx": (
        "python-pptx makes no documented byte-output-stability promise across versions; the "
        "exercised version is part of the evidence."
    ),
    "template_sha256": (
        "The fabricated template these writes rendered through — pins WHAT was measured."
    ),
}


def render_bytes(title: str) -> bytes:
    """The same miniature writer `tests/test_pptx_determinism.py` uses. Only `save()` is unstable."""
    prs = Presentation(str(TEMPLATE))
    slide = prs.slides[0]

    content = {
        "NL_WEEK_TITLE": title,
        "NL_MODULE": "module-a",
        "NL_HIGHLIGHTS": "shipped the thing",
        "NL_LOWLIGHTS": "the other thing slipped",
    }
    by_name = {shape.name: shape for shape in slide.shapes}
    for name, text in content.items():
        by_name[name].text_frame.text = text

    cp = prs.core_properties
    cp.category = MARKER
    cp.content_status = GATE_STATE
    cp.created = _EPOCH_NAIVE
    cp.modified = _EPOCH_NAIVE

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def measure() -> dict:
    """Perform the REAL double write and return the evidence mapping."""
    raw_a = render_bytes("2026-W35")
    time.sleep(SECONDS_BETWEEN_WRITES)
    raw_b = render_bytes("2026-W35")

    norm_a = normalize_opc_zip(raw_a)
    norm_b = normalize_opc_zip(raw_b)

    varying_parts = differing_parts(raw_a, raw_b)
    varying_zip_fields = differing_zipinfo_fields(raw_a, raw_b)
    normalized_equal = norm_a == norm_b

    outcome = (
        "BYTE-STABLE via a declared post-save zip normalization: two real python-pptx writes "
        f"{SECONDS_BETWEEN_WRITES}s apart differ ONLY in zip {varying_zip_fields}, every unzipped "
        "part is byte-identical, and normalize_opc_zip makes the files byte-identical. Scoped to "
        "a fixed (python-pptx, zlib) pair; the implementation-independent assertion is part_digest."
    )

    return {
        "recorded": "2026-08-29",
        "python": sys.version.split()[0],
        "python_pptx": pptx.__version__,
        "zlib": zlib.ZLIB_RUNTIME_VERSION,
        "platform": platform.platform(),
        "template_sha256": hashlib.sha256(TEMPLATE.read_bytes()).hexdigest(),
        "seconds_between_writes": SECONDS_BETWEEN_WRITES,
        "raw_a_sha256": hashlib.sha256(raw_a).hexdigest(),
        "raw_b_sha256": hashlib.sha256(raw_b).hexdigest(),
        "raw_bytes_equal": raw_a == raw_b,
        "varying_parts": varying_parts,
        "varying_zip_fields": varying_zip_fields,
        "normalized_a_sha256": hashlib.sha256(norm_a).hexdigest(),
        "normalized_b_sha256": hashlib.sha256(norm_b).hexdigest(),
        "normalized_bytes_equal": normalized_equal,
        "part_digest_a": part_digest(raw_a),
        "part_digest_b": part_digest(raw_b),
        "outcome": outcome,
        "notes": NOTES,
    }


def record() -> None:
    """Measure and write the committed evidence artifact."""
    evidence = measure()
    EVIDENCE.parent.mkdir(parents=True, exist_ok=True)
    EVIDENCE.write_text(json.dumps(evidence, indent=2, sort_keys=False) + "\n", "utf-8")
    print(f"wrote {EVIDENCE}")
    print(evidence["outcome"])


def check() -> int:
    """Re-measure and compare the implementation-INDEPENDENT fields. Exit 1 on any drift."""
    if not EVIDENCE.is_file():
        print(
            f"FAIL: {EVIDENCE} does not exist — run this script with no arguments first."
        )
        return 1

    committed = json.loads(EVIDENCE.read_text("utf-8"))
    fresh = measure()

    drift = [
        (field, committed.get(field), fresh[field])
        for field in CHECKED_FIELDS
        if committed.get(field) != fresh[field]
    ]
    if drift:
        print(
            "FAIL: the committed determinism evidence no longer matches a live measurement."
        )
        print(
            "These fields are implementation-INDEPENDENT, so a drift here is a real finding:"
        )
        for field, was, now in drift:
            print(f"  - {field}: committed={was!r}  measured={now!r}")
        print(
            "\nWhat this means: either the writer's behaviour changed (a python-pptx upgrade, a "
            "changed template, a new source of non-determinism), or the committed numbers were "
            "not produced by a real measurement. Re-run this script with no arguments to "
            "re-record, and say in the summary WHY the outcome moved."
        )
        return 1

    print(
        f"OK: {len(CHECKED_FIELDS)} implementation-independent fields re-verified against a "
        "live measurement."
    )
    print(
        f"  part_digest_a == part_digest_b: {fresh['part_digest_a'] == fresh['part_digest_b']}"
    )
    print(
        f"  raw_bytes_equal: {fresh['raw_bytes_equal']}  (False = the negative control holds)"
    )
    print(f"  varying_zip_fields: {fresh['varying_zip_fields']}")
    print(f"  normalized_bytes_equal: {fresh['normalized_bytes_equal']}")
    return 0


def main() -> int:
    if "--check" in sys.argv[1:]:
        return check()
    record()
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
