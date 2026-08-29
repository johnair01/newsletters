"""Author the fabricated weekly `.pptx` template the determinism spike renders through.

THE FIXTURE. One deck, `template.pptx`, one slide on the blank layout, five explicitly named
shapes:

| Shape name        | Why it is there |
|-------------------|-----------------|
| `NL_WEEK_TITLE`   | renderer slot — the week label |
| `NL_MODULE`       | renderer slot — the module this weekly is for |
| `NL_HIGHLIGHTS`   | renderer slot — authored highlights |
| `NL_LOWLIGHTS`    | renderer slot — authored lowlights (the absence a weekly is tempted to hide) |
| `Footer`          | DELIBERATELY UNPREFIXED — a decorative shape that is NOT a slot |

THE UNPREFIXED SHAPE IS LOAD-BEARING. The `NL_` reserved prefix (01-RESEARCH §Pattern 2) exists so
an operator's logo, footer and page number are not mistaken for unfilled slots. A prefix only earns
its keep if a non-slot shape exists to prove it. `Footer` is that shape: Phase 2's fail-loud contract
(the locked decision D-03 — named placeholders, fail loud on missing AND unknown names) needs a
fixture it can test in BOTH directions — "an `NL_` slot left empty must raise" and "a non-`NL_`
shape left alone must NOT raise" — and this deck is it.

THIRD-PARTY CORE PROPERTIES ARE SCRUBBED (01-RESEARCH §Pitfall 6). `Presentation()` with no argument
loads python-pptx's bundled default template, whose `docProps/core.xml` ships
``last_modified_by = "Steve Canny"``, ``comments = "generated using python-pptx"`` and
``created/modified = 2013-01-27``. Committing that into this repo would put a third party's name and
a foreign tool's marketing string into our fixture corpus — cheap to fix, embarrassing to ship. Every
core property is therefore overwritten with a fabricated, neutral value and the timestamps pinned to
one fixed instant. `test_pptx_determinism.py` asserts the absence against the written
`docProps/core.xml`, so a regeneration that forgot to scrub fails loudly.

BYTE-REPRODUCIBILITY. The saved bytes are routed through `newsletters.pptx_writer.normalize_opc_zip`
BEFORE the file is written, so the committed binary is itself already normalized: rerunning this script on
identical inputs reproduces the committed file byte-for-byte within one (python-pptx, zlib) pair, and
`normalize_opc_zip(committed) == committed` holds as an idempotence check anyone can run. No `now()`,
no `random` — the only clock python-pptx would consult is the zip entry timestamp, and the normalizer
removes it.

THE FABRICATED CONTENT IS PLACEHOLDER TEXT, NOT SAMPLE CONTENT. The slot bodies carry a short
"filled by the renderer" marker rather than realistic weekly prose: this is a TEMPLATE, and any
narrative that looked real would risk being mistaken for authored evidence.

Run:  .venv/bin/python tests/fixtures/weekly/_author_template.py
"""

from __future__ import annotations

import io
import pathlib
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt

from newsletters.pptx_writer import normalize_opc_zip

HERE = pathlib.Path(__file__).parent
TEMPLATE = HERE / "template.pptx"

# A fixed instant pinned into docProps/core.xml so the deck embeds no wall-clock. Arbitrary and
# stable — mirrors the `_FIXED` convention in tests/fixtures/pptx/_author_fixtures.py.
#
# DELIBERATELY NOT `EPOCH_ZERO`, and it must stay that way — this is a FALSIFIABILITY CONTROL
# (02-01-PLAN P-05, which REJECTS the IN-04 proposal to consolidate the two literals). The writer
# is specified to pin `dcterms:created`/`.modified` to EPOCH_ZERO; if the TEMPLATE already carried
# EPOCH_ZERO, the marker read-back assertion in plan 02-02 would pass on a deck whose timestamps
# the writer never touched. A template date the writer must overwrite is what makes that assertion
# able to fail.
_FIXED = datetime(2026, 1, 1, 0, 0, 0)

# Fabricated, neutral core properties. Every one of these overwrites a stock-template value; the
# strings name THIS repo's fixture, never a person and never a third-party tool.
_CORE_PROPERTIES = {
    "author": "newsletters fixture author",
    "last_modified_by": "newsletters fixture author",
    "title": "Weekly template (fabricated fixture)",
    "subject": "Determinism spike fixture",
    "comments": "Fabricated fixture template; regenerate with _author_template.py",
    "category": "",
    "content_status": "",
    "keywords": "",
    "identifier": "",
    "language": "",
    "version": "",
}

# (name, left, top, width, height, body) in inches. The four NL_ slots plus one unprefixed
# decorative shape, authored in a fixed order so shape order is deterministic.
_SHAPES = [
    ("NL_WEEK_TITLE", 0.5, 0.4, 9.0, 1.0, "NL_WEEK_TITLE"),
    ("NL_MODULE", 0.5, 1.4, 9.0, 0.6, "NL_MODULE"),
    ("NL_HIGHLIGHTS", 0.5, 2.2, 4.3, 3.4, "NL_HIGHLIGHTS"),
    ("NL_LOWLIGHTS", 5.2, 2.2, 4.3, 3.4, "NL_LOWLIGHTS"),
    # NOT a renderer slot: no NL_ prefix. Proves the reserved prefix discriminates.
    ("Footer", 0.5, 6.0, 9.0, 0.4, "fabricated fixture deck - not a renderer slot"),
]


def build_template_bytes() -> bytes:
    """Build the fabricated template deck and return its NORMALIZED bytes."""
    prs = Presentation()

    core = prs.core_properties
    for field, value in _CORE_PROPERTIES.items():
        setattr(core, field, value)
    core.created = _FIXED
    core.modified = _FIXED
    core.revision = 1

    # slide_layouts[6] is the BLANK layout in python-pptx's default template
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    for name, left, top, width, height, body in _SHAPES:
        box = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        box.name = name
        frame = box.text_frame
        frame.text = body
        frame.paragraphs[0].runs[0].font.size = Pt(18)

    buf = io.BytesIO()
    prs.save(buf)
    # Normalize BEFORE the file is written: the committed artifact is already normalized, so
    # `normalize_opc_zip(committed) == committed` and a re-run is byte-reproducible.
    return normalize_opc_zip(buf.getvalue())


def main() -> None:
    raw = build_template_bytes()
    TEMPLATE.write_bytes(raw)
    print(f"wrote {TEMPLATE} ({len(raw)} bytes)")


if __name__ == "__main__":
    main()
