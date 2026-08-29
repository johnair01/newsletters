"""Weekly Spec authoring path — proof suite (ported from ``tests/test_casespec.py``).

Three committed fixtures under ``tests/fixtures/weekly/`` drive the LIVE validator/loader/composer
(``src/newsletters/weeklyspec.py``) end to end. Assertions are on STRUCTURE and INVARIANTS; where
a concrete value is needed it is read from the SAME parsed fixture the loader reads, so the test
tracks the corpus rather than freezing a magic string.

* ``weekly-full.yml`` — all eight keys: an ordinary quoted highlight AND a block-scalar one, a
  lowlight, two recognitions (one sourced, one not), two team members, three assets (complete /
  deep-link-required / provenance-incomplete), and a ``config:`` subtree whose values must NEVER
  become claims. One person string appears in BOTH ``recognitions:`` and ``team:`` — the
  span-swap regression's non-vacuity.
* ``weekly-sparse.yml`` — only ``week`` + ``module``: every other key must be disclosed in
  ``Distillation.missing[]``, never fabricated.
* ``weekly-editorial-bait.yml`` — six highlights arranged as a summarizable, an out-of-order and
  a mergeable pair. It exists to BAIT a composer into "improving" the author's lines; the
  composer must carry all six separately, in file order, byte-identical.

The invariants proven here (the Case Spec's eight, plus two the weekly earns):

1. **Schema validation** — teaching errors on a non-mapping doc, an unknown TOP-LEVEL key, an
   unknown key inside a recognition / a team member / an asset, and a type-coerced scalar.
2. **Trace faithfulness** — every claim content-addressed via ``Trace.from_source``, non-stale,
   re-sliceable from the live transcript, and entailed by the LIVE
   ``SpanContainmentFaithfulness`` gate (the strict half — never the structural fallback).
3. **The block-scalar item region** — asserted as the gate's own rule
   (``_normalize(text) in _normalize(span)``), never ``transcript[start:end] == text``.
4. **The span-swap regression** — a value duplicated across two sections traces to ITS OWN LINE.
5. **missing[] honesty** — every absent key is disclosed, in schema order, including "no
   lowlights"; no claim exists for an absent field.
6. **Narrative verbatim** — authored lines reach the typed spec byte-identical to the parsed YAML.
7. **Config never claimed** — no ``config:`` leaf appears in any claim; the subtree is carried.
8. **Recognition evidence** — absent / unresolvable ``source:`` ⇒ ``evidence == []`` + a named
   disclosure; a resolved one ⇒ exactly one SPAN-LESS ``Trace``. Never a fabricated span.
9. **Root containment** — a path outside ``root`` RAISES (a refusal, not a ``missing[]`` entry).
10. **Determinism + read-only** — two loads are byte-identical, the JSON round-trips, and the
    fixture's mtime and size are unchanged after a load.
11. **Asset routing** — every row of ``docs/weekly-spec.md`` §"The routing" proved BOTH ways:
    each refusal carries the spec's exact disclosure string while a well-formed asset in the same
    document still places; a path escaping the root (directly or through a symlink) RAISES and
    contributes nothing to ``missing[]``; the image is hashed and never decoded.
12. **The composer** — a Draft ``Surface(REPORT)`` at ``EPOCH_ZERO`` in a fixed, asserted block
    order, byte-identical across two composes, that never advances the gate; plus the
    editorialization guard: every block string is authored or a declared connective constant,
    with a planted-paraphrase non-vacuity arm.
13. **The deck (SC-5)** — ``weekly_slots`` emits exactly the four ``NL_`` keys the committed
    template declares, in a fixed order, every value a ``list[str]`` whose every line is either
    authored or a member of ``surface.missing``; an empty section's single line IS its own
    disclosure (with a non-vacuity arm proving the membership self-check refuses anything else);
    and the composed weekly — full AND sparse — renders through Phase 2's writer to a
    deterministic, marked, Draft-watermarked deck, every property asserted by reopening the
    WRITTEN bytes, with the Surface unchanged and the committed template untouched.
"""

from __future__ import annotations

import hashlib
import pathlib

import pytest

import newsletters.weeklyspec
from newsletters._yaml_loader import load_config
from newsletters.adapters._timestamps import EPOCH_ZERO
from newsletters.compose import addressed, compute_delta
from newsletters.distill.faithfulness import SpanContainmentFaithfulness, _normalize
from newsletters.pptx_writer import (
    DRAFT_STATUS,
    MARKER,
    WATERMARK_NAME,
    part_digest,
    render_surface_pptx,
    render_surface_pptx_bytes,
)
from newsletters.semantic import Claim, KpiItem, ReviewState, Source, Trace
from newsletters.swimlane import SectionBinding
from newsletters.templates import REPORT
from newsletters.weeklyspec import (
    CONNECTIVE_CONSTANTS,
    HIGHLIGHTS_SLOT,
    LOWLIGHTS_SLOT,
    MODULE_SLOT,
    WEEK_TITLE_SLOT,
    AuthoredAsset,
    AuthoredMember,
    AuthoredRecognition,
    WeeklySpec,
    WeeklySpecLoad,
    _validate,
    build_weekly_report,
    load_weekly_spec,
    weekly_slots,
)

REPO_ROOT = pathlib.Path(__file__).resolve().parent.parent
FIXTURE_DIR = REPO_ROOT / "tests" / "fixtures" / "weekly"
FULL = "weekly-full.yml"
SPARSE = "weekly-sparse.yml"
BAIT = "weekly-editorial-bait.yml"
FIXTURES = [FULL, SPARSE, BAIT]
AUTHOR = "author-x"

# The eight keys, in schema order. Duplicated here ON PURPOSE: importing the module's own
# ``_KNOWN_KEYS`` would make this assertion vacuous (it would compare the constant to itself).
SCHEMA_KEYS = [
    "week",
    "module",
    "highlights",
    "lowlights",
    "recognitions",
    "team",
    "assets",
    "config",
]


def _load(name: str, **kwargs) -> WeeklySpecLoad:
    return load_weekly_spec(FIXTURE_DIR / name, root=REPO_ROOT, **kwargs)


def _parsed(name: str) -> dict:
    return load_config((FIXTURE_DIR / name).read_text(encoding="utf-8"))


def _line_of(transcript: str, offset: int) -> int:
    """1-based line number of a character offset — the unit the span-swap proof asserts in."""
    return transcript[:offset].count("\n") + 1


def _line_containing(transcript: str, needle: str) -> int:
    """1-based number of the FIRST line containing ``needle`` (the fixture authors it once)."""
    for index, line in enumerate(transcript.splitlines(), start=1):
        if needle in line:
            return index
    raise AssertionError(f"{needle!r} is not on any line of the fixture")


def _validated(text: str) -> dict:
    """Parse + validate one authored document, the way the loader will."""
    return _validate(load_config(text))


# --------------------------------------------------------------------------- #
# 1 — schema validation: teaching errors at BOTH levels, never a silent drop
# --------------------------------------------------------------------------- #


def test_non_mapping_document_names_the_eight_key_schema() -> None:
    """A document that is not a mapping fails loudly, naming the schema and the spec doc."""
    with pytest.raises(ValueError) as excinfo:
        _validated("- just\n- a list\n")
    message = str(excinfo.value)
    assert "YAML mapping" in message
    assert "docs/weekly-spec.md" in message
    for key in SCHEMA_KEYS:
        assert repr(key) in message, f"{key} not named in the schema error: {message}"


def test_unknown_top_level_key_refuses_to_drop_authored_content() -> None:
    """``highlight:`` for ``highlights:`` fails loudly — a typo must not lose a lowlight."""
    with pytest.raises(ValueError) as excinfo:
        _validated('week: "W1"\nhighlight:\n  - "one"\n')
    message = str(excinfo.value)
    assert "unknown Weekly Spec field" in message
    assert "'highlight'" in message
    assert "Refusing to drop authored content silently." in message
    for key in SCHEMA_KEYS:
        assert repr(key) in message


@pytest.mark.parametrize(
    ("document", "container", "offender"),
    [
        (
            'recognitions:\n  - person: "P"\n    resaon: "typo"\n',
            "recognitions[0]",
            "resaon",
        ),
        (
            'team:\n  - name: "N"\n    rol: "typo"\n',
            "team[0]",
            "rol",
        ),
        (
            'assets:\n  k:\n    file: "f"\n    fodler: "typo"\n',
            "assets['k']",
            "fodler",
        ),
    ],
)
def test_unknown_nested_key_fails_loud_naming_container_and_key(
    document: str, container: str, offender: str
) -> None:
    """A mistyped field INSIDE a container earns the same refusal as a top-level typo.

    Rule 1 names the top level, but a dropped ``resaon:`` loses authored content exactly as
    silently as a dropped ``lowlight:`` would. Same teaching voice, same refusal.
    """
    with pytest.raises(ValueError) as excinfo:
        _validated(document)
    message = str(excinfo.value)
    assert repr(offender) in message
    assert container in message
    assert "Refusing to drop authored content silently." in message


@pytest.mark.parametrize(
    ("document", "field", "actual"),
    [
        ("week: 42\n", "week", "int"),
        ("module: yes\n", "module", "bool"),
        ('highlights:\n  - "ok"\n  - 42\n', "highlights[1]", "int"),
        ('recognitions:\n  - person: yes\n', "recognitions[0].person", "bool"),
        ('team:\n  - name: "N"\n    lines:\n      - 42\n', "team[0].lines[0]", "int"),
        ('assets:\n  k:\n    sha256: 42\n', "assets['k'].sha256", "int"),
    ],
)
def test_type_coerced_scalar_names_the_field_its_type_and_the_fix(
    document: str, field: str, actual: str
) -> None:
    """A bare ``yes`` / ``42`` stops being the author's text, so the loader refuses it."""
    with pytest.raises(ValueError) as excinfo:
        _validated(document)
    message = str(excinfo.value)
    assert repr(field) in message
    assert f"got {actual}" in message
    assert "quote the value so YAML cannot type-coerce it" in message


@pytest.mark.parametrize("name", FIXTURES)
def test_both_committed_fixtures_parse_and_validate(name: str) -> None:
    """The corpus every later Phase-3 proof is written against is itself schema-clean."""
    parsed = _validate(_parsed(name))
    assert isinstance(parsed, dict) and parsed
    assert all(key in SCHEMA_KEYS for key in parsed), sorted(parsed)
    # File order IS schema order in the corpus — the forward cursor depends on the walk being
    # able to follow the document; a fixture authored out of schema order would hide that.
    assert list(parsed) == [k for k in SCHEMA_KEYS if k in parsed]


def test_full_fixture_carries_the_shapes_the_phase_is_proved_against() -> None:
    """Non-vacuity: the full fixture really does carry every shape the later tests need.

    Without this, a fixture that quietly lost its block scalar or its duplicated person would
    leave the block-scalar and span-swap proofs passing for the wrong reason.
    """
    parsed = _parsed(FULL)
    assert list(parsed) == SCHEMA_KEYS, "the full fixture must author all eight keys"

    # Two highlights: one an ordinary scalar, one a BLOCK scalar (multi-line, so its folded
    # value is not a verbatim substring of the file).
    highlights = parsed["highlights"]
    assert len(highlights) == 2
    assert "\n" not in highlights[0], "the first highlight must be an ordinary scalar"
    assert "\n" in highlights[1], "the second highlight must be a block scalar"

    assert len(parsed["lowlights"]) == 1

    # Two recognitions — exactly one of them sourced (rule 6's both-halves case).
    sourced = [r for r in parsed["recognitions"] if r.get("source")]
    assert len(parsed["recognitions"]) == 2 and len(sourced) == 1

    # The duplicated person string: named in BOTH recognitions and team.
    people = {r["person"] for r in parsed["recognitions"]}
    members = {m["name"] for m in parsed["team"]}
    assert people & members, "a person must appear in BOTH sections for the span-swap proof"

    # Three assets: complete / deep-link-required / provenance-incomplete.
    assets = parsed["assets"]
    assert len(assets) == 3
    complete = [
        k for k, r in assets.items() if all(r.get(f) for f in ("folder", "date", "event"))
    ]
    assert len(complete) == 2
    values_no_link = [
        k
        for k, r in assets.items()
        if r.get("stands_in_for") == "values" and not r.get("link")
    ]
    assert len(values_no_link) == 1, "the deep-link row must be exercised"
    incomplete = [
        k
        for k, r in assets.items()
        if not all(r.get(f) for f in ("folder", "date", "event"))
    ]
    assert len(incomplete) == 1, "one asset must be missing a provenance minimum"

    assert parsed["config"], "config must be non-empty or the never-claimed guard is vacuous"


def test_sparse_fixture_authors_only_the_two_scalars() -> None:
    """The sparse fixture's whole job is ABSENCE — six keys unwritten, including lowlights."""
    parsed = _parsed(SPARSE)
    assert list(parsed) == ["week", "module"]
    for key in SCHEMA_KEYS[2:]:
        assert key not in parsed


# --------------------------------------------------------------------------- #
# 2 — the authored-side models: permissive by design, so absences are disclosable
# --------------------------------------------------------------------------- #


def test_authored_models_are_permissive_so_an_incomplete_entry_is_disclosable() -> None:
    """An incomplete authored entry must be REPRESENTABLE long enough to be disclosed.

    This is the load half of the load/place seam. A recognition with no source is still credit
    owed; an asset missing a provenance minimum must be carried far enough to NAME the missing
    field. The strict shape (``semantic.AssetRecord``, whose minimums are non-optional) is
    deliberately not used here — that one guards PLACEMENT, where an incomplete record must be
    unrepresentable. Both directions are asserted so neither can be quietly relaxed.
    """
    # Permissive on the authored side: nothing is required, everything defaults empty.
    assert AuthoredRecognition().model_dump() == {
        "person": "",
        "reason": "",
        "source": "",
        "evidence": [],
    }
    assert AuthoredMember().lines == []
    assert AuthoredAsset(key="k").folder == ""

    # ...and the strict shape refuses exactly what the permissive one carries.
    from pydantic import ValidationError

    from newsletters.semantic import AssetRecord

    with pytest.raises(ValidationError) as excinfo:
        AssetRecord(key="k", file="f", sha256="s")  # no folder / date / event
    missing_fields = {
        error["loc"][0] for error in excinfo.value.errors() if error["type"] == "missing"
    }
    assert {"folder", "date", "event"} <= missing_fields, missing_fields


def test_weekly_spec_defaults_are_empty_never_fabricated() -> None:
    """A ``WeeklySpec`` built from nothing carries nothing — no placeholder, no invented value."""
    spec = WeeklySpec()
    assert spec.week == "" and spec.module == ""
    assert spec.highlights == [] and spec.lowlights == []
    assert spec.recognitions == [] and spec.team == [] and spec.assets == []
    assert spec.config == {}


# --------------------------------------------------------------------------- #
# 3 — trace faithfulness: content-addressed spans, entailed by the LIVE gate
# --------------------------------------------------------------------------- #


@pytest.mark.parametrize("name", FIXTURES)
def test_trace_faithfulness(name: str) -> None:
    """Every claim is content-addressed, non-stale, re-sliceable, and STRICTLY entailed.

    Because every trace ``is_addressed``, ``SpanContainmentFaithfulness`` takes its strict branch
    (normalized claim text contained in the normalized span) — the structural un-addressed
    fallback is provably not what passes here.
    """
    gate = SpanContainmentFaithfulness()
    load = _load(name)
    source = load.source
    assert load.distillation.claims, f"{name}: expected at least one claim"
    for claim in load.distillation.claims:
        assert claim.is_traced, f"{name}: {claim.text!r} is untraced"
        assert gate.entails(claim), f"{name}: {claim.text!r} fails the live gate"
        for trace in claim.evidence:
            assert trace.is_addressed, f"{name}: {claim.text!r} not content-addressed"
            assert trace.start is not None and trace.end is not None
            # the recorded window re-slices the LIVE transcript to the stored span
            assert source.transcript[trace.start : trace.end] == trace.span
            # the span REALLY contains the claim (the strict rule, asserted directly)
            assert _normalize(claim.text) in _normalize(trace.span)
            assert trace.content_hash == source.content_hash()
            assert not trace.is_stale_against(source)


def test_block_scalar_highlight_is_traced_to_its_own_item_region() -> None:
    """The block-scalar highlight's span is the raw ITEM REGION — asserted as the gate's rule.

    A block scalar's FOLDED value is not a substring of the file at all (the raw text carries the
    indentation the folding removed), so "a real span" here means "the region that CONTAINS your
    text", not "the exact bytes of your text". The honest assertion is therefore the gate's own —
    ``_normalize(text) in _normalize(span)`` — plus ``transcript[start:end] == span``. Asserting
    ``transcript[start:end] == text`` would be false for a correct implementation.

    The regression this pins (inherited from the Case Spec): the block item's region must not
    swallow its sibling, so the ordinary highlight is still located verbatim and minted, never
    falsely disclosed as unlocatable.
    """
    load = _load(FULL)
    transcript = load.source.transcript
    plain, block = _parsed(FULL)["highlights"]
    claims = [c for c in load.distillation.claims if c.topics == ["highlights"]]
    texts = [c.text for c in claims]

    assert texts == [plain, block], "both highlights must be traced claims, in file order"
    assert not any("could not be located" in note for note in load.distillation.missing)

    plain_trace = claims[0].evidence[0]
    assert transcript[plain_trace.start : plain_trace.end] == plain, "exact-find span"

    block_trace = claims[1].evidence[0]
    span = transcript[block_trace.start : block_trace.end]
    assert span == block_trace.span
    assert _normalize(block) in _normalize(span)
    assert span != block, "a block scalar's folded value is NOT its raw region"
    assert plain not in span, "the item region must not swallow its sibling"


# --------------------------------------------------------------------------- #
# 4 — the span-swap regression: a duplicated value traces to ITS OWN line
# --------------------------------------------------------------------------- #


def test_duplicated_person_traces_to_its_own_line_not_its_twins() -> None:
    """The person named in BOTH ``recognitions:`` and ``team:`` gets two claims on two LINES.

    03-RESEARCH proved by execution that minting a later field before an earlier one silently
    SWAPS the spans of a duplicated value — and that BOTH claims still pass the faithfulness
    gate, because the text is identical. No gate can catch it. So this test asserts on LINE
    NUMBERS, which is the only thing that differs, and it was observed RED once against a
    deliberately reversed walk order (recorded in 03-02-SUMMARY.md).
    """
    load = _load(FULL)
    transcript = load.source.transcript
    parsed = _parsed(FULL)

    duplicated = {r["person"] for r in parsed["recognitions"]} & {
        m["name"] for m in parsed["team"]
    }
    assert len(duplicated) == 1, "the fixture must name exactly one person in both sections"
    person = duplicated.pop()

    recognition_claims = [
        c
        for c in load.distillation.claims
        if c.topics == ["recognitions.person"] and c.text == person
    ]
    team_claims = [
        c
        for c in load.distillation.claims
        if c.topics == ["team.name"] and c.text == person
    ]
    assert len(recognition_claims) == 1 and len(team_claims) == 1

    recognition_line = _line_of(transcript, recognition_claims[0].evidence[0].start)
    team_line = _line_of(transcript, team_claims[0].evidence[0].start)

    assert recognition_line == _line_containing(transcript, f'person: "{person}"')
    assert team_line == _line_containing(transcript, f'name: "{person}"')
    assert recognition_line < team_line, "the spans are swapped — file order was not honored"


def test_claim_spans_ascend_strictly_in_document_order() -> None:
    """The forward cursor's signature: every claim's span starts after the previous one ends.

    The whole-document version of the assertion above — it catches an out-of-order walk anywhere,
    not only where the corpus happens to duplicate a value.
    """
    load = _load(FULL)
    previous_end = -1
    for claim in load.distillation.claims:
        trace = claim.evidence[0]
        assert trace.start >= previous_end, (
            f"claim {claim.text[:40]!r} starts at {trace.start}, before the previous claim "
            f"ended at {previous_end} — the walk left file order"
        )
        previous_end = trace.end


def test_value_in_a_full_line_comment_never_captures_the_span(
    tmp_path: pathlib.Path,
) -> None:
    """WR-01 (03-review): the forward exact-find must skip YAML comments.

    The review proved by execution that a value appearing verbatim in a comment BETWEEN the
    cursor and its own field pinned the claim's span to the COMMENT — a mis-attribution the
    faithfulness gate cannot catch, because the text is identical. The weekly is the first spec
    kind where authors realistically write comments beside people's names, so the regression is
    pinned here, through the weekly loader, against the review's own bait document.
    """
    text = (
        "highlights:\n"
        '  - "A line."\n'
        "# reviewer note: Devi R. did great work this week\n"
        "team:\n"
        '  - name: "Devi R."\n'
    )
    path = tmp_path / "weekly-comment-bait.yml"
    path.write_text(text, encoding="utf-8")
    load = load_weekly_spec(path, root=tmp_path)
    transcript = load.source.transcript

    claims = [c for c in load.distillation.claims if c.topics == ["team.name"]]
    assert len(claims) == 1 and claims[0].text == "Devi R."
    trace = claims[0].evidence[0]
    assert transcript[trace.start : trace.end] == "Devi R."
    assert _line_of(transcript, trace.start) == 5, (
        "the span landed on the comment (line 3), not the authored field (line 5) — "
        "a comment vouching for a team member"
    )


def test_value_in_an_inline_comment_never_captures_the_span(
    tmp_path: pathlib.Path,
) -> None:
    """The inline variant: a ``#`` comment after a value on an earlier line is still a comment."""
    text = (
        "highlights:\n"
        '  - "kept" # note: Kira Nerys reviewed this\n'
        "team:\n"
        '  - name: "Kira Nerys"\n'
    )
    path = tmp_path / "weekly-inline-comment-bait.yml"
    path.write_text(text, encoding="utf-8")
    load = load_weekly_spec(path, root=tmp_path)
    transcript = load.source.transcript

    claims = [c for c in load.distillation.claims if c.topics == ["team.name"]]
    assert len(claims) == 1 and claims[0].text == "Kira Nerys"
    trace = claims[0].evidence[0]
    assert transcript[trace.start : trace.end] == "Kira Nerys"
    assert _line_of(transcript, trace.start) == 4, (
        "the span landed on the inline comment (line 2), not the authored field (line 4)"
    )


# --------------------------------------------------------------------------- #
# 5 — missing[] honesty: every absence disclosed, in schema order
# --------------------------------------------------------------------------- #


def test_sparse_spec_discloses_every_absent_key_in_schema_order() -> None:
    """Six unwritten keys, disclosed by name and in schema order — including "no lowlights".

    Rule 4 singles out lowlights as "precisely the absence a weekly is tempted to hide", so the
    assertion names it rather than trusting a loop to have covered it.
    """
    load = _load(SPARSE)
    parsed = _parsed(SPARSE)
    absent_keys = [k for k in SCHEMA_KEYS if k not in parsed]
    assert absent_keys == SCHEMA_KEYS[2:], "the sparse fixture must omit six keys"

    disclosed_order = [
        key
        for key in SCHEMA_KEYS
        if any(f"{key!r} is absent" in note for note in load.distillation.missing)
    ]
    # `config` absence is NOT a gap — a weekly with no org slots is simply fully portable.
    assert disclosed_order == [k for k in absent_keys if k != "config"]
    assert any(
        "'lowlights' is absent" in note for note in load.distillation.missing
    ), load.distillation.missing

    # No fabrication: every claim's topic is a key the author actually wrote.
    authored = set(parsed)
    for claim in load.distillation.claims:
        assert claim.topics and claim.topics[0].split(".")[0] in authored


def test_full_spec_discloses_its_per_item_absences() -> None:
    """Absences INSIDE a carried entry are disclosed too — the entry is kept, the gap is named."""
    load = _load(FULL)
    notes = load.distillation.missing
    # The second team member authors no photo; the member and their lines are still carried.
    assert any("'team[1].photo' is absent" in note for note in notes), notes
    assert load.spec.team[1].name and load.spec.team[1].lines
    # ...and nothing was invented to fill the gap.
    assert load.spec.team[1].photo == ""


def test_blank_list_items_are_dropped_with_a_positional_disclosure(
    tmp_path: pathlib.Path,
) -> None:
    """WR-05 (03-review): a PRESENT-but-blank item inside a non-empty list is disclosed.

    The banner promises "anything absent, empty or unlocatable lands in missing[]" (rule 4),
    but a blank item inside a non-empty list (``["kept", ""]``, a ``~`` null entry, a blank
    team line) was filtered out with NO disclosure — proven by execution in the review. The
    author put a line there; its emptiness is now disclosed by position, never silently eaten.
    """
    text = (
        "highlights:\n"
        '  - "kept"\n'
        '  - ""\n'
        "  - ~\n"
        "team:\n"
        '  - name: "Kira Nerys"\n'
        "    lines:\n"
        '      - "carried"\n'
        '      - ""\n'
    )
    path = tmp_path / "weekly-blank-items.yml"
    path.write_text(text, encoding="utf-8")
    load = load_weekly_spec(path, root=tmp_path)
    notes = load.distillation.missing

    # Nothing real is lost, and nothing blank is fabricated into the spec...
    assert load.spec.highlights == ["kept"]
    assert load.spec.team[0].lines == ["carried"]
    # ...and every dropped blank is disclosed by its authored position.
    assert any("'highlights[1]' is absent" in note for note in notes), notes
    assert any("'highlights[2]' is absent" in note for note in notes), notes
    assert any("'team[0].lines[1]' is absent" in note for note in notes), notes
    # The kept lines are NOT disclosed — disclosure tracks the blanks, not the list.
    assert not any("'highlights[0]'" in note for note in notes), notes
    assert not any("'team[0].lines[0]'" in note for note in notes), notes


# --------------------------------------------------------------------------- #
# 6 — narrative reaches the typed spec byte-verbatim
# --------------------------------------------------------------------------- #


def test_authored_narrative_is_carried_byte_verbatim() -> None:
    """Every authored line equals the parsed YAML value exactly — never summarised or merged."""
    load = _load(FULL)
    parsed = _parsed(FULL)

    assert load.spec.week == parsed["week"]
    assert load.spec.module == parsed["module"]
    assert load.spec.highlights == parsed["highlights"]
    assert load.spec.lowlights == parsed["lowlights"]

    for authored, carried in zip(parsed["recognitions"], load.spec.recognitions):
        assert carried.person == authored["person"]
        assert carried.reason == authored["reason"]
    for authored, carried in zip(parsed["team"], load.spec.team):
        assert carried.name == authored["name"]
        assert carried.role == authored.get("role", "")
        assert carried.lines == authored.get("lines", [])

    # The assets keep their authored key AND their file order.
    assert [a.key for a in load.spec.assets] == list(parsed["assets"])
    for asset in load.spec.assets:
        authored_record = parsed["assets"][asset.key]
        for field, value in authored_record.items():
            assert getattr(asset, field) == value

    # Every carried line is also a traced claim — the record and its evidence cannot drift.
    claim_texts = {c.text for c in load.distillation.claims}
    for line in parsed["highlights"] + parsed["lowlights"]:
        assert line in claim_texts


# --------------------------------------------------------------------------- #
# 7 — config values are org slots: carried, NEVER claimed
# --------------------------------------------------------------------------- #


def _config_leaves(node: object) -> list[str]:
    if isinstance(node, dict):
        return [leaf for v in node.values() for leaf in _config_leaves(v)]
    if isinstance(node, list):
        return [leaf for v in node for leaf in _config_leaves(v)]
    return [] if node is None else [str(node)]


def test_config_never_in_claims_but_always_carried() -> None:
    """No ``config:`` leaf appears in any claim text; the subtree is carried exactly (non-vacuous)."""
    load = _load(FULL)
    leaves = _config_leaves(_parsed(FULL)["config"])
    assert leaves, "fixture must declare config values for this guard to be non-vacuous"

    for leaf in leaves:
        for claim in load.distillation.claims:
            assert leaf not in claim.text, f"config value {leaf!r} leaked into {claim.text!r}"
        for note in load.distillation.missing:
            assert leaf not in note, f"config value {leaf!r} leaked into a disclosure"

    # ...but the slots are carried (not lost) — they stay CONFIG, available for binding.
    assert load.spec.config == _parsed(FULL)["config"]


# --------------------------------------------------------------------------- #
# 8 — recognition evidence: a pointer, a disclosure, or both halves. Never a lie.
# --------------------------------------------------------------------------- #


def test_recognition_without_source_is_carried_and_disclosed() -> None:
    """Rule 6, both halves: the credit survives AND its missing evidence is named."""
    load = _load(FULL)
    parsed = _parsed(FULL)
    index, authored = next(
        (i, r) for i, r in enumerate(parsed["recognitions"]) if not r.get("source")
    )
    carried = load.spec.recognitions[index]

    assert carried.person == authored["person"] and carried.reason == authored["reason"]
    assert carried.evidence == [], "no source ⇒ no evidence, never a fabricated Trace"
    assert any(
        f"'recognitions[{index}].source' is absent" in note
        for note in load.distillation.missing
    ), load.distillation.missing
    # The author's own word IS separately evidenced — as claims at real spans of this file.
    claim_texts = {c.text for c in load.distillation.claims}
    assert authored["person"] in claim_texts and authored["reason"] in claim_texts


def test_resolved_source_becomes_exactly_one_span_less_trace() -> None:
    """A ``source:`` that names a KNOWN Source is a POINTER — one Trace, deliberately span-less.

    The loader has not read that source's text, so minting a span for it would be fabrication.
    """
    parsed = _parsed(FULL)
    source_id = next(r["source"] for r in parsed["recognitions"] if r.get("source"))
    known = Source(id=source_id, context="test:known", transcript="unread by the loader")

    load = _load(FULL, known_sources=[known])
    index = next(i for i, r in enumerate(parsed["recognitions"]) if r.get("source"))
    evidence = load.spec.recognitions[index].evidence

    assert len(evidence) == 1
    trace = evidence[0]
    assert trace.source_id == source_id
    assert trace.span == "" and trace.start is None and trace.end is None
    assert not trace.is_addressed, "a pointer is not a content-addressed span"
    assert not any(
        "does not resolve" in note for note in load.distillation.missing
    ), "a resolved id must not also be disclosed as unresolvable"


def test_unresolvable_source_is_disclosed_by_name_never_traced() -> None:
    """An id that names nothing gets the ABSENT-source treatment plus the id in the disclosure."""
    load = _load(FULL)  # no known_sources — the authored id resolves to nothing
    parsed = _parsed(FULL)
    index, authored = next(
        (i, r) for i, r in enumerate(parsed["recognitions"]) if r.get("source")
    )

    assert load.spec.recognitions[index].evidence == []
    note = next(n for n in load.distillation.missing if "does not resolve" in n)
    assert repr(authored["source"]) in note, note
    assert repr(authored["person"]) in note, note
    # The recognition itself is fully carried — dropping it would erase credit.
    assert load.spec.recognitions[index].person == authored["person"]
    assert load.spec.recognitions[index].reason == authored["reason"]


# --------------------------------------------------------------------------- #
# 9 — root containment is a REFUSAL, not an absence
# --------------------------------------------------------------------------- #


def test_spec_outside_root_raises_and_never_lands_in_missing(
    tmp_path: pathlib.Path,
) -> None:
    """A path escaping ``root`` raises. ``missing[]`` is for absent content, not refused requests.

    Rule 7: "This is not a ``missing[]`` case; it is a refusal." Proven by pairing the raise with
    a CONSTRUCTING arm — the same file loads cleanly once ``root`` legitimately contains it, so
    the refusal is the containment check and not a broken loader.
    """
    outside = tmp_path / "outside.yml"
    outside.write_text('week: "W1"\nmodule: "M"\n', encoding="utf-8")
    inner = tmp_path / "inner"
    inner.mkdir()

    with pytest.raises(ValueError):
        load_weekly_spec(outside, root=inner)

    load = load_weekly_spec(outside, root=tmp_path)
    assert load.spec.week == "W1"
    assert not any("outside" in note for note in load.distillation.missing)


# --------------------------------------------------------------------------- #
# 10 — determinism, lossless round-trip, and a loader that writes nothing
# --------------------------------------------------------------------------- #


@pytest.mark.parametrize("name", FIXTURES)
def test_lossless_roundtrip_and_determinism(name: str) -> None:
    """Raw file carried verbatim; two loads and the JSON round-trip are byte-identical."""
    raw = (FIXTURE_DIR / name).read_text(encoding="utf-8")
    first = _load(name)
    assert first.source.transcript == raw
    assert first.source.context == f"weekly-spec:tests/fixtures/weekly/{name}"

    second = _load(name)
    assert first.model_dump_json() == second.model_dump_json()

    dumped = first.model_dump_json()
    assert WeeklySpecLoad.model_validate_json(dumped).model_dump_json() == dumped


@pytest.mark.parametrize("name", FIXTURES)
def test_the_loader_writes_nothing(name: str) -> None:
    """READ-ONLY, asserted on the filesystem: the fixture's mtime and size survive a load."""
    path = FIXTURE_DIR / name
    before = path.stat()
    _load(name)
    after = path.stat()
    assert (after.st_mtime_ns, after.st_size) == (before.st_mtime_ns, before.st_size)


# --------------------------------------------------------------------------- #
# 11 — asset routing: every row of docs/weekly-spec.md §"The routing", both ways
#
# The four expected strings below are copied BYTE-FOR-BYTE from the spec's routing table and
# are deliberately NOT imported from ``weeklyspec`` — importing the module's own constants
# would compare each string to itself and assert nothing about the contract.
# --------------------------------------------------------------------------- #

_PROVENANCE_DISCLOSURE = (
    "asset {key!r}: provenance field {field!r} is absent — the minimum is folder + date "
    "+ event label; disclosed, never placed"
)
_DEEP_LINK_DISCLOSURE = (
    "asset {key!r}: a screenshot standing in for values requires a deep link to the "
    "report; disclosed, never placed"
)
_CONTENT_ADDRESS_DISCLOSURE = (
    "asset {key!r}: file {file!r} does not match its recorded content address — refusing "
    "to place a file that is not the one the record describes"
)
_PHOTO_DISCLOSURE = (
    "team member {name!r}: photo key {key!r} names no placed asset — the member is "
    "carried, the photo is not"
)

# The committed 1×1 PNG's real bytes, and a DIFFERENT file built from them. The pair drives the
# substitution row: a record describing A while B sits on disk.
_IMAGE_A = (FIXTURE_DIR / "assets" / "bay-cycle-throughput.png").read_bytes()
_SHA_A = hashlib.sha256(_IMAGE_A).hexdigest()
_IMAGE_B = _IMAGE_A + b"\n# a wholly different file"
_SHA_B = hashlib.sha256(_IMAGE_B).hexdigest()

# The complete record every case starts from (field order = the schema's).
_BASE_RECORD = {
    "file": "assets/shot.png",
    "sha256": _SHA_A,
    "folder": "Review pack",
    "date": "2374-08-24",
    "event": "Friday review",
}
# A second, always-well-formed asset. Its job is NON-VACUITY: every refusal case below also
# asserts that THIS one places, so no assertion can pass on a loader that places nothing.
_CLEAN_KEY = "clean-shot"
_CLEAN_RECORD = {
    "file": "assets/clean.png",
    "sha256": _SHA_A,
    "folder": "Review pack",
    "date": "2374-08-24",
    "event": "Friday review",
}


def _assets_yaml(records: dict) -> str:
    """Author an ``assets:`` subtree (mapping key → record) in the given order."""
    lines = ["assets:"]
    for key, record in records.items():
        lines.append(f"  {key}:")
        for field, value in record.items():
            lines.append(f'    {field}: "{value}"')
    return "\n".join(lines) + "\n"


def _write_weekly(
    root: pathlib.Path, records: dict, *, team: str = "", files: dict = None
) -> pathlib.Path:
    """Author one weekly + its asset files under ``root``; return the spec path."""
    (root / "assets").mkdir(parents=True, exist_ok=True)
    for name, payload in (files if files is not None else {}).items():
        (root / name).write_bytes(payload)
    path = root / "weekly.yml"
    path.write_text(
        'week: "W1"\nmodule: "M"\nhighlights:\n  - "One authored line."\n'
        + team
        + _assets_yaml(records),
        encoding="utf-8",
    )
    return path


def _asset_keys(load: WeeklySpecLoad) -> list[str]:
    return [block.asset.key for block in load.assets]


# Each case: (id, the mutation applied to _BASE_RECORD, the files on disk, the expected
# disclosure). "shot" is the asset under test; _CLEAN_KEY is always well-formed.
_ROUTING_CASES = [
    (
        "provenance-folder-absent",
        {"drop": "folder"},
        {"assets/shot.png": _IMAGE_A},
        _PROVENANCE_DISCLOSURE.format(key="shot", field="folder"),
    ),
    (
        "provenance-date-absent",
        {"drop": "date"},
        {"assets/shot.png": _IMAGE_A},
        _PROVENANCE_DISCLOSURE.format(key="shot", field="date"),
    ),
    (
        "provenance-event-absent",
        {"drop": "event"},
        {"assets/shot.png": _IMAGE_A},
        _PROVENANCE_DISCLOSURE.format(key="shot", field="event"),
    ),
    (
        "values-screenshot-without-deep-link",
        {"set": {"stands_in_for": "values"}},
        {"assets/shot.png": _IMAGE_A},
        _DEEP_LINK_DISCLOSURE.format(key="shot"),
    ),
    (
        "file-absent-on-disk",
        {"set": {"file": "assets/gone.png"}},
        {"assets/shot.png": _IMAGE_A},
        _CONTENT_ADDRESS_DISCLOSURE.format(key="shot", file="assets/gone.png"),
    ),
    (
        "content-substituted-record-describes-A-B-on-disk",
        {"set": {}},  # record keeps A's hash; B's bytes are written to A's path
        {"assets/shot.png": _IMAGE_B},
        _CONTENT_ADDRESS_DISCLOSURE.format(key="shot", file="assets/shot.png"),
    ),
]


@pytest.mark.parametrize(
    "case_id,mutation,files,expected",
    _ROUTING_CASES,
    ids=[case[0] for case in _ROUTING_CASES],
)
def test_asset_routing_refuses_and_discloses_exactly(
    tmp_path: pathlib.Path, case_id: str, mutation: dict, files: dict, expected: str
) -> None:
    """One row of the routing table: no ``AssetBlock``, and the spec's EXACT disclosure.

    Both ways, in one test: the malformed asset is refused with the spec's wording, and the
    well-formed one alongside it still places. Without the second half every assertion here
    would pass on a loader that places nothing at all.
    """
    record = dict(_BASE_RECORD)
    if "drop" in mutation:
        record.pop(mutation["drop"])
    record.update(mutation.get("set", {}))
    on_disk = dict(files)
    on_disk["assets/clean.png"] = _IMAGE_A

    path = _write_weekly(
        tmp_path, {"shot": record, _CLEAN_KEY: _CLEAN_RECORD}, files=on_disk
    )
    load = load_weekly_spec(path, root=tmp_path)

    assert "shot" not in _asset_keys(load), f"{case_id}: a refused asset was placed"
    assert expected in load.distillation.missing, (
        f"{case_id}: the spec's exact disclosure is absent.\n"
        f"expected: {expected!r}\ngot: {load.distillation.missing!r}"
    )
    # NON-VACUITY: the well-formed asset in the same document DID place.
    assert _CLEAN_KEY in _asset_keys(load), f"{case_id}: nothing placed — proof is vacuous"


def test_asset_routing_places_a_clean_record_with_its_provenance_and_a_trace(
    tmp_path: pathlib.Path,
) -> None:
    """The happy row: exactly one ``AssetBlock``, a populated record, and ≥1 trace INTO it."""
    record = dict(_BASE_RECORD)
    record["link"] = "https://example.invalid/report/1"
    record["stands_in_for"] = "values"  # the link is present, so the row is satisfied
    record["caption"] = "A caption the author wrote."
    path = _write_weekly(tmp_path, {"shot": record}, files={"assets/shot.png": _IMAGE_A})
    load = load_weekly_spec(path, root=tmp_path)

    assert _asset_keys(load) == ["shot"]
    block = load.assets[0]
    assert block.kind == "asset"
    assert block.asset.file == record["file"]
    assert block.asset.sha256 == _SHA_A
    assert (block.asset.folder, block.asset.date, block.asset.event) == (
        record["folder"],
        record["date"],
        record["event"],
    )
    assert block.asset.link == record["link"]
    assert block.asset.stands_in_for == "values"
    assert block.caption == record["caption"]

    # The evidence is a real, content-addressed span of the RECORD itself — the provenance
    # claim is the asset's evidence (the record is the evidence, never the image).
    assert len(block.evidence) >= 1
    for trace in block.evidence:
        assert trace.is_addressed, "an asset trace must be content-addressed"
        assert trace.source_id == load.source.id
        sliced = load.source.transcript[trace.start : trace.end]
        assert _SHA_A in sliced, "the trace must address the record's own sha256 span"

    # No routing disclosure fired for a clean asset.
    for note in load.distillation.missing:
        assert "asset 'shot'" not in note, note


@pytest.mark.parametrize(
    "photo,places,disclosed",
    [
        ("no-such-key", False, True),  # names no assets: entry at all
        ("shot", False, True),  # names an asset that was NOT placed
        (_CLEAN_KEY, True, False),  # names a PLACED asset — survives (non-vacuity)
    ],
    ids=["photo-names-nothing", "photo-names-an-unplaced-asset", "photo-names-a-placed-asset"],
)
def test_asset_routing_resolves_team_photo_against_placed_assets(
    tmp_path: pathlib.Path, photo: str, places: bool, disclosed: bool
) -> None:
    """The two reference rows, plus the surviving case that keeps them non-vacuous."""
    name = "Crewman Nog"
    team = f'team:\n  - name: "{name}"\n    role: "Bay Systems"\n    photo: "{photo}"\n'
    broken = dict(_BASE_RECORD)
    broken.pop("folder")  # "shot" is authored but NEVER placed
    path = _write_weekly(
        tmp_path,
        {"shot": broken, _CLEAN_KEY: _CLEAN_RECORD},
        team=team,
        files={"assets/shot.png": _IMAGE_A, "assets/clean.png": _IMAGE_A},
    )
    load = load_weekly_spec(path, root=tmp_path)

    # The member is CARRIED in every row — only the photo is at stake.
    assert [m.name for m in load.spec.team] == [name]
    assert load.spec.team[0].photo == photo, "the spec carries what the author wrote"

    note = _PHOTO_DISCLOSURE.format(name=name, key=photo)
    if disclosed:
        assert note in load.distillation.missing, load.distillation.missing
    else:
        assert note not in load.distillation.missing
    assert (photo in _asset_keys(load)) is places


def test_asset_routing_root_escape_raises_and_never_reaches_missing(
    tmp_path: pathlib.Path,
) -> None:
    """A path escaping the root is a REFUSAL, not an absence — and never a `missing[]` note."""
    root = tmp_path / "proj"
    root.mkdir()
    (tmp_path / "outside.png").write_bytes(_IMAGE_A)
    record = dict(_BASE_RECORD)
    record["file"] = "../outside.png"
    path = _write_weekly(root, {"shot": record}, files={})

    with pytest.raises(ValueError) as excinfo:
        load_weekly_spec(path, root=root)

    message = str(excinfo.value)
    assert "shot" in message, "the refusal must name the asset key"
    assert "../outside.png" in message, "the refusal must name the offending path"
    # It is a refusal, in the refusal's own words — never one of the missing[] disclosures.
    assert "disclosed, never placed" not in message
    assert "does not match its recorded content address" not in message

    # CONSTRUCTING ARM: the identical record loads once the root legitimately contains the
    # file — so the raise above is the containment check, not a broken placement path.
    inside = _write_weekly(
        root,
        {"shot": {**_BASE_RECORD, "file": "assets/shot.png"}},
        files={"assets/shot.png": _IMAGE_A},
    )
    load = load_weekly_spec(inside, root=root)
    assert _asset_keys(load) == ["shot"]


def test_asset_routing_refuses_a_symlink_resolving_outside_the_root(
    tmp_path: pathlib.Path,
) -> None:
    """The non-obvious half: ``Path.resolve()`` follows symlinks BEFORE the containment test.

    An in-root symlink whose target is outside the root is refused exactly like a literal
    ``../`` escape — a containment check that tested the authored string would miss this.
    """
    root = tmp_path / "proj"
    (root / "assets").mkdir(parents=True)
    target = tmp_path / "outside.png"
    target.write_bytes(_IMAGE_A)
    (root / "assets" / "link.png").symlink_to(target)

    record = dict(_BASE_RECORD)
    record["file"] = "assets/link.png"
    path = _write_weekly(root, {"shot": record}, files={})

    with pytest.raises(ValueError) as excinfo:
        load_weekly_spec(path, root=root)
    message = str(excinfo.value)
    assert "shot" in message and "assets/link.png" in message
    assert str(target) in message, "the refusal names the RESOLVED target, not just the link"


def test_asset_routing_hashes_the_image_and_never_decodes_it(
    tmp_path: pathlib.Path,
) -> None:
    """T-03-10: the loader hashes bytes. No image library is reachable, and no byte is written."""
    source = (
        pathlib.Path(newsletters.weeklyspec.__file__).read_text(encoding="utf-8").splitlines()
    )
    code = "\n".join(line for line in source if not line.lstrip().startswith("#"))
    for token in ("PIL", "Pillow", "imghdr"):
        assert token not in code, f"{token!r} appears in weeklyspec.py — the image is HASHED only"
    assert "hashlib.sha256" in code, "the content address must come from hashlib.sha256"

    # And the file itself is untouched by a load that hashes it.
    path = _write_weekly(tmp_path, {"shot": dict(_BASE_RECORD)}, files={"assets/shot.png": _IMAGE_A})
    image = tmp_path / "assets" / "shot.png"
    before = image.stat()
    load = load_weekly_spec(path, root=tmp_path)
    after = image.stat()
    assert _asset_keys(load) == ["shot"], "non-vacuity: the load actually read the image"
    assert (after.st_mtime_ns, after.st_size) == (before.st_mtime_ns, before.st_size)


def test_unknown_stands_in_for_value_fails_loud(tmp_path: pathlib.Path) -> None:
    """``stands_in_for`` is author-declared AND closed: an unknown kind is a teaching refusal.

    ``AssetRecord.stands_in_for`` is ``Literal["values"] | None``; carrying an unknown string
    that far would surface as a Pydantic ``ValidationError`` at placement instead of naming the
    typo, so the schema refuses it where every other typo is refused.
    """
    record = dict(_BASE_RECORD)
    record["stands_in_for"] = "trends"
    path = _write_weekly(tmp_path, {"shot": record}, files={"assets/shot.png": _IMAGE_A})
    with pytest.raises(ValueError) as excinfo:
        load_weekly_spec(path, root=tmp_path)
    message = str(excinfo.value)
    assert "stands_in_for" in message and "'trends'" in message and "'values'" in message


def test_committed_full_fixture_places_exactly_its_one_complete_asset() -> None:
    """The committed corpus routes all three of its assets, and exactly one survives."""
    load = _load(FULL)
    assert _asset_keys(load) == ["bay-cycle-throughput"]
    disclosures = "\n".join(load.distillation.missing)
    assert (
        _DEEP_LINK_DISCLOSURE.format(key="crew-rota-board") in disclosures
    ), disclosures
    assert (
        _PROVENANCE_DISCLOSURE.format(key="crew-manifest-scan", field="folder") in disclosures
    ), disclosures
    # The team photo names the one asset that DID place, so it is not disclosed.
    assert "photo key" not in disclosures


# --------------------------------------------------------------------------- #
# 12 — build_weekly_report: fixed order, Draft on arrival, and no editorializing
# --------------------------------------------------------------------------- #

# The fixed block order, pinned as a list of discriminator values. It is part of the
# determinism claim: a composer free to reorder blocks is a composer forming an opinion about
# what matters this week.
_FULL_BLOCK_ORDER = ["prose", "narrative", "narrative", "recognitions", "team", "asset"]

# Keys that carry STRUCTURE or PROVENANCE rather than display text. The editorialization scan
# skips them on purpose: `kind`/`tone` are discriminators, `topics` are schema field names, and a
# trace's ids/hashes/offsets address the record instead of being read from it.
_STRUCTURAL_KEYS = {
    "kind",
    "tone",
    "topics",
    "evidence",
    "source_id",
    "content_hash",
    "start",
    "end",
    "confidence",
}


def _report(name: str = FULL, **kwargs):
    return build_weekly_report(_load(name), author=AUTHOR, **kwargs)


def _display_strings(block) -> list[str]:
    """Every string a reader could SEE in one block (structure and provenance excluded)."""
    out: list[str] = []

    def walk(node) -> None:
        if isinstance(node, dict):
            for key, value in node.items():
                if key not in _STRUCTURAL_KEYS:
                    walk(value)
        elif isinstance(node, list):
            for value in node:
                walk(value)
        elif isinstance(node, str) and node:
            out.append(node)

    walk(block.model_dump(mode="json"))
    return out


def _blocks_of(surface, kind: str) -> list:
    return [b for b in surface.blocks if b.kind == kind]


def _unauthored(surface, transcript: str) -> list[str]:
    """Every block string the AUTHOR did not write and the composer did not declare.

    Containment is tested through the faithfulness gate's OWN normal form (``_normalize``:
    case-folded, whitespace-collapsed) rather than a raw ``in``. A YAML block scalar's parsed
    value is folded — its line breaks and indentation differ from the file — so a raw substring
    test would flag the author's own multi-line highlight as editorializing while still letting a
    reformatted paraphrase through. The gate's normal form is the honest comparator: it forgives
    whitespace and nothing else.
    """
    haystack = _normalize(transcript)
    return [
        text
        for block in surface.blocks
        for text in _display_strings(block)
        if _normalize(text) not in haystack and text not in CONNECTIVE_CONSTANTS
    ]


def test_weekly_report_is_a_draft_report_surface_at_epoch_zero() -> None:
    """SC-5's composition half: Draft on arrival, REPORT template, EPOCH_ZERO, author bylined."""
    surface = _report()
    assert surface.template is REPORT
    assert surface.review.state is ReviewState.DRAFT
    assert not surface.is_published
    assert surface.created == EPOCH_ZERO
    assert surface.byline == [AUTHOR] and surface.review.author == AUTHOR
    assert surface.traces == [_load(FULL).source]


def test_weekly_report_block_order_is_fixed() -> None:
    """The order is asserted, not assumed — and the assets keep their `assets:` file order."""
    surface = _report()
    assert [b.kind for b in surface.blocks] == _FULL_BLOCK_ORDER
    narratives = _blocks_of(surface, "narrative")
    assert [b.tone for b in narratives] == ["highlight", "lowlight"]
    assert [b.asset.key for b in _blocks_of(surface, "asset")] == ["bay-cycle-throughput"]


def test_weekly_report_omits_empty_blocks_and_the_absence_is_disclosed() -> None:
    """A weekly with nothing but its two scalars gets the lead and NOTHING else — honestly."""
    surface = _report(SPARSE)
    assert [b.kind for b in surface.blocks] == ["prose"], surface.blocks
    disclosures = "\n".join(surface.missing)
    for key in ("highlights", "lowlights", "recognitions", "team", "assets"):
        assert f"field '{key}' is absent or empty" in disclosures, key


def test_weekly_report_title_and_eyebrow_are_carried_verbatim_never_joined() -> None:
    """Identity is the author's two strings, each whole — a join would be authored connective text."""
    parsed = _parsed(FULL)
    surface = _report()
    assert surface.title == parsed["week"]
    assert surface.eyebrow == parsed["module"]


def test_weekly_report_falls_back_without_inventing_when_identity_is_absent(
    tmp_path: pathlib.Path,
) -> None:
    """No week/module: the title comes from the FILE, the eyebrow from a declared constant."""
    path = tmp_path / "bay-week-nine.yml"
    path.write_text('highlights:\n  - "One line."\n', encoding="utf-8")
    surface = build_weekly_report(
        load_weekly_spec(path, root=tmp_path), author=AUTHOR
    )
    assert surface.title == "Bay Week Nine"
    assert surface.eyebrow in CONNECTIVE_CONSTANTS


def test_weekly_report_cannot_publish_without_the_gate() -> None:
    """publish() without a recorded approval raises, and the state does NOT advance."""
    surface = _report()
    with pytest.raises(ValueError, match="No auto-publish"):
        surface.publish()
    assert surface.review.state is ReviewState.DRAFT, "a failed publish must not advance"

    # The gate — a recorded approval — is the one legitimate path, proving the refusal above
    # is the POLICY and not a broken publish.
    surface.publish(reviewer=AUTHOR)
    assert surface.is_published and surface.review.reviewer == AUTHOR


def test_weekly_composer_has_no_gate_advancing_call() -> None:
    """T-03-09: the module cannot auto-publish, because it names no gate-advancing method."""
    code = pathlib.Path(newsletters.weeklyspec.__file__).read_text(encoding="utf-8")
    for call in (".publish(", ".approve(", ".open_pull_request("):
        assert call not in code, f"{call!r} appears in weeklyspec.py — the composer returns Draft"


@pytest.mark.parametrize("name", FIXTURES)
def test_two_composes_of_two_loads_are_byte_identical(name: str) -> None:
    """SC-5: composing the same file twice, through two independent loads, is byte-identical."""
    first = build_weekly_report(_load(name), author=AUTHOR)
    second = build_weekly_report(_load(name), author=AUTHOR)
    assert first.model_dump_json() == second.model_dump_json()


def test_unaddressed_binding_claims_never_reach_a_claims_block() -> None:
    """The shared trust predicate, applied: untraced / un-addressed claims go to missing[]."""
    load = _load(FULL)
    good = next(c for c in load.distillation.claims if addressed(c))
    untraced = Claim(text="planted untraced binding claim")
    unaddressed = Claim(
        text="planted un-addressed binding claim",
        evidence=[Trace(source_id=load.source.id)],  # no content_hash -> not addressed
    )
    binding = SectionBinding(
        heading="Bay throughput", claims=[good, untraced, unaddressed]
    )
    surface = build_weekly_report(load, author=AUTHOR, bindings=[binding])

    claims_blocks = _blocks_of(surface, "claims")
    assert len(claims_blocks) == 1
    kept = {c.text for c in claims_blocks[0].claims}
    assert good.text in kept, "non-vacuity: the addressed claim DID reach the block"
    assert untraced.text not in kept and unaddressed.text not in kept
    assert untraced.text in surface.missing and unaddressed.text in surface.missing
    # A binding with no KPIs discloses the omitted strip rather than rendering an empty one.
    assert "section 'Bay throughput' declares no KPIs — strip omitted" in surface.missing
    assert not _blocks_of(surface, "kpi")


def test_weekly_report_missing_is_deduped_in_order() -> None:
    """IN-01 (03-review): a binding note repeating a loader disclosure reads ONCE, not twice.

    `compose_module_report` already runs the shared order-preserving dedup; the weekly composer
    now does too. Nothing is removed — only repeats — and order is preserved.
    """
    load = _load(SPARSE)
    duplicate = load.distillation.missing[0]
    binding = SectionBinding(heading="Bay throughput", missing=[duplicate])
    surface = build_weekly_report(load, author=AUTHOR, bindings=[binding])
    assert surface.missing.count(duplicate) == 1, surface.missing
    # Order-preserving: the deduped list is the original with repeats dropped, not a re-sort.
    seen: list[str] = []
    for note in load.distillation.missing:
        if note not in seen:
            seen.append(note)
    assert [n for n in surface.missing if n in seen] == seen


def test_kpi_delta_comes_from_compute_delta_and_is_never_re_derived() -> None:
    """The composer imports the ONE delta derivation; the strip sits before the claims block."""
    load = _load(FULL)
    endpoint_source = Source(
        id="endpoints", context="test", transcript="9 then 12", timestamp=EPOCH_ZERO
    )
    start = Claim(text="9", evidence=[Trace.from_source(endpoint_source, 0, 1)])
    close = Claim(text="12", evidence=[Trace.from_source(endpoint_source, 7, 9)])
    binding = SectionBinding(
        heading="Bay throughput",
        kpi_items=[KpiItem(label="Bay cycles", value="12")],
        kpi_endpoints=[[start, close]],
        claims=[start, close],
    )
    surface = build_weekly_report(load, author=AUTHOR, bindings=[binding])

    strips = _blocks_of(surface, "kpi")
    assert len(strips) == 1 and strips[0].heading == "Bay throughput"
    item = strips[0].items[0]
    assert (item.label, item.value) == ("Bay cycles", "12")
    assert (item.delta, item.dir) == compute_delta(start.text, close.text)
    kinds = [b.kind for b in surface.blocks]
    assert kinds.index("kpi") < kinds.index("claims"), "the strip precedes its claims"


def test_team_photo_is_none_unless_its_asset_was_placed(
    tmp_path: pathlib.Path,
) -> None:
    """The composed member carries photo=None for an unplaced key — and the key when placed."""
    surface = _report()
    members = _blocks_of(surface, "team")[0].members
    by_name = {m.name: m for m in members}
    placed = {b.asset.key for b in _blocks_of(surface, "asset")}
    photos = {m.name: m.photo for m in members}
    assert any(p in placed for p in photos.values() if p), "non-vacuity: one photo survived"
    for member in members:
        assert member.photo is None or member.photo in placed
        assert member.lines, f"{member.name} lost their authored lines"
    assert by_name, "the team block must carry members"


def test_composer_carries_baited_lines_separately_in_order_byte_identical() -> None:
    """The bait fixture: nothing summarized, nothing sorted, nothing merged."""
    parsed = _parsed(BAIT)
    surface = _report(BAIT)
    highlights = [b for b in _blocks_of(surface, "narrative") if b.tone == "highlight"]
    assert len(highlights) == 1, "one block per tone — never one merged narrative"
    rendered = [item.text for item in highlights[0].items]
    assert rendered == parsed["highlights"], "lines were reordered, merged or rewritten"
    assert len(rendered) == 6, "the bait must carry all three pairs"
    for item in highlights[0].items:
        assert item.claim is not None and item.claim.text == item.text
        assert addressed(item.claim), "each carried line stands on its own addressed span"


@pytest.mark.parametrize("name", FIXTURES)
def test_every_block_string_is_authored_or_a_declared_connective_constant(
    name: str,
) -> None:
    """FAITHFUL, NOT SUGGESTIVE, enforced: no text reaches a block that the author did not write.

    The only exceptions are the declared connective constants — the lead and the section labels,
    each numeral-free and fact-free — which are listed in ``CONNECTIVE_CONSTANTS`` so a composer
    that starts authoring prose has to change a declared allowlist to do it.
    """
    load = _load(name)
    surface = build_weekly_report(load, author=AUTHOR)
    for text in _unauthored(surface, load.source.transcript):
        raise AssertionError(
            f"block string {text!r} is neither authored in the spec file nor a declared "
            f"connective constant — the composer editorialized"
        )


def test_editorialization_guard_detects_a_planted_paraphrase() -> None:
    """NON-VACUITY: plant a paraphrase into a composed surface; the scan above must fire.

    Mirrors ``test_abstraction_guard.py::test_guard_detects_planted_leak`` — a guard whose
    firing has never been observed is a vibe, not a proof.
    """
    load = _load(BAIT)
    surface = build_weekly_report(load, author=AUTHOR)
    dumped = surface.model_dump(mode="json")
    paraphrase = "In short, the rota rewrite was the week's big win."
    planted = False
    for block in dumped["blocks"]:
        if block["kind"] == "narrative":
            block["items"][0]["text"] = paraphrase
            planted = True
            break
    assert planted, "the fixture must compose at least one narrative block"

    tampered = type(surface).model_validate(dumped)
    offenders = _unauthored(tampered, load.source.transcript)
    assert paraphrase in offenders, offenders
    # ...and the untampered surface is clean, so the guard is discriminating, not indiscriminate.
    assert _unauthored(surface, load.source.transcript) == []


def test_connective_constants_author_no_facts() -> None:
    """The allowlist's own contract: no numeral, and short enough to be connective only."""
    assert CONNECTIVE_CONSTANTS, "an empty allowlist would make the guard trivially strict"
    for text in CONNECTIVE_CONSTANTS:
        assert not any(ch.isdigit() for ch in text), f"{text!r} carries a numeral"


# --------------------------------------------------------------------------- #
# 13 — the deck: `weekly_slots` and the render through Phase 2's writer (SC-5)
#
# The slots half needs no optional extra and runs on every install. The render half needs
# python-pptx (the non-AI `[pptx]` extra), so it is guarded by a skipif — NOT by a module-level
# `importorskip`, which would skip this whole 70-plus-test module on a bare install and hide the
# authoring proofs behind an extra they do not use. The `weekly` CI job installs `[pptx]` and
# asserts `0 skipped`, so these execute there rather than skipping quietly (W21).
# --------------------------------------------------------------------------- #

try:  # noqa: SIM105 — the guard needs the bound name, not just the suppression
    import pptx as _pptx
except ImportError:  # pragma: no cover — exercised on a bare install, not in the weekly job
    _pptx = None

requires_pptx = pytest.mark.skipif(
    _pptx is None, reason="optional [pptx] extra (python-pptx) not installed"
)

# READ-ONLY. Nothing in this module writes to FIXTURE_DIR; every render goes to `tmp_path`.
COMMITTED_TEMPLATE = FIXTURE_DIR / "template.pptx"

# The template's deliberately UNPREFIXED shape and its authored body (tests/fixtures/weekly/
# _author_template.py). It is not a renderer slot, so a render must leave it exactly as authored.
FOOTER_SHAPE = "Footer"
FOOTER_TEXT = "fabricated fixture deck - not a renderer slot"


def _slots(name: str = FULL) -> tuple:
    """One load, its composed surface, and the slots derived from the pair."""
    load = _load(name)
    surface = build_weekly_report(load, author=AUTHOR)
    return load, surface, weekly_slots(load, surface)


def _shape_text(prs, name: str) -> str:
    """The text of the named shape in a REOPENED deck (never the writer's return value)."""
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.name == name:
                return shape.text_frame.text
    raise AssertionError(f"the written deck has no shape named {name!r}")


@pytest.mark.parametrize("name", FIXTURES)
def test_weekly_slots_emit_exactly_the_four_declared_keys_in_order(name: str) -> None:
    """The shape of the mapping ``bind_slots`` will validate — four keys, fixed order, list[str].

    All four keys are ALWAYS emitted, including for a section the author left empty: the template
    declares the shape, and ``bind_slots`` refuses an ``NL_``-prefixed shape with no matching
    content, so omitting the key would make a weekly with no lowlights fail to render at all.

    Values are explicit ``list[str]`` because a bare ``str`` is itself a ``Sequence[str]`` of its
    characters (02-review WR-03) — the type is the assertion, not a style preference.
    """
    _, _, slots = _slots(name)
    assert list(slots) == [
        WEEK_TITLE_SLOT,
        MODULE_SLOT,
        HIGHLIGHTS_SLOT,
        LOWLIGHTS_SLOT,
    ], "the slot mapping's key ORDER is part of the determinism claim"
    for key, lines in slots.items():
        assert key.startswith("NL_"), f"{key!r} is not a reserved-prefix slot name"
        assert isinstance(lines, list), f"{key!r} carries {type(lines).__name__}, not a list"
        assert lines, f"{key!r} is empty — `fill_slot` refuses an empty slot"
        assert all(isinstance(line, str) for line in lines), lines
        assert any(line.strip() for line in lines), f"{key!r} is all blank — a blank box"


@pytest.mark.parametrize("name", FIXTURES)
def test_weekly_slots_are_pure_and_repeatable(name: str) -> None:
    """Two derivations of one weekly are equal, key order included — half of "same deck twice"."""
    load = _load(name)
    surface = build_weekly_report(load, author=AUTHOR)
    first = weekly_slots(load, surface)
    second = weekly_slots(load, surface)
    assert first == second
    assert list(first) == list(second)
    # ...and the derivation READ the surface: it is `model_dump()`-identical afterwards.
    assert surface.model_dump() == build_weekly_report(load, author=AUTHOR).model_dump()


@pytest.mark.parametrize("name", FIXTURES)
def test_every_slot_line_is_authored_or_the_surface_own_disclosure(name: str) -> None:
    """The content contract: no third source of text can reach a slide.

    Authored lines are compared through the faithfulness gate's own normal form (see
    :func:`_unauthored` — a YAML block scalar's parsed value is folded, so a raw ``in`` would flag
    the author's own multi-line highlight). Everything else must be a member of
    ``surface.missing`` — the disclosure the loader wrote, not prose the composer invented.
    """
    load, surface, slots = _slots(name)
    haystack = _normalize(load.source.transcript)
    disclosed = set(surface.missing)
    for key, lines in slots.items():
        for line in lines:
            assert _normalize(line) in haystack or line in disclosed, (
                f"slot {key!r} carries {line!r}, which is neither the author's words nor a "
                f"recorded disclosure — that is composer-invented content on a slide"
            )


def test_sparse_weekly_empty_sections_carry_their_own_disclosure_line() -> None:
    """SC-5's real test: the weekly with no lowlights still renders, and says so on the slide.

    "No lowlights" is the absence a weekly is most tempted to hide (docs/weekly-spec.md rule 4).
    The single line on that slide IS the honesty-panel entry — not "Nothing to report", not an em
    dash, not an empty box.
    """
    _, surface, slots = _slots(SPARSE)
    for key, authored in ((HIGHLIGHTS_SLOT, "highlights"), (LOWLIGHTS_SLOT, "lowlights")):
        assert len(slots[key]) == 1, f"{key!r} should carry exactly the one disclosure line"
        line = slots[key][0]
        assert line in surface.missing, f"{line!r} is not in the honesty panel"
        assert authored in line, f"{line!r} does not name the section it discloses"
    # The authored halves of the same weekly are still the author's own words.
    assert slots[WEEK_TITLE_SLOT] == [_parsed(SPARSE)["week"]]
    assert slots[MODULE_SLOT] == [_parsed(SPARSE)["module"]]


def test_weekly_slots_refuse_a_disclosure_line_the_surface_never_recorded() -> None:
    """NON-VACUITY for the membership self-check — the mechanism that keeps the line honest.

    Strip the lowlights disclosure out of a composed surface and re-derive: the composer must
    REFUSE rather than write a line the record does not carry. Without this check the disclosure
    branch would be free to drift into any string at all, which is exactly the invented-prose
    failure the branch exists to prevent (T-03-20).
    """
    load = _load(SPARSE)
    surface = build_weekly_report(load, author=AUTHOR)
    dumped = surface.model_dump(mode="json")
    before = len(dumped["missing"])
    dumped["missing"] = [m for m in dumped["missing"] if "'lowlights'" not in m]
    assert len(dumped["missing"]) < before, "the fixture must disclose absent lowlights"
    tampered = type(surface).model_validate(dumped)

    with pytest.raises(ValueError) as excinfo:
        weekly_slots(load, tampered)
    message = str(excinfo.value)
    assert LOWLIGHTS_SLOT in message
    assert "'lowlights'" in message
    assert "composer-invented" in message


# --- the render itself: every assertion made by reopening the WRITTEN bytes ---------------- #


@requires_pptx
def test_full_weekly_renders_byte_identically_twice_in_process() -> None:
    """SC-5 determinism, both assertions, each for its own stated reason.

    ``bytes_a == bytes_b`` is the IN-PROCESS claim: one interpreter, one zlib, so the container is
    comparable and a difference would mean the writer embedded a clock. ``part_digest`` is the
    implementation-INDEPENDENT claim — DEFLATE output varies between zlib and zlib-ng, so this is
    the assertion a cross-environment gate (Phase 4's committed==fresh) must use.

    Nothing here compares a rendered deck to the TEMPLATE: reopening and re-saving the template
    already yields a different digest, a python-pptx load-path property and not a regression
    (Phase 2-03 decision).
    """
    _, surface, slots = _slots(FULL)
    first = render_surface_pptx_bytes(
        surface, template=COMMITTED_TEMPLATE, slots=slots
    )
    second = render_surface_pptx_bytes(
        surface, template=COMMITTED_TEMPLATE, slots=slots
    )
    assert first == second, "two in-process renders differ — the writer embedded a clock"
    assert part_digest(first) == part_digest(second), (
        "the part-content digests differ, so the DIFFERENCE IS IN THE CONTENT, not the zip "
        "container — this is the assertion a cross-environment gate inherits"
    )


@pytest.mark.parametrize("name", [FULL, SPARSE])
@requires_pptx
def test_weekly_deck_reads_back_marked_watermarked_and_slot_faithful(
    name: str, tmp_path: pathlib.Path
) -> None:
    """The written deck, REOPENED: the marker, the Draft watermark, and every slot line.

    The sparse weekly is rendered too, and that is the point of the parametrization: it is the
    weekly whose lowlights slot carries only the disclosure, so this proves the empty-section
    mechanism produces a real, openable deck rather than a refusal.
    """
    _, surface, slots = _slots(name)
    out_path = tmp_path / f"{name}.pptx"

    render_surface_pptx(
        surface, template=COMMITTED_TEMPLATE, slots=slots, out_path=out_path
    )

    written = _pptx.Presentation(str(out_path))  # reopen the FILE that was written
    core = written.core_properties
    assert core.category == MARKER, core.category
    assert core.content_status == DRAFT_STATUS, core.content_status
    assert core.identifier == surface.id, core.identifier

    for index, slide in enumerate(written.slides):
        names = [shape.name for shape in slide.shapes]
        assert WATERMARK_NAME in names, (
            f"slide {index} of a DRAFT weekly carries no watermark — an unreviewed page that "
            f"does not look unreviewed. Shapes: {names}"
        )

    for key, lines in slots.items():
        assert _shape_text(written, key) == "\n".join(lines), (
            f"slot {key!r} did not read back as the lines it was given"
        )
    assert _shape_text(written, FOOTER_SHAPE) == FOOTER_TEXT, (
        "the unprefixed shape was written — the operator's footer is not a renderer slot"
    )


@requires_pptx
def test_rendering_never_advances_the_gate_or_mutates_the_surface(
    tmp_path: pathlib.Path,
) -> None:
    """The product's hardest rule, on the deck path: a render is a READ of the gate.

    Asserted in the direction that can actually fail — the whole ``model_dump()`` before and
    after, not just ``review.state`` — because a writer that touched any other field would still
    be writing back into the reviewed record.
    """
    _, surface, slots = _slots(FULL)
    before = surface.model_dump()

    render_surface_pptx(
        surface,
        template=COMMITTED_TEMPLATE,
        slots=slots,
        out_path=tmp_path / "gate.pptx",
    )

    assert surface.model_dump() == before, "the render mutated the Surface"
    assert surface.review.state is ReviewState.DRAFT
    assert not surface.is_published


@requires_pptx
def test_the_committed_template_is_never_written_by_a_render(
    tmp_path: pathlib.Path,
) -> None:
    """READ-ONLY: the fixture template's bytes, mtime and size survive a render untouched."""
    before_bytes = COMMITTED_TEMPLATE.read_bytes()
    before_stat = COMMITTED_TEMPLATE.stat()

    _, surface, slots = _slots(FULL)
    render_surface_pptx(
        surface,
        template=COMMITTED_TEMPLATE,
        slots=slots,
        out_path=tmp_path / "readonly.pptx",
    )

    assert COMMITTED_TEMPLATE.read_bytes() == before_bytes
    assert COMMITTED_TEMPLATE.stat().st_mtime == before_stat.st_mtime
    assert COMMITTED_TEMPLATE.stat().st_size == before_stat.st_size
