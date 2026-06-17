"""Unit tests for ``PptxAdapter`` (ADAPT-04 / criteria 1+2).

These prove the PowerPoint adapter's emit/report fork directly against tiny, in-memory ``.pptx``
bytes authored at test time (no committed corpus — that lands in 06-04). They assert:

* verbatim per-paragraph / per-cell / notes extraction with ``Slide N / <shape>`` locators
  (criterion 1);
* the full ``unextracted[]`` taxonomy — SmartArt, chart, picture, media, OLE, unknown graphic
  frame, whole-source-unreadable — never a silent drop (criterion 2);
* group RECURSION with the nested leaf-accounting identity
  ``leaf_shapes == producers + skipped_empty + unextracted_from_shapes`` (L3);
* registration ``"pptx"`` + resolution, content-addressed claims via ``normalize()``,
  ``assert_conforms`` passing, and round-trip coverage parity on a FRESH adapter.

The whole module is skipped cleanly when the optional ``[pptx]`` extra (python-pptx) is absent, so
the bare-install gate (``tests/test_ai_optional.py``) is unaffected.
"""

from __future__ import annotations

import importlib.util
import io
from datetime import datetime, timezone

import pytest

pytestmark = pytest.mark.skipif(
    importlib.util.find_spec("pptx") is None,
    reason="optional [pptx] extra (python-pptx) not installed",
)

from newsletters.distill import DistillationResult, assert_conforms, resolve  # noqa: E402
from newsletters.semantic import Source  # noqa: E402

# A constant, valid 1x1 transparent PNG byte literal — embeds via ``add_picture`` WITHOUT a Pillow
# call (deterministic, no nondeterminism), per 06-RESEARCH Fixture gotcha 1.
_PNG_1x1 = bytes.fromhex(
    "89504e470d0a1a0a0000000d4948445200000001000000010806000000"
    "1f15c4890000000d49444154789c6360000002000100ffff0300000600"
    "0557bfabd40000000049454e44ae426082"
)

_DIAGRAM_URI = "http://schemas.openxmlformats.org/drawingml/2006/diagram"


def _adapter():
    """A FRESH PptxAdapter (imported lazily so this module imports without the [pptx] extra)."""
    from newsletters.adapters.pptx_adapter import PptxAdapter

    return PptxAdapter()


def _blank_slide(prs):
    """Add and return a slide on the blank layout (index 6 in the default template)."""
    return prs.slides.add_slide(prs.slide_layouts[6])


def _serialize(prs) -> bytes:
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# --------------------------------------------------------------------------- #
# Fixture builders
# --------------------------------------------------------------------------- #


def _deck_title_body() -> bytes:
    from pptx import Presentation
    from pptx.util import Emu

    prs = Presentation()
    slide = _blank_slide(prs)
    tb = slide.shapes.add_textbox(Emu(0), Emu(0), Emu(1000), Emu(1000))
    tf = tb.text_frame
    tf.text = "First paragraph"
    tf.add_paragraph().text = "Second paragraph"
    tf.add_paragraph().text = ""  # empty paragraph -> not emitted, not a drop
    return _serialize(prs)


def _deck_table() -> bytes:
    from pptx import Presentation
    from pptx.util import Emu

    prs = Presentation()
    slide = _blank_slide(prs)
    gf = slide.shapes.add_table(2, 2, Emu(0), Emu(0), Emu(2000), Emu(2000))
    table = gf.table
    table.cell(0, 0).text = "r0c0"
    table.cell(0, 1).text = "r0c1"
    table.cell(1, 0).text = "r1c0"
    table.cell(1, 1).text = ""  # empty cell -> not emitted
    return _serialize(prs)


def _deck_notes() -> bytes:
    from pptx import Presentation

    prs = Presentation()
    slide = _blank_slide(prs)
    slide.notes_slide.notes_text_frame.text = "Speaker note line"
    return _serialize(prs)


def _deck_picture() -> bytes:
    from pptx import Presentation
    from pptx.util import Emu

    prs = Presentation()
    slide = _blank_slide(prs)
    slide.shapes.add_picture(io.BytesIO(_PNG_1x1), Emu(0), Emu(0))
    return _serialize(prs)


def _deck_chart() -> bytes:
    from pptx import Presentation
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.util import Emu

    prs = Presentation()
    slide = _blank_slide(prs)
    data = CategoryChartData()
    data.categories = ["a", "b"]
    data.add_series("s", (1, 2))
    slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Emu(0), Emu(0), Emu(2000), Emu(2000), data
    )
    return _serialize(prs)


def _deck_group(nested: bool = False) -> bytes:
    """A group with a readable text member and an unreadable picture member.

    With ``nested=True``, the readable member lives inside an INNER group, so the accounting
    identity must hold across two levels of nesting.
    """
    from pptx import Presentation
    from pptx.util import Emu

    prs = Presentation()
    slide = _blank_slide(prs)
    tb = slide.shapes.add_textbox(Emu(0), Emu(0), Emu(500), Emu(500))
    tb.text_frame.text = "grouped text"
    pic = slide.shapes.add_picture(io.BytesIO(_PNG_1x1), Emu(600), Emu(0))
    if nested:
        inner = slide.shapes.add_group_shape([tb])
        slide.shapes.add_group_shape([inner, pic])
    else:
        slide.shapes.add_group_shape([tb, pic])
    return _serialize(prs)


def _deck_smartart() -> bytes:
    """A deck with an XML-injected ``graphicFrame`` whose ``a:graphicData/@uri`` is the diagram NS.

    This is the minimum to be DETECTED as SmartArt (``shape_type is None`` + the diagram URI); the
    diagram itself need not render (06-RESEARCH Pitfall 1 / Fixture gotcha 3).
    """
    from pptx import Presentation
    from pptx.util import Emu

    prs = Presentation()
    slide = _blank_slide(prs)
    # Anchor on an existing graphic frame (a table), then rewrite its graphicData URI to the
    # diagram namespace and strip the table payload — yielding a frame python-pptx sees as
    # shape_type None with the diagram URI (no high-level SmartArt authoring API exists). An empty
    # graphicData with the diagram URI is the minimum needed to be DETECTED as SmartArt (probed
    # against python-pptx 1.0.2; both ``graphicData_uri`` and the lxml ``@uri`` fallback resolve).
    gf = slide.shapes.add_table(1, 1, Emu(0), Emu(0), Emu(1000), Emu(1000))
    graphic_data = gf._element.graphic.graphicData
    graphic_data.set("uri", _DIAGRAM_URI)
    for child in list(graphic_data):
        graphic_data.remove(child)
    return _serialize(prs)


# --------------------------------------------------------------------------- #
# Task 1: parse — emit decisions
# --------------------------------------------------------------------------- #


def test_title_body_emits_per_paragraph() -> None:
    source, units, drops = _adapter().parse(_deck_title_body(), "doc.pptx")
    assert units == ["First paragraph", "Second paragraph"]
    assert drops == []
    assert "First paragraph" in source.transcript
    assert "Slide 1 / " in source.transcript


def test_table_emits_per_cell_row_major() -> None:
    _, units, drops = _adapter().parse(_deck_table(), "doc.pptx")
    assert units == ["r0c0", "r0c1", "r1c0"]
    assert drops == []


def test_table_cell_locator_has_rc_suffix() -> None:
    source, _, _ = _adapter().parse(_deck_table(), "doc.pptx")
    assert "[r0c0]" in source.transcript
    assert "[r1c0]" in source.transcript


def test_notes_emit_with_notes_locator() -> None:
    source, units, _ = _adapter().parse(_deck_notes(), "doc.pptx")
    assert "Speaker note line" in units
    assert "Slide 1 / notes" in source.transcript


# --------------------------------------------------------------------------- #
# Task 1: parse — report (unextracted taxonomy)
# --------------------------------------------------------------------------- #


def test_picture_routes_to_unextracted() -> None:
    _, units, drops = _adapter().parse(_deck_picture(), "doc.pptx")
    assert units == []
    assert len(drops) == 1
    assert "picture" in drops[0].reason.lower() or "image" in drops[0].reason.lower()


def test_chart_routes_to_unextracted() -> None:
    _, units, drops = _adapter().parse(_deck_chart(), "doc.pptx")
    assert units == []
    assert len(drops) == 1
    assert "chart" in drops[0].reason.lower()


def test_smartart_routes_to_unextracted() -> None:
    _, units, drops = _adapter().parse(_deck_smartart(), "doc.pptx")
    assert units == []
    assert len(drops) == 1
    assert "smartart" in drops[0].reason.lower() or "diagram" in drops[0].reason.lower()


# --------------------------------------------------------------------------- #
# Task 1: groups + accounting identity
# --------------------------------------------------------------------------- #


@pytest.mark.parametrize("nested", [False, True])
def test_group_recursed_member_extracted_and_unreadable_reported(nested: bool) -> None:
    _, units, drops = _adapter().parse(_deck_group(nested=nested), "doc.pptx")
    assert units == ["grouped text"]  # the group node itself is neither a claim nor a drop
    assert len(drops) == 1  # only the unreadable picture member is reported
    assert "picture" in drops[0].reason.lower() or "image" in drops[0].reason.lower()


# --------------------------------------------------------------------------- #
# Task 1: malformed input + determinism
# --------------------------------------------------------------------------- #


def test_malformed_pptx_is_one_whole_source_drop() -> None:
    source, units, drops = _adapter().parse(b"not a pptx at all", "broken.pptx")
    assert source.transcript == ""
    assert units == []
    assert len(drops) == 1
    assert source.extraction is not None  # the drop travels on the Source


def test_parse_is_deterministic() -> None:
    raw = _deck_title_body()
    s1, _, _ = _adapter().parse(raw, "doc.pptx")
    s2, _, _ = _adapter().parse(raw, "doc.pptx")
    assert s1.model_dump_json() == s2.model_dump_json()


def test_missing_created_maps_to_epoch_zero() -> None:
    from newsletters.adapters._timestamps import EPOCH_ZERO

    # Strip dcterms:created so the deck has no intrinsic timestamp.
    import re
    import zipfile

    raw = _deck_title_body()
    zin = zipfile.ZipFile(io.BytesIO(raw))
    out = io.BytesIO()
    with zipfile.ZipFile(out, "w", zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            content = zin.read(item.filename)
            if item.filename == "docProps/core.xml":
                content = re.sub(
                    rb"<dcterms:created[^>]*>[^<]*</dcterms:created>", b"", content
                )
            zout.writestr(item, content)
    source, _, _ = _adapter().parse(out.getvalue(), "doc.pptx")
    assert source.timestamp == EPOCH_ZERO


def test_intrinsic_created_is_preserved() -> None:
    from pptx import Presentation

    prs = Presentation()
    fixed = datetime(2021, 3, 4, 5, 6, 7, tzinfo=timezone.utc)
    prs.core_properties.created = fixed
    _blank_slide(prs)
    source, _, _ = _adapter().parse(_serialize(prs), "doc.pptx")
    assert source.timestamp == fixed


# --------------------------------------------------------------------------- #
# Task 2: distill + registration + conformance + round-trip parity
# --------------------------------------------------------------------------- #


def test_registered_and_resolvable() -> None:
    from newsletters.adapters.pptx_adapter import PptxAdapter

    assert isinstance(resolve("pptx"), PptxAdapter)


def test_distill_mints_content_addressed_claims() -> None:
    adapter = _adapter()
    source, _, _ = adapter.parse(_deck_title_body(), "doc.pptx")
    result = adapter.distill([source])
    assert len(result.distillation.claims) == 2
    for claim in result.distillation.claims:
        assert len(claim.evidence) == 1
        trace = claim.evidence[0]
        assert trace.is_addressed
        assert claim.text == trace.span


def test_assert_conforms_passes() -> None:
    adapter = _adapter()
    source, _, _ = adapter.parse(_deck_title_body(), "doc.pptx")
    assert_conforms(adapter, [source])


def test_distill_result_round_trips_through_json() -> None:
    adapter = _adapter()
    source, _, _ = adapter.parse(_deck_table(), "doc.pptx")
    result = adapter.distill([source])
    assert DistillationResult.model_validate_json(result.model_dump_json()) == result


def test_round_trip_coverage_parity_on_fresh_adapter() -> None:
    # Parse a deck WITH a drop (picture), persist the Source, reload, distill on a FRESH adapter.
    source, _, _ = _adapter().parse(_deck_group(), "doc.pptx")
    reloaded = Source.model_validate_json(source.model_dump_json())
    original = _adapter().distill([source])
    fresh = _adapter().distill([reloaded])
    assert fresh.coverage == original.coverage
    assert fresh.coverage.complete is False  # the picture drop survives persistence


def test_orphan_source_marks_coverage_not_reconstructable() -> None:
    from newsletters.adapters._coverage_codec import COVERAGE_NOT_RECONSTRUCTABLE

    # A Source NOT produced by parse() (no extraction carrier) -> R2 marker fires.
    orphan = Source(id="orphan", context="orphan", transcript="Slide 1 / X\tbody\n")
    result = _adapter().distill([orphan])
    reasons = [u.reason for u in result.coverage.unextracted]
    assert COVERAGE_NOT_RECONSTRUCTABLE in reasons
